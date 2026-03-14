"""
build_tm10_decks.py
-------------------
Generates two TM-10 slide decks:
  1. mss_platform_overview.pptx  — new deck, 10 slides
  2. Patches army_data_orientation_v1 slides 12-14 (writes a new v2 PPTX)

Run from repo root:
  python3 scripts/build_tm10_decks.py
"""

import copy
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ── Paths ─────────────────────────────────────────────────────────────────────
REPO = Path(__file__).resolve().parents[1]
SOURCE_PPTX = REPO / "maven_training/pdf/Army_Data_Orientation_PROTO.pptx"
OUT_DIR = REPO / "maven_training/source_material/course_portal/downloads"
OUT_DIR.mkdir(parents=True, exist_ok=True)

# ── Design system ─────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x17, 0x32, 0x5C)   # dark navy — header bars
TEAL   = RGBColor(0x1A, 0x78, 0x6E)   # teal accent
BLUE   = RGBColor(0x26, 0x4D, 0x7E)   # medium blue — secondary text
GOLD   = RGBColor(0xC8, 0x9A, 0x00)   # gold — callouts
RED    = RGBColor(0xBD, 0x2E, 0x2E)   # red — problems/before
GREEN  = RGBColor(0x2A, 0x6E, 0x3F)   # green — success/after
PURPLE = RGBColor(0x6B, 0x32, 0x8A)   # purple — applications layer
DARK   = RGBColor(0x0D, 0x1F, 0x3C)   # near-black — quote boxes
LGRAY  = RGBColor(0xF5, 0xF6, 0xF8)   # light gray — content panels
LBLUE  = RGBColor(0xD0, 0xE4, 0xF5)   # light blue
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
BLACK  = RGBColor(0x00, 0x00, 0x00)
FOOTER_COLOR = RGBColor(0x18, 0x33, 0x5F)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.50)
FOOTER_TEXT = "USAREUR-AF  ·  C2DAO  ·  MSS-TM10  ·  MAR 2026  ·  UNCLASSIFIED"
CLASSIF = "UNCLASSIFIED"

# ── Helper functions ──────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs

def blank_slide(prs):
    layout = prs.slide_layouts[6]   # completely blank
    return prs.slides.add_slide(layout)

def add_rect(slide, l, t, w, h, fill_color, line_color=None):
    shape = slide.shapes.add_shape(1, Emu(l), Emu(t), Emu(w), Emu(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape

def add_textbox(slide, l, t, w, h, text, font_size=12, bold=False,
                color=WHITE, align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(Emu(l), Emu(t), Emu(w), Emu(h))
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.color.rgb = color
    return txb

def add_multi(slide, l, t, w, h, lines, font_size=11, color=BLUE,
              bold_first=False, line_spacing=None):
    """Add a textbox with multiple paragraph lines."""
    txb = slide.shapes.add_textbox(Emu(l), Emu(t), Emu(w), Emu(h))
    tf  = txb.text_frame
    tf.word_wrap = True
    first = True
    for text in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = text
        run.font.size  = Pt(font_size)
        run.font.bold  = bold_first and (text == lines[0])
        run.font.color.rgb = color
    return txb

def add_classif_banners(slide):
    """Top and bottom UNCLASSIFIED banners."""
    # top banner (thin purple strip — matches existing decks)
    add_textbox(slide, 0, 0,
                int(SLIDE_W * 0.914), int(Inches(0.31)),
                CLASSIF, 8, False, FOOTER_COLOR, PP_ALIGN.LEFT)
    # bottom banner
    add_textbox(slide, 0, int(SLIDE_H - Inches(0.31)),
                int(SLIDE_W * 0.914), int(Inches(0.31)),
                CLASSIF, 8, False, FOOTER_COLOR, PP_ALIGN.LEFT)

def add_footer(slide):
    add_textbox(slide,
                int(Inches(0.3)), int(SLIDE_H - Inches(0.38)),
                int(Inches(12.5)), int(Inches(0.28)),
                FOOTER_TEXT, 8, False, FOOTER_COLOR, PP_ALIGN.LEFT)

def slide_title(slide, title_text):
    """Slide section title — small, bold, top-left."""
    add_textbox(slide,
                int(Inches(0.4)), int(Inches(0.35)),
                int(Inches(11)), int(Inches(0.45)),
                title_text, 18, True, NAVY, PP_ALIGN.LEFT)

def header_bar(slide, l, t, w, h, label, fill=NAVY):
    add_rect(slide, l, t, w, h, fill)
    add_textbox(slide, l + int(Inches(0.12)), t + int(Inches(0.04)),
                w - int(Inches(0.2)), h - int(Inches(0.05)),
                label, 10, True, WHITE, PP_ALIGN.LEFT)

def content_panel(slide, l, t, w, h):
    add_rect(slide, l, t, w, h, LGRAY)

def subtitle_text(slide, text):
    add_textbox(slide,
                int(Inches(0.4)), int(Inches(0.88)),
                int(Inches(12.2)), int(Inches(0.35)),
                text, 11, False, BLUE, PP_ALIGN.LEFT)

# ─────────────────────────────────────────────────────────────────────────────
# DECK 1: MSS PLATFORM OVERVIEW (10 slides)
# ─────────────────────────────────────────────────────────────────────────────

def build_mss_platform_overview():
    prs = new_prs()

    # ── Slide 1: Title ────────────────────────────────────────────────────────
    s = blank_slide(prs)
    add_classif_banners(s)
    # dark background block
    add_rect(s, 0, int(Inches(1.8)), int(SLIDE_W), int(Inches(4.5)), NAVY)
    add_textbox(s, int(Inches(1.5)), int(Inches(2.2)), int(Inches(10)), int(Inches(1.1)),
                "MSS PLATFORM OVERVIEW", 40, True, WHITE, PP_ALIGN.CENTER)
    add_textbox(s, int(Inches(1.5)), int(Inches(3.35)), int(Inches(10)), int(Inches(0.55)),
                "Your Guide to Using MSS", 20, False, LGRAY, PP_ALIGN.CENTER)
    add_textbox(s, int(Inches(1.5)), int(Inches(4.05)), int(Inches(10)), int(Inches(0.4)),
                "UNCLASSIFIED  ·  TM-10  ·  No Technical Background Required",
                12, False, GOLD, PP_ALIGN.CENTER)
    add_footer(s)

    # ── Slide 2: What is MSS? ─────────────────────────────────────────────────
    s = blank_slide(prs)
    add_classif_banners(s)
    slide_title(s, "WHAT IS MSS?")
    subtitle_text(s, "The Mission Support System — and why your unit is using it.")

    # Three columns
    col_w = int(Inches(3.7))
    col_h = int(Inches(4.5))
    col_t = int(Inches(1.35))
    gaps  = int(Inches(0.3))

    for i, (label, fill, lines) in enumerate([
        ("WHAT IT REPLACES", RED, [
            "Hundreds of separate Army business systems — each with its own login,",
            "its own format, and its own version of the truth.",
            "",
            "▸  Manual Excel readiness reports",
            "▸  Emailed PowerPoint briefings",
            "▸  Phone calls to get a status update",
            "▸  Data that was stale before it was printed",
        ]),
        ("WHAT MSS IS", TEAL, [
            "MSS is your unit's live data environment — a single platform where",
            "readiness, personnel, equipment, and operations data are always current.",
            "",
            "▸  One login. One platform.",
            "▸  Real-time data — no more stale spreadsheets",
            "▸  Your unit's picture and the larger formation",
            "▸  AI tools that answer your questions",
        ]),
        ("WHAT YOU'LL FIND IN IT", NAVY, [
            "The information your unit generates and depends on — all in one place.",
            "",
            "▸  Unit readiness and equipment status",
            "▸  Personnel accountability",
            "▸  Training records and schedules",
            "▸  Logistics and supply status",
            "▸  Commander dashboards and reports",
        ]),
    ]):
        col_l = int(Inches(0.3)) + i * (col_w + gaps)
        header_bar(s, col_l, col_t, col_w, int(Inches(0.45)), label, fill)
        content_panel(s, col_l, col_t + int(Inches(0.45)), col_w, col_h - int(Inches(0.45)))
        add_multi(s, col_l + int(Inches(0.15)), col_t + int(Inches(0.55)),
                  col_w - int(Inches(0.25)), col_h - int(Inches(0.65)),
                  lines, 10.5, DARK)

    add_footer(s)

    # ── Slide 3: Navigating the Platform ──────────────────────────────────────
    s = blank_slide(prs)
    add_classif_banners(s)
    slide_title(s, "NAVIGATING THE PLATFORM")
    subtitle_text(s, "Finding what you need — quickly.")

    # Left: step list | Right: tips panel
    left_w  = int(Inches(7.8))
    right_w = int(Inches(4.5))
    left_l  = int(Inches(0.3))
    right_l = int(Inches(8.4))
    panel_t = int(Inches(1.3))
    panel_h = int(Inches(5.5))

    header_bar(s, left_l, panel_t, left_w, int(Inches(0.45)), "HOW TO FIND INFORMATION", NAVY)
    content_panel(s, left_l, panel_t + int(Inches(0.45)), left_w, panel_h - int(Inches(0.45)))
    add_multi(s, left_l + int(Inches(0.18)), panel_t + int(Inches(0.62)),
              left_w - int(Inches(0.3)), panel_h - int(Inches(0.75)),
              [
                  "STEP 1 — Log In",
                  "Open your browser and navigate to your unit's MSS URL.",
                  "Enter your CAC credentials when prompted. MFA may be required.",
                  "",
                  "STEP 2 — Find Your Application",
                  "The home screen shows applications your unit has access to.",
                  "Look for readiness dashboards, personnel apps, or logistics views.",
                  "If you don't see what you need — ask your unit data manager.",
                  "",
                  "STEP 3 — Search for Specific Items",
                  "Use the search bar to find a specific soldier, vehicle, or piece of equipment.",
                  "Search by name, ID, or unit designation.",
                  "",
                  "STEP 4 — Use Filters",
                  "Most dashboards have filters at the top or side.",
                  "Filter by unit, date range, status, or location to narrow your view.",
                  "Filters change what you see — they do not change the data.",
              ], 10.5, DARK)

    header_bar(s, right_l, panel_t, right_w, int(Inches(0.45)), "QUICK TIPS", TEAL)
    content_panel(s, right_l, panel_t + int(Inches(0.45)), right_w, panel_h - int(Inches(0.45)))
    add_multi(s, right_l + int(Inches(0.15)), panel_t + int(Inches(0.62)),
              right_w - int(Inches(0.25)), panel_h - int(Inches(0.75)),
              [
                  "▸  Bookmark your most-used apps",
                  "",
                  "▸  Your access is tied to your role — if something is missing, contact your unit data manager",
                  "",
                  "▸  Don't share login credentials under any circumstances",
                  "",
                  "▸  If the page looks wrong or blank — try refreshing before calling for help",
                  "",
                  "▸  Navigation panels on the left or top vary by app — explore them",
              ], 10.5, DARK)

    add_footer(s)

    # ── Slide 4: Reading a Dashboard ──────────────────────────────────────────
    s = blank_slide(prs)
    add_classif_banners(s)
    slide_title(s, "READING A DASHBOARD")
    subtitle_text(s, "Dashboards show live data — here is how to interpret what you see.")

    # Four widget type boxes in a 2x2 grid
    bw = int(Inches(5.9))
    bh = int(Inches(2.5))
    bt = int(Inches(1.3))
    bl = int(Inches(0.3))
    gap = int(Inches(0.25))

    for row in range(2):
        for col in range(2):
            idx = row * 2 + col
            label, fill, lines = [
                ("STATUS INDICATORS", TEAL, [
                    "Color-coded tiles (green / amber / red) show health at a glance.",
                    "Green = meets standard.  Amber = degraded.  Red = non-mission-capable.",
                    "Hover or click a tile to see detail behind the status.",
                ]),
                ("CHARTS AND GRAPHS", NAVY, [
                    "Bar charts, pie charts, and trend lines summarize data visually.",
                    "The title tells you what is being measured.",
                    "Look for date labels — confirm the data is current before briefing it.",
                ]),
                ("DATA TABLES", BLUE, [
                    "Tables show individual records — one row per soldier, vehicle, or event.",
                    "Column headers tell you what each field contains.",
                    "Click a column header to sort.  Use row filters to narrow results.",
                ]),
                ("FILTERS AND CONTROLS", GOLD, [
                    "Dropdowns, date pickers, and checkboxes let you adjust your view.",
                    "Filters never change the underlying data — they only change what you see.",
                    "Reset filters to return to the default view.",
                ]),
            ][idx]
            lx = bl + col * (bw + gap)
            ty = bt + row * (bh + gap)
            header_bar(s, lx, ty, bw, int(Inches(0.4)), label, fill)
            content_panel(s, lx, ty + int(Inches(0.4)), bw, bh - int(Inches(0.4)))
            add_multi(s, lx + int(Inches(0.15)), ty + int(Inches(0.52)),
                      bw - int(Inches(0.25)), bh - int(Inches(0.58)),
                      lines, 10.5, DARK)

    add_footer(s)

    # ── Slide 5: Using the AI Assistant ───────────────────────────────────────
    s = blank_slide(prs)
    add_classif_banners(s)
    slide_title(s, "USING THE AI ASSISTANT (AIP)")
    subtitle_text(s, "AIP answers questions about your unit's data — in plain language.")

    left_w  = int(Inches(5.8))
    right_w = int(Inches(6.3))
    panel_t = int(Inches(1.3))
    panel_h = int(Inches(5.5))

    header_bar(s, int(Inches(0.3)), panel_t, left_w, int(Inches(0.45)), "WHAT YOU CAN ASK", GREEN)
    content_panel(s, int(Inches(0.3)), panel_t + int(Inches(0.45)), left_w, panel_h - int(Inches(0.45)))
    add_multi(s, int(Inches(0.48)), panel_t + int(Inches(0.62)),
              left_w - int(Inches(0.3)), panel_h - int(Inches(0.75)),
              [
                  "AIP reads your unit's data and answers in plain English.",
                  "Ask it like you would ask a knowledgeable colleague.",
                  "",
                  'Examples:',
                  '"What is the FMC rate for 1st platoon?"',
                  '"Which vehicles are NMC this week?"',
                  '"How many soldiers are current on APFT?"',
                  '"Show me readiness trends for the last 30 days."',
                  '"Who are the non-deployable personnel in Bravo Company?"',
                  "",
                  "BLUF: If you would ask it in a battle update brief,",
                  "you can ask AIP first.",
              ], 10.5, DARK)

    header_bar(s, int(Inches(6.45)), panel_t, right_w, int(Inches(0.45)),
               "WHAT TO WATCH OUT FOR", RED)
    content_panel(s, int(Inches(6.45)), panel_t + int(Inches(0.45)), right_w, panel_h - int(Inches(0.45)))
    add_multi(s, int(Inches(6.62)), panel_t + int(Inches(0.62)),
              right_w - int(Inches(0.3)), panel_h - int(Inches(0.75)),
              [
                  "AIP is only as good as the data behind it.",
                  "If the data is wrong or missing, the answer will be wrong too.",
                  "",
                  "▸  ALWAYS verify critical answers against source data",
                  "   before briefing a commander.",
                  "",
                  "▸  AIP cannot access classified data above your",
                  "   current network and access level.",
                  "",
                  "▸  Do not enter classified information into the",
                  "   AI chat window as a query.",
                  "",
                  "▸  Do not use AIP to make final decisions on",
                  "   personnel or readiness — it informs, it does not decide.",
                  "",
                  "▸  If an answer seems wrong — trust your instinct.",
                  "   Flag it to your unit data manager.",
              ], 10.5, DARK)

    add_footer(s)

    # ── Slide 6: Submitting Forms and Actions ─────────────────────────────────
    s = blank_slide(prs)
    add_classif_banners(s)
    slide_title(s, "SUBMITTING FORMS AND ACTIONS")
    subtitle_text(s, "How to enter data, submit requests, and record events in MSS.")

    # Two columns: How To + What Happens After
    col_w = int(Inches(5.8))
    panel_t = int(Inches(1.3))
    panel_h = int(Inches(4.5))

    header_bar(s, int(Inches(0.3)), panel_t, col_w, int(Inches(0.45)), "HOW TO SUBMIT", NAVY)
    content_panel(s, int(Inches(0.3)), panel_t + int(Inches(0.45)), col_w, panel_h - int(Inches(0.45)))
    add_multi(s, int(Inches(0.48)), panel_t + int(Inches(0.62)),
              col_w - int(Inches(0.3)), panel_h - int(Inches(0.75)),
              [
                  "1.  Open the form or action panel in your app.",
                  "",
                  "2.  Fill in required fields.",
                  "    Fields marked with * or highlighted in red are required.",
                  "",
                  "3.  Review before submitting.",
                  "    Check dates, unit designations, and status codes.",
                  "    Errors are harder to fix after submission.",
                  "",
                  "4.  Click Submit (or Save).",
                  "    You will see a confirmation message.",
                  "    If you see a red error — read it before trying again.",
                  "",
                  "5.  Verify the record updated.",
                  "    Refresh the dashboard to confirm your entry appears.",
              ], 10.5, DARK)

    header_bar(s, int(Inches(6.45)), panel_t, col_w, int(Inches(0.45)),
               "WHAT HAPPENS AFTER YOU SUBMIT", TEAL)
    content_panel(s, int(Inches(6.45)), panel_t + int(Inches(0.45)), col_w, panel_h - int(Inches(0.45)))
    add_multi(s, int(Inches(6.62)), panel_t + int(Inches(0.62)),
              col_w - int(Inches(0.3)), panel_h - int(Inches(0.75)),
              [
                  "Your entry goes into a live system — not a form that sits in an inbox.",
                  "",
                  "▸  Data is visible to others in your chain immediately",
                  "▸  It feeds dashboards the commander sees",
                  "▸  It may trigger automated alerts or workflows",
                  "",
                  "This is why accuracy matters.",
                  "One wrong status entry shows up everywhere.",
              ], 10.5, DARK)

    # Warning callout box
    add_rect(s, int(Inches(0.3)), int(Inches(6.05)),
             int(Inches(12.5)), int(Inches(0.75)), DARK)
    add_textbox(s, int(Inches(0.5)), int(Inches(6.12)),
                int(Inches(12.2)), int(Inches(0.6)),
                "WARNING: If you submit incorrect data — do not try to fix it by re-submitting different data. "
                "Contact your unit data manager immediately. Overwriting data incorrectly can corrupt records "
                "upstream.",
                10, False, GOLD, PP_ALIGN.LEFT)

    add_footer(s)

    # ── Slide 7: Data Classification and Handling ─────────────────────────────
    s = blank_slide(prs)
    add_classif_banners(s)
    slide_title(s, "DATA CLASSIFICATION AND HANDLING")
    subtitle_text(s, "MSS operates across classification levels. Know the rules before you act.")

    # Three columns: Can Do / Can't Do / If Unsure
    col_w = int(Inches(3.7))
    col_h = int(Inches(5.0))
    panel_t = int(Inches(1.3))

    for i, (label, fill, lines) in enumerate([
        ("YOU CAN", GREEN, [
            "▸  View data within your access level",
            "▸  Filter, sort, and export to approved formats",
            "   when your unit's policy allows",
            "▸  Brief information at the classification level",
            "   it was accessed on",
            "▸  Screenshot an UNCLASSIFIED dashboard for use",
            "   in UNCLASSIFIED products",
            "▸  Share a link to a dashboard with others",
            "   who have equivalent access",
        ]),
        ("YOU CANNOT", RED, [
            "▸  Copy classified data to an unclassified system",
            "▸  Email screenshots of classified dashboards",
            "▸  Share your login with anyone — ever",
            "▸  Download raw data to a personal device",
            "▸  Enter data into a network above your clearance",
            "▸  Discuss classified dashboard details on",
            "   unclassified channels",
        ]),
        ("IF UNSURE — STOP", GOLD, [
            "Classification handling is not a judgment call.",
            "",
            "If you are not certain whether an action is authorized:",
            "",
            "▸  Do not proceed",
            "▸  Contact your unit S6 or data manager",
            "▸  Consult your unit's data handling SOP",
            "",
            "Violations are reportable security incidents.",
            "When in doubt — ask first.",
        ]),
    ]):
        col_l = int(Inches(0.3)) + i * (col_w + int(Inches(0.3)))
        header_bar(s, col_l, panel_t, col_w, int(Inches(0.45)), label, fill)
        content_panel(s, col_l, panel_t + int(Inches(0.45)), col_w, col_h - int(Inches(0.45)))
        add_multi(s, col_l + int(Inches(0.15)), panel_t + int(Inches(0.62)),
                  col_w - int(Inches(0.25)), col_h - int(Inches(0.65)),
                  lines, 10.5, DARK)

    add_footer(s)

    # ── Slide 8: When Something Looks Wrong ───────────────────────────────────
    s = blank_slide(prs)
    add_classif_banners(s)
    slide_title(s, "WHEN SOMETHING LOOKS WRONG")
    subtitle_text(s, "You are not expected to fix data problems. You are expected to report them.")

    # Left: What to report | Right: Who to call
    left_w  = int(Inches(7.5))
    right_w = int(Inches(5.0))
    panel_t = int(Inches(1.3))
    panel_h = int(Inches(4.8))

    header_bar(s, int(Inches(0.3)), panel_t, left_w, int(Inches(0.45)),
               "WHAT 'SOMETHING WRONG' LOOKS LIKE", NAVY)
    content_panel(s, int(Inches(0.3)), panel_t + int(Inches(0.45)), left_w, panel_h - int(Inches(0.45)))
    add_multi(s, int(Inches(0.48)), panel_t + int(Inches(0.62)),
              left_w - int(Inches(0.3)), panel_h - int(Inches(0.75)),
              [
                  "Data accuracy problems:",
                  "▸  A soldier appears in the wrong unit",
                  "▸  A vehicle shows FMC but you know it's NMC",
                  "▸  A record is missing that should be there",
                  "▸  Numbers don't match what your section submitted",
                  "",
                  "Platform problems:",
                  "▸  Page doesn't load or shows an error",
                  "▸  Dashboard is blank where it had data yesterday",
                  "▸  A form won't submit and the error message doesn't help",
                  "▸  You're seeing data you don't believe you should have access to",
              ], 10.5, DARK)

    header_bar(s, int(Inches(8.1)), panel_t, right_w, int(Inches(0.45)),
               "WHO TO CONTACT", TEAL)
    content_panel(s, int(Inches(8.1)), panel_t + int(Inches(0.45)), right_w, panel_h - int(Inches(0.45)))
    add_multi(s, int(Inches(8.28)), panel_t + int(Inches(0.62)),
              right_w - int(Inches(0.25)), panel_h - int(Inches(0.75)),
              [
                  "Data accuracy issue:",
                  "→ Your unit data manager",
                  "   (S6 data shop or MSS POC)",
                  "",
                  "Platform / technical issue:",
                  "→ Your unit S6",
                  "   or MSS help desk",
                  "",
                  "Security / access issue:",
                  "→ Your unit S2 or S6",
                  "   Do NOT wait. Report immediately.",
                  "",
                  "NOT SURE WHICH?",
                  "→ Start with your unit data manager.",
              ], 10.5, DARK)

    # Bottom callout
    add_rect(s, int(Inches(0.3)), int(Inches(6.35)),
             int(Inches(12.5)), int(Inches(0.5)), DARK)
    add_textbox(s, int(Inches(0.5)), int(Inches(6.4)),
                int(Inches(12.2)), int(Inches(0.42)),
                "BLUF: Reporting a data problem is a contribution — not a complaint. "
                "Every error you flag makes the commander's picture more accurate.",
                10, False, GOLD, PP_ALIGN.LEFT)

    add_footer(s)

    # ── Slide 9: What You Cannot Do — And Who Can ─────────────────────────────
    s = blank_slide(prs)
    add_classif_banners(s)
    slide_title(s, "WHAT YOU CANNOT DO — AND WHO CAN")
    subtitle_text(s, "MSS has a clear boundary between users and builders. Know which side you're on.")

    # Two columns: User (you) | Builder
    col_w = int(Inches(5.9))
    panel_t = int(Inches(1.3))
    panel_h = int(Inches(4.7))

    header_bar(s, int(Inches(0.3)), panel_t, col_w, int(Inches(0.45)),
               "WHAT A TM-10 USER DOES", NAVY)
    content_panel(s, int(Inches(0.3)), panel_t + int(Inches(0.45)), col_w, panel_h - int(Inches(0.45)))
    add_multi(s, int(Inches(0.48)), panel_t + int(Inches(0.62)),
              col_w - int(Inches(0.3)), panel_h - int(Inches(0.75)),
              [
                  "▸  Navigates to existing dashboards and apps",
                  "▸  Reads and interprets data",
                  "▸  Filters, sorts, and exports within authorized limits",
                  "▸  Submits data through existing forms and actions",
                  "▸  Uses the AI assistant for data queries",
                  "▸  Reports problems to the appropriate contact",
                  "",
                  "You do not need to understand how the data gets there.",
                  "You need to know how to use it accurately.",
              ], 10.5, DARK)

    header_bar(s, int(Inches(7.05)), panel_t, col_w, int(Inches(0.45)),
               "WHAT A BUILDER OR DEVELOPER DOES", BLUE)
    content_panel(s, int(Inches(7.05)), panel_t + int(Inches(0.45)), col_w, panel_h - int(Inches(0.45)))
    add_multi(s, int(Inches(7.22)), panel_t + int(Inches(0.62)),
              col_w - int(Inches(0.3)), panel_h - int(Inches(0.75)),
              [
                  "▸  Builds the dashboards and apps you use",
                  "▸  Configures data pipelines that feed the platform",
                  "▸  Creates object types and data models",
                  "▸  Manages ontology and data governance",
                  "▸  Grants and revokes access",
                  "▸  Troubleshoots technical platform issues",
                  "",
                  "If you need a new view, a new report, or a new app —",
                  "submit a request to your unit data manager.",
                  "Do not try to build it yourself without authorization.",
              ], 10.5, DARK)

    # Divider label
    add_textbox(s, int(Inches(5.95)), panel_t + int(Inches(1.8)),
                int(Inches(1.1)), int(Inches(1.2)),
                "←  USER\n\nBUILDER  →", 9, True, NAVY, PP_ALIGN.CENTER)

    add_footer(s)

    # ── Slide 10: Your Role ───────────────────────────────────────────────────
    s = blank_slide(prs)
    add_classif_banners(s)

    # Full dark background
    add_rect(s, 0, int(Inches(0.6)), int(SLIDE_W), int(Inches(6.6)), NAVY)

    add_textbox(s, int(Inches(1.0)), int(Inches(1.1)), int(Inches(11)), int(Inches(0.8)),
                "YOUR ROLE", 36, True, WHITE, PP_ALIGN.CENTER)

    add_textbox(s, int(Inches(1.5)), int(Inches(2.0)), int(Inches(10)), int(Inches(0.7)),
                "You are the last link in the data chain — and the most important one.",
                20, False, GOLD, PP_ALIGN.CENTER)

    add_multi(s, int(Inches(1.8)), int(Inches(2.95)), int(Inches(9.5)), int(Inches(2.5)),
              [
                  "The Army has built the infrastructure. The pipelines are running.",
                  "The dashboards are live. The AI is reading real data.",
                  "",
                  "None of it matters if the information doesn't reach the commander accurately.",
                  "",
                  "▸  Enter data correctly — someone is briefing it to leadership.",
                  "▸  Read dashboards critically — question what doesn't look right.",
                  "▸  Use AIP as a tool — not as a substitute for your judgment.",
                  "▸  Report problems — your observation keeps the picture clean.",
              ], 13, LGRAY, False, PP_ALIGN.LEFT)

    add_multi(s, int(Inches(1.8)), int(Inches(5.55)), int(Inches(9.5)), int(Inches(0.6)),
              ["TM-10 Complete  ·  Next: TM-20 Builder Course"], 11, GOLD)

    add_footer(s)

    out_path = OUT_DIR / "mss_platform_overview.pptx"
    prs.save(str(out_path))
    print(f"Saved: {out_path}")
    return out_path


# ─────────────────────────────────────────────────────────────────────────────
# DECK 2: Patch Army Data Orientation slides 12-14
# ─────────────────────────────────────────────────────────────────────────────

def patch_army_data_orientation():
    """
    Opens Army_Data_Orientation_PROTO.pptx and rewrites slides 12-14
    for a TM-10 end-user audience, saves as army_data_orientation_v2.pptx.

    Changes:
      Slide 12 (Where You Fit In): Move "YOU" marker to Applications/Decisions layer.
      Slide 13 (What You'll Learn): Replace builder scope with end-user task scope.
      Slide 14 (Your Mission): Replace "build the architecture" with consumer mission.
    """
    prs = Presentation(str(SOURCE_PPTX))

    # Slide indices are 0-based; slides 12-14 → indices 11-13
    # We'll delete and replace slides 12, 13, 14 (indices 11-13).
    # Easiest: remove shapes from those slides and rebuild content in place.

    def clear_slide(slide):
        """Remove all non-placeholder shapes."""
        sp_tree = slide.shapes._spTree
        to_remove = []
        for sp in sp_tree:
            if sp.tag.endswith('}sp') or sp.tag.endswith('}grpSp'):
                name = sp.find('.//{http://schemas.openxmlformats.org/drawingml/2006/spreadsheetDrawing}nvPr')
                to_remove.append(sp)
        for sp in list(sp_tree):
            tag = sp.tag.split('}')[-1]
            if tag in ('sp', 'pic', 'cxnSp', 'grpSp'):
                sp_tree.remove(sp)

    # ── Slide 12: WHERE YOU FIT IN (index 11) ────────────────────────────────
    s12 = prs.slides[11]
    clear_slide(s12)

    add_classif_banners(s12)
    slide_title(s12, "WHERE YOU FIT IN")
    subtitle_text(s12, "The enterprise platforms are in place. Your job is to use them — accurately.")

    # Vertical stack: SOURCE → INTEGRATION → ONTOLOGY → APPLICATIONS (YOU) → DECISIONS (YOU)
    stack = [
        ("SOURCE SYSTEMS",   BLUE,  "ERPs, logistics, HR, readiness feeds — the raw data sources"),
        ("DATA INTEGRATION", BLUE,  "Pipelines that clean, transform, and connect data — built by data engineers"),
        ("ONTOLOGY",         BLUE,  "Object types, link types, properties — the Army's data model"),
        ("APPLICATIONS",     TEAL,  "Dashboards, workflows, decision-support tools — THIS IS WHERE YOU WORK"),
        ("DECISIONS",        NAVY,  "Commanders act on real-time, trusted data — THIS IS WHAT YOU ENABLE"),
    ]

    box_h = int(Inches(0.72))
    box_w = int(Inches(7.0))
    box_l = int(Inches(0.8))
    start_t = int(Inches(1.3))
    gap = int(Inches(0.12))

    for i, (label, fill, desc) in enumerate(stack):
        t = start_t + i * (box_h + gap)
        is_you = i >= 3
        actual_fill = fill if is_you else LGRAY
        label_fill  = fill if is_you else BLUE
        add_rect(s12, box_l, t, box_w, box_h, actual_fill)
        add_textbox(s12, box_l + int(Inches(0.15)), t + int(Inches(0.08)),
                    int(Inches(3.0)), int(Inches(0.38)),
                    label, 12, True,
                    WHITE if is_you else BLUE, PP_ALIGN.LEFT)
        add_textbox(s12, box_l + int(Inches(0.15)), t + int(Inches(0.38)),
                    box_w - int(Inches(0.3)), int(Inches(0.28)),
                    desc, 9.5, False,
                    WHITE if is_you else BLUE, PP_ALIGN.LEFT)
        # Downward arrow between boxes
        if i < len(stack) - 1:
            add_textbox(s12, box_l + int(Inches(3.2)), t + box_h,
                        int(Inches(0.5)), gap + int(Inches(0.05)),
                        "↓", 14, False, BLUE, PP_ALIGN.CENTER)

    # Right callout
    add_rect(s12, int(Inches(8.2)), int(Inches(3.8)),
             int(Inches(4.5)), int(Inches(2.1)), DARK)
    add_multi(s12, int(Inches(8.38)), int(Inches(3.95)),
              int(Inches(4.2)), int(Inches(1.85)),
              [
                  "YOU OPERATE HERE",
                  "",
                  "Applications surface the data that builders have",
                  "structured and engineers have pipelined.",
                  "",
                  "Your job: use it accurately, report what's wrong,",
                  "and make sure the commander's picture is right.",
              ], 10.5, GOLD if True else WHITE)

    add_footer(s12)

    # ── Slide 13: WHAT YOU'LL DO IN MSS (index 12) ────────────────────────────
    s13 = prs.slides[12]
    clear_slide(s13)

    add_classif_banners(s13)
    slide_title(s13, "WHAT YOU'LL DO IN MSS")
    subtitle_text(s13, "These are the tasks you will perform as an MSS user.")

    # Three columns of tasks
    col_w = int(Inches(3.9))
    col_h = int(Inches(5.1))
    panel_t = int(Inches(1.3))
    gaps = int(Inches(0.28))

    for i, (label, fill, lines) in enumerate([
        ("FIND & READ INFORMATION", NAVY, [
            "Navigate to your unit's applications",
            "Read readiness and status dashboards",
            "Search for specific soldiers, vehicles, or events",
            "Filter and sort data to find what you need",
            "Understand what the data is telling you",
        ]),
        ("SUBMIT & RECORD DATA", TEAL, [
            "Enter status updates through authorized forms",
            "Submit action requests through workflows",
            "Record events and observations accurately",
            "Verify your entry appeared correctly",
            "Never guess — if you don't know the value, find out",
        ]),
        ("USE AI & REPORT PROBLEMS", BLUE, [
            "Use the AI assistant (AIP) to query unit data",
            "Verify AI answers before briefing them",
            "Identify and report data accuracy problems",
            "Contact your unit data manager with issues",
            "Follow your unit's data handling and security SOP",
        ]),
    ]):
        col_l = int(Inches(0.3)) + i * (col_w + gaps)
        header_bar(s13, col_l, panel_t, col_w, int(Inches(0.45)), label, fill)
        content_panel(s13, col_l, panel_t + int(Inches(0.45)), col_w, col_h - int(Inches(0.45)))
        for j, line in enumerate(lines):
            add_textbox(s13,
                        col_l + int(Inches(0.15)),
                        panel_t + int(Inches(0.62)) + j * int(Inches(0.75)),
                        col_w - int(Inches(0.25)), int(Inches(0.7)),
                        f"▸  {line}", 10.5, False, DARK, PP_ALIGN.LEFT)

    add_footer(s13)

    # ── Slide 14: YOUR MISSION (index 13) ────────────────────────────────────
    s14 = prs.slides[13]
    clear_slide(s14)

    add_classif_banners(s14)
    add_rect(s14, 0, int(Inches(0.6)), int(SLIDE_W), int(Inches(6.6)), NAVY)

    add_textbox(s14, int(Inches(1.0)), int(Inches(1.15)), int(Inches(11)), int(Inches(0.8)),
                "YOUR MISSION", 36, True, WHITE, PP_ALIGN.CENTER)

    add_textbox(s14, int(Inches(1.5)), int(Inches(2.05)), int(Inches(10)), int(Inches(0.65)),
                "Use the data architecture that enables enterprise capability — accurately.",
                20, False, GOLD, PP_ALIGN.CENTER)

    checks = [
        "The Army has built live data infrastructure. Your job is to use it — not to build it.",
        "Every form you submit, every status you enter, feeds a commander's real-time picture.",
        "Accurate data in means accurate decisions out. That is the mission.",
    ]
    for i, text in enumerate(checks):
        ty = int(Inches(3.0)) + i * int(Inches(0.9))
        add_rect(s14, int(Inches(1.2)), ty, int(Inches(0.55)), int(Inches(0.55)), GREEN)
        add_textbox(s14, int(Inches(1.4)), ty + int(Inches(0.07)),
                    int(Inches(0.3)), int(Inches(0.4)), "✓", 16, True, WHITE, PP_ALIGN.CENTER)
        add_textbox(s14, int(Inches(1.9)), ty + int(Inches(0.07)),
                    int(Inches(9.5)), int(Inches(0.5)),
                    text, 13, False, WHITE, PP_ALIGN.LEFT)

    add_footer(s14)

    out_path = OUT_DIR / "army_data_orientation_v2.pptx"
    prs.save(str(out_path))
    print(f"Saved: {out_path}")
    return out_path


# ─────────────────────────────────────────────────────────────────────────────
if __name__ == "__main__":
    print("Building TM-10 decks...")
    p1 = build_mss_platform_overview()
    p2 = patch_army_data_orientation()
    print("\nDone.")
    print(f"  {p1}")
    print(f"  {p2}")
PYEOF