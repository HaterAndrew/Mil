"""
Build Maven Builder Training Force Rollout one-pager slide (.pptx).
Run: python3 scripts/build_rollout_slide.py
Output: maven_training/mss_info_app/Maven_Rollout_Plan.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Palette (matches mss_info_app navy/gold scheme) ─────────────────────────
NAVY_DARK  = RGBColor(0x07, 0x16, 0x28)   # --navy-dark  — header bar
NAVY       = RGBColor(0x0C, 0x23, 0x40)   # --navy       — table header row
NAVY_LIGHT = RGBColor(0x16, 0x3A, 0x6C)   # --navy-light — phase cell accent
GOLD       = RGBColor(0xC8, 0x97, 0x1A)   # --gold       — accent text / timeline
GOLD_LIGHT = RGBColor(0xE0, 0xB8, 0x40)   # --gold-light — highlight
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE  = RGBColor(0xF3, 0xF5, 0xFA)   # --off-white  — alternating row fill
LIGHT_GRAY = RGBColor(0xE0, 0xE4, 0xEF)   # --gray-100   — table borders
DARK_TEXT  = RGBColor(0x0A, 0x16, 0x28)   # --gray-900   — body text
GREEN_BANNER = RGBColor(0x1E, 0x6B, 0x2A) # banner green

# ── Slide dimensions (widescreen 16:9) ──────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

slide_layout = prs.slide_layouts[6]   # blank
slide = prs.slides.add_slide(slide_layout)

def add_rect(slide, l, t, w, h, fill_rgb=None, line_rgb=None, line_width_pt=0):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    shape.line.fill.background()
    if fill_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
    else:
        shape.fill.background()
    if line_rgb and line_width_pt:
        from pptx.util import Pt as Pt2
        shape.line.color.rgb = line_rgb
        shape.line.width = Pt2(line_width_pt)
    else:
        shape.line.fill.background()
    return shape

def add_textbox(slide, l, t, w, h, text, font_size=11, bold=False,
                color=DARK_TEXT, align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Arial"
    return txBox

# ── Background ───────────────────────────────────────────────────────────────
add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=OFF_WHITE)

# ── Header bar ───────────────────────────────────────────────────────────────
add_rect(slide, 0, 0, 13.33, 1.0, fill_rgb=NAVY_DARK)

# Title
add_textbox(slide, 0.3, 0.05, 9.0, 0.55,
            "MAVEN BUILDER TRAINING — FORCE ROLLOUT PLAN",
            font_size=20, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

# Subtitle / date
add_textbox(slide, 0.3, 0.53, 9.0, 0.28,
            "USAREUR-AF Operational Data Team  |  MAR 2026",
            font_size=9, bold=False, color=GOLD_LIGHT, align=PP_ALIGN.LEFT)

# Seal placeholder (right side of header)
add_rect(slide, 11.8, 0.22, 1.2, 0.72, fill_rgb=NAVY_LIGHT)
add_textbox(slide, 11.8, 0.30, 1.2, 0.55,
            "USAREUR-AF\nOD TEAM",
            font_size=7, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

# ── Table header row ─────────────────────────────────────────────────────────
col_x     = [0.15, 1.05, 4.55, 9.10]
col_w     = [0.85, 3.45, 4.50, 4.10]
row_h     = 0.44
header_y  = 1.10

hdr_labels = ["PHASE", "ACTION", "DECISION / OWNER", "SUCCESS STANDARD"]

for i, (lbl, x, w) in enumerate(zip(hdr_labels, col_x, col_w)):
    add_rect(slide, x, header_y, w, row_h, fill_rgb=NAVY)
    add_textbox(slide, x + 0.05, header_y + 0.06, w - 0.1, row_h - 0.08,
                lbl, font_size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ── Table rows ────────────────────────────────────────────────────────────────
rows = [
    (
        "1\nCOMMAND\nAUTHORITY",
        "Issue training directive (EXORD or policy memo) designating Maven Builder Training as required for identified MOSs/roles.",
        "CDR decision: mandatory vs. encouraged by role.\nG3/J3 staff action.",
        "Signed directive published to all subordinate units."
    ),
    (
        "2\nCADRE\nSTANDUP",
        "Identify and assign 5 cadre: 1 OIC + 4 SMEs covering TM-10/20/30/40 tracks. Develop cadre SOP and instructor certification standard.",
        "G6/G2/S6 nominations.\nOIC appointed by CDR.",
        "Cadre certified; SOP signed; instructors completed TM-40 series."
    ),
    (
        "3\nAUDIENCE\nSEQUENCING",
        "Phase training: TM-10 (all staff) → TM-20/30 (data-adjacent billets) → TM-40 series (specialists by MOS).",
        "G3 decision: phased by echelon or MOS priority.\nUnit S3s build into ATP.",
        "100% TM-10 completion; TM-40 track assignments published."
    ),
    (
        "4\nRESOURCE\n& ACCESS",
        "Provision Foundry seats, schedule lab time/VTC blocks, integrate into unit training calendars. Coordinate with G6 for network access.",
        "G6 action: license provisioning.\nG3: calendar deconfliction.",
        "All students have Foundry access NLT training start date."
    ),
    (
        "5\nTRACK\n& REPORT",
        "Establish completion tracking in DTMS/ATRRS. Set Go/No-Go checkpoints per MTP. Report completion status at monthly CDR training brief.",
        "G1/S1 action: reporting format.\nCDR sets completion deadline (recommended D+180).",
        "Monthly completion report to CDR; Go/No-Go recorded per student."
    ),
    (
        "6\nFEEDBACK\n& SUSTAIN",
        "Conduct AAR after first cohort; update curriculum and cadre SOP quarterly. Integrate into new-hire onboarding and unit OPORDs.",
        "Cadre OIC owns curriculum updates.\nCDR approves policy changes.",
        "Quarterly AAR complete; curriculum version controlled; onboarding updated."
    ),
]

row_fills = [OFF_WHITE, WHITE, OFF_WHITE, WHITE, OFF_WHITE, WHITE]
data_y = header_y + row_h
r_h = 0.80   # row height

for r_idx, (phase, action, decision, standard) in enumerate(rows):
    y = data_y + r_idx * r_h
    fill = row_fills[r_idx]

    # Row background
    add_rect(slide, 0.15, y, 12.95, r_h, fill_rgb=fill,
             line_rgb=LIGHT_GRAY, line_width_pt=0.5)

    # Phase cell (navy accent left strip)
    add_rect(slide, 0.15, y, 0.85, r_h, fill_rgb=NAVY_LIGHT)
    add_textbox(slide, 0.16, y + 0.04, 0.83, r_h - 0.08,
                phase, font_size=7.5, bold=True, color=GOLD,
                align=PP_ALIGN.CENTER)

    # Action cell
    add_textbox(slide, 1.10, y + 0.05, 3.38, r_h - 0.10,
                action, font_size=8, bold=False, color=DARK_TEXT)

    # Decision/Owner cell
    add_textbox(slide, 4.60, y + 0.05, 4.38, r_h - 0.10,
                decision, font_size=8, bold=False, color=DARK_TEXT)

    # Standard cell
    add_textbox(slide, 9.15, y + 0.05, 3.90, r_h - 0.10,
                standard, font_size=8, bold=False, color=DARK_TEXT)

# ── Timeline footer bar ───────────────────────────────────────────────────────
footer_y = data_y + len(rows) * r_h + 0.05
add_rect(slide, 0.15, footer_y, 12.95, 0.30, fill_rgb=NAVY_DARK)
add_textbox(slide, 0.20, footer_y + 0.04, 12.80, 0.24,
            "RECOMMENDED TIMELINE:  Directive signed NLT D+30  |  Cadre designated D+45  |  First cohort begins D+90  |  Full force completion D+180",
            font_size=8, bold=True, color=GOLD, align=PP_ALIGN.CENTER)



# ── Save ──────────────────────────────────────────────────────────────────────
out_path = "maven_training/mss_info_app/Maven_Rollout_Plan.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
