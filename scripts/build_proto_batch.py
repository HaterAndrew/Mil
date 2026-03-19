#!/usr/bin/env python3
"""
Batch prototype builder: 3 training decks reformatted to USAREUR-AF template.
Run from repo root: python3 scripts/build_proto_batch.py
Outputs:
  maven_training/pdf/MSS_Training_Progression_PROTO.pptx
  maven_training/pdf/MSS_Project_Overview_PROTO.pptx
  maven_training/pdf/MSS_PROGRAM_OVERVIEW_PROTO.pptx
"""

import os
import shutil
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml

# ── Paths ──────────────────────────────────────────────────────────────────────
# Template path: set via PPTX_TEMPLATE env var, or fall back to default location.
TEMPLATE = Path(os.environ.get(
    "PPTX_TEMPLATE",
    Path(__file__).parent.parent / "maven_training" / "source_material" / "USAREUR-AF PPT Template.pptx"
))
OUT_DIR  = Path("maven_training/pdf")

# ── Palette ────────────────────────────────────────────────────────────────────
NAVY      = RGBColor(0x17, 0x32, 0x5C)   # primary header
NAVY2     = RGBColor(0x18, 0x33, 0x5F)   # slight variant
NAVY_MID  = RGBColor(0x26, 0x4D, 0x7E)   # lighter band (audience rows)
NAVY_LT   = RGBColor(0x3A, 0x64, 0x9E)   # prereq / secondary bands
DARK_BG   = RGBColor(0x0D, 0x1F, 0x3C)   # BLUF / timeline bars
PURPLE    = RGBColor(0x70, 0x30, 0xA0)   # classification text
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE = RGBColor(0xF5, 0xF6, 0xF8)
GRAY_LT   = RGBColor(0xE2, 0xE8, 0xF0)
GRAY_MID  = RGBColor(0xC8, 0xD2, 0xDE)
GOLD      = RGBColor(0xC8, 0x9A, 0x00)   # stat highlights

CLASSIFICATION = ""


# ── Low-level helpers ──────────────────────────────────────────────────────────

def solid_fill_xml(rgb: RGBColor) -> str:
    ns = 'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"'
    return f'<a:solidFill {ns}><a:srgbClr val="{str(rgb)}"/></a:solidFill>'


def set_shape_fill(shape, rgb: RGBColor):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb


def set_shape_no_line(shape):
    shape.line.fill.background()


def set_cell_bg(cell, rgb: RGBColor):
    """Solid fill for a table cell."""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for tag in ("a:solidFill", "a:noFill", "a:gradFill"):
        el = tcPr.find(qn(tag))
        if el is not None:
            tcPr.remove(el)
    tcPr.append(parse_xml(solid_fill_xml(rgb)))


def set_textbox_fill(tb, rgb: RGBColor):
    """Solid fill on a text box's spPr."""
    spPr = tb._element.spPr
    # Remove existing fills
    for tag in (qn("a:solidFill"), qn("a:noFill"), qn("a:gradFill")):
        el = spPr.find(tag)
        if el is not None:
            spPr.remove(el)
    spPr.append(parse_xml(solid_fill_xml(rgb)))


# ── Mid-level drawing helpers ──────────────────────────────────────────────────

def rect(slide, x, y, w, h, fill: RGBColor, no_line=True):
    """Add a solid rectangle, return shape."""
    shp = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    set_shape_fill(shp, fill)
    if no_line:
        set_shape_no_line(shp)
    shp.text_frame.clear()   # ensure no stray text
    return shp


def text_box(slide, x, y, w, h, lines, *, font_size=9, bold=False,
             color: RGBColor = WHITE, align=PP_ALIGN.LEFT, wrap=True,
             fill: RGBColor = None, v_anchor=None):
    """Add a text box with one or more lines (list[str] or str split on \\n).
    Returns the text box shape.
    """
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    if fill:
        set_textbox_fill(tb, fill)
    tf = tb.text_frame
    tf.word_wrap = wrap
    if v_anchor:
        tf.vertical_anchor = v_anchor
    if isinstance(lines, str):
        lines = lines.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb


def clear_and_set(slide_shape, text, font_size=16, bold=True,
                  color: RGBColor = NAVY, align=PP_ALIGN.LEFT):
    """Replace all text in an existing slide shape."""
    tf = slide_shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color


def remove_shape(shape):
    el = shape._element
    el.getparent().remove(el)


def add_bottom_banner(slide, text=CLASSIFICATION):
    """Add classification banner at bottom of content slide (mirrors top)."""
    tb = slide.shapes.add_textbox(
        Inches(0), Inches(7.25), Inches(13.33), Inches(0.25)
    )
    tf = tb.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(9)
    run.font.bold = True
    run.font.color.rgb = PURPLE


def prep_template(output_path: Path) -> Presentation:
    """Copy template and return open Presentation."""
    shutil.copy(TEMPLATE, output_path)
    return Presentation(str(output_path))


def update_template_chrome(slide, header_text: str, classification=CLASSIFICATION):
    """Update the standard template chrome on a content slide."""
    to_remove = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        txt  = shape.text_frame.text.strip()
        name = shape.name
        if name == "TextBox 2" or txt == "Header":
            clear_and_set(shape, header_text, font_size=15, bold=True,
                          color=NAVY, align=PP_ALIGN.LEFT)
        elif txt in ("CUI", "UNCLASSIFIED", "Classification"):
            for p in shape.text_frame.paragraphs:
                for run in p.runs:
                    run.text = classification
        elif name == "TextBox 8":
            to_remove.append(shape)
    for s in to_remove:
        remove_shape(s)
    # Add bottom banner (slide 2 only has top in template)
    add_bottom_banner(slide, classification)


def update_title_chrome(slide, title_lines: list[str], poc: str,
                        classification=CLASSIFICATION):
    """Update the template title slide."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        txt  = shape.text_frame.text.strip()
        name = shape.name
        if name == "Title 1":
            tf = shape.text_frame
            tf.clear()
            for i, line in enumerate(title_lines):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.alignment = PP_ALIGN.CENTER
                run = p.add_run()
                run.text = line
                run.font.size = Pt(26 if i == 0 else 18)
                run.font.bold = (i == 0)
        elif name.startswith("hlSlideMaster") and txt in ("Classification", "CUI", "UNCLASSIFIED"):
            for p in shape.text_frame.paragraphs:
                for run in p.runs:
                    run.text = classification
        elif name == "TextBox 3" and txt == "Classification":
            for p in shape.text_frame.paragraphs:
                for run in p.runs:
                    run.text = classification
        elif name == "TextBox 5" and txt == "POC":
            tf = shape.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = poc
            run.font.size = Pt(11)
            run.font.color.rgb = NAVY


# ══════════════════════════════════════════════════════════════════════════════
# PROTOTYPE 1 — MSS Training Progression
# ══════════════════════════════════════════════════════════════════════════════

PROGRESSION_CARDS = [
    {
        "tm":       "TM-10",
        "title":    "MAVEN USER",
        "audience": "ALL PERSONNEL",
        "target":   "Every Soldier, Officer & Civilian",
        "prereq":   "None",
        "bullets":  [
            "Navigate the MSS portal",
            "Access datasets & dashboards",
            "Run saved queries / filters",
            "Export and share products",
            "Understand data provenance",
        ],
    },
    {
        "tm":       "TM-20",
        "title":    "BUILDER",
        "audience": "ALL STAFF",
        "target":   "Light Builders  (no coding required)",
        "prereq":   "Prereq: TM-10",
        "bullets":  [
            "Build Workshop dashboards",
            "Create no-code transforms",
            "Design basic data pipelines",
            "Configure object views",
            "Publish products to unit portals",
        ],
    },
    {
        "tm":       "TM-30",
        "title":    "ADVANCED BUILDER",
        "audience": "DATA-ADJACENT SPECIALISTS",
        "target":   "17/25-series · S6/G6 · G2",
        "prereq":   "Prereq: TM-20",
        "bullets":  [
            "Author Python/SQL transforms",
            "Model ontology object types",
            "Build AIP Logic workflows",
            "Performance-tune pipelines",
            "Govern datasets & lineage",
        ],
    },
]


def build_progression_content(slide):
    """3-tier card layout + specialist panel."""
    update_template_chrome(slide, "MSS TRAINING PROGRESSION")

    # Layout constants
    Y0     = 0.75   # card top
    H_CARD = 5.85   # card height → bottom at 6.60"
    W_CARD = 3.00
    W_ARR  = 0.26   # arrow gutter
    W_SPEC = 3.17
    X_START = 0.20

    card_x = [
        X_START,                        # TM-10
        X_START + W_CARD + W_ARR,       # TM-20
        X_START + 2*(W_CARD + W_ARR),   # TM-30
    ]
    x_spec = X_START + 3*W_CARD + 3*W_ARR  # specialist panel

    # ── Draw the 3 tier cards ────────────────────────────────────────────────
    for i, card in enumerate(PROGRESSION_CARDS):
        x = card_x[i]

        # Drop shadow
        rect(slide, x + 0.04, Y0 + 0.04, W_CARD, H_CARD, DARK_BG)

        # Card body (off-white)
        rect(slide, x, Y0, W_CARD, H_CARD, OFF_WHITE)

        # Navy header band (TM number + course name)
        rect(slide, x, Y0, W_CARD, 0.88, NAVY)
        text_box(slide, x + 0.08, Y0 + 0.04, W_CARD - 0.16, 0.42,
                 card["tm"], font_size=26, bold=True, color=WHITE,
                 align=PP_ALIGN.LEFT, wrap=False)
        text_box(slide, x + 0.08, Y0 + 0.48, W_CARD - 0.16, 0.32,
                 card["title"], font_size=12, bold=False, color=GRAY_LT,
                 align=PP_ALIGN.LEFT, wrap=False)

        # Audience band
        rect(slide, x, Y0 + 0.88, W_CARD, 0.52, NAVY_MID)
        text_box(slide, x + 0.08, Y0 + 0.90, W_CARD - 0.16, 0.24,
                 card["audience"], font_size=8, bold=True, color=WHITE,
                 align=PP_ALIGN.LEFT, wrap=False)
        text_box(slide, x + 0.08, Y0 + 1.13, W_CARD - 0.16, 0.22,
                 card["target"], font_size=7.5, bold=False, color=GRAY_LT,
                 align=PP_ALIGN.LEFT, wrap=False)

        # Prereq line
        rect(slide, x, Y0 + 1.40, W_CARD, 0.24, NAVY_LT)
        text_box(slide, x + 0.08, Y0 + 1.42, W_CARD - 0.16, 0.20,
                 card["prereq"], font_size=7.5, bold=False, color=GRAY_LT,
                 align=PP_ALIGN.LEFT, wrap=False)

        # Bullet list
        bullet_y = Y0 + 1.68
        for bullet in card["bullets"]:
            text_box(slide, x + 0.10, bullet_y, W_CARD - 0.20, 0.30,
                     f"▸  {bullet}", font_size=8, color=NAVY,
                     align=PP_ALIGN.LEFT)
            bullet_y += 0.70

    # ── Arrows between cards ─────────────────────────────────────────────────
    for i in range(2):
        ax = card_x[i] + W_CARD + 0.04
        text_box(slide, ax, Y0 + H_CARD/2 - 0.18, W_ARR - 0.04, 0.36,
                 "▶", font_size=18, bold=True, color=NAVY_MID,
                 align=PP_ALIGN.CENTER, wrap=False)

    # ── Specialist panel ─────────────────────────────────────────────────────
    x = x_spec

    rect(slide, x + 0.04, Y0 + 0.04, W_SPEC, H_CARD, DARK_BG)
    rect(slide, x, Y0, W_SPEC, H_CARD, OFF_WHITE)
    rect(slide, x, Y0, W_SPEC, 0.52, DARK_BG)
    text_box(slide, x + 0.08, Y0 + 0.08, W_SPEC - 0.16, 0.38,
             "SPECIALIST TRACKS", font_size=11, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, wrap=False)

    # TM-40 section
    rect(slide, x + 0.06, Y0 + 0.56, W_SPEC - 0.12, 0.24, NAVY)
    text_box(slide, x + 0.10, Y0 + 0.58, W_SPEC - 0.20, 0.20,
             "TM-40 SERIES  (Prereq: TM-30)", font_size=8, bold=True,
             color=WHITE, wrap=False)

    wff_tracks = [
        ("TM-40A", "Intelligence"),
        ("TM-40B", "Fires"),
        ("TM-40C", "Movement & Maneuver"),
        ("TM-40D", "Sustainment"),
        ("TM-40E", "Protection"),
        ("TM-40F", "Mission Command"),
    ]
    spec_tracks = [
        ("TM-40G", "ORSA"),
        ("TM-40H", "AI Engineer"),
        ("TM-40M", "ML Engineer"),
        ("TM-40J", "Program Manager"),
        ("TM-40K", "Knowledge Manager"),
        ("TM-40L", "Software Engineer"),
    ]
    all_tracks = wff_tracks + spec_tracks
    ty = Y0 + 0.84
    for track_id, track_name in all_tracks:
        text_box(slide, x + 0.10, ty, 0.72, 0.26,
                 track_id, font_size=7.5, bold=True, color=NAVY, wrap=False)
        text_box(slide, x + 0.84, ty, W_SPEC - 0.94, 0.26,
                 track_name, font_size=7.5, color=NAVY, wrap=False)
        ty += 0.26

    # Divider
    rect(slide, x + 0.06, ty + 0.02, W_SPEC - 0.12, 0.02, NAVY_MID)

    # TM-50 section
    rect(slide, x + 0.06, ty + 0.08, W_SPEC - 0.12, 0.24, NAVY)
    text_box(slide, x + 0.10, ty + 0.10, W_SPEC - 0.20, 0.20,
             "TM-50 SERIES  (Prereq: TM-40)", font_size=8, bold=True,
             color=WHITE, wrap=False)
    text_box(slide, x + 0.10, ty + 0.36, W_SPEC - 0.20, 0.48,
             "Advanced / expert continuation for\nTM-40G–M tracks only.",
             font_size=7.5, color=NAVY)

    # Footer
    text_box(slide, 0.10, 6.63, 13.13, 0.20,
             "USAREUR-AF  ·  C2DAO  ·  MSS-POI-001  ·  MAR 2026",
             font_size=7, color=NAVY2, align=PP_ALIGN.CENTER, wrap=False)


def build_progression():
    out = OUT_DIR / "MSS_Training_Progression_PROTO.pptx"
    prs = prep_template(out)
    update_title_chrome(prs.slides[0],
                        ["MSS TRAINING PROGRESSION",
                         "Maven Smart System  —  Core Learning Path"],
                        "USAREUR-AF Operational Data Team  |  C2DAO  |  MAR 2026")
    build_progression_content(prs.slides[1])
    prs.save(str(out))
    print(f"Saved → {out}")


# ══════════════════════════════════════════════════════════════════════════════
# PROTOTYPE 2 — MSS Project Overview
# ══════════════════════════════════════════════════════════════════════════════

BASELINE_TIERS = [
    {
        "tm": "TM-10",
        "label": "ALL PERSONNEL",
        "sub":   "Every Soldier & Civilian",
        "why":   "WHY CRITICAL:  No access = no adoption.",
        "body": (
            "CAC login, project navigation, Workshop apps, "
            "authorized Actions, ContainerShip data views, "
            "export, and data provenance basics."
        ),
    },
    {
        "tm": "TM-20",
        "label": "ALL STAFF",
        "sub":   "Light Builders — no coding",
        "why":   "WHY CRITICAL:  Units must own their data products.",
        "body": (
            "Pipeline Builder (visual, no-code), Workshop dashboards, "
            "forms, Object Type views, branch publishing."
        ),
    },
    {
        "tm": "TM-30",
        "label": "DATA-ADJACENT SPECIALISTS",
        "sub":   "17/25-series · S6 · G2",
        "why":   "WHY CRITICAL:  Maven's power is in complex, connected products.",
        "body": (
            "Multi-source pipelines, complex Workshop apps, "
            "Ontology architecture (Object Types, Links), "
            "Python/SQL transforms, AIP Logic."
        ),
    },
]

STATS = [
    ("3",      "Mandatory\nBaseline TMs"),
    ("100%",   "MOS\nCoverage"),
    ("88 hrs", "TM-10→30\nInstruction"),
    ("3-tier", "Progression\nNo-Code→Code"),
    ("18",     "Courses\nPublished"),
    ("160+",   "PDFs\nPublished"),
]


def build_overview_content(slide):
    """BLUF + 3-column baseline tiers + specialist panels + stats."""
    update_template_chrome(slide,
                           "MAVEN SMART SYSTEM — COMMAND-WIDE ADOPTION STRATEGY")

    Y0 = 0.76  # content start (below header)

    # ── BLUF bar ──────────────────────────────────────────────────────────────
    rect(slide, 0.10, Y0, 13.13, 0.34, DARK_BG)
    text_box(slide, 0.20, Y0 + 0.03, 13.00, 0.28,
             "BLUF:  Maven adoption fails at the unit level when personnel cannot "
             "access, navigate, and build within the system. This training infrastructure "
             "is the prerequisite for command-wide adoption.",
             font_size=8, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, wrap=True)

    Y1 = Y0 + 0.38   # below BLUF

    # Section labels
    text_box(slide, 0.10, Y1, 7.60, 0.22,
             "THE BASELINE TIERS  (mandatory for all personnel)",
             font_size=8, bold=True, color=NAVY, wrap=False)
    text_box(slide, 7.84, Y1, 5.29, 0.22,
             "SPECIALIST / ADVANCED",
             font_size=8, bold=True, color=NAVY, wrap=False)

    # Vertical divider
    rect(slide, 7.76, Y1, 0.04, 4.80, GRAY_MID)

    Y2 = Y1 + 0.26   # columns start

    # ── 3 Baseline columns ────────────────────────────────────────────────────
    col_w = 2.48
    col_gap = 0.02
    for j, tier in enumerate(BASELINE_TIERS):
        cx = 0.10 + j * (col_w + col_gap)

        # Header band
        rect(slide, cx, Y2, col_w, 0.60, NAVY)
        text_box(slide, cx + 0.06, Y2 + 0.04, col_w - 0.12, 0.34,
                 tier["tm"], font_size=22, bold=True, color=WHITE,
                 align=PP_ALIGN.LEFT, wrap=False)
        text_box(slide, cx + 0.06, Y2 + 0.38, col_w - 0.12, 0.20,
                 tier["sub"], font_size=7, color=GRAY_LT,
                 align=PP_ALIGN.LEFT, wrap=False)

        # Audience band
        rect(slide, cx, Y2 + 0.60, col_w, 0.28, NAVY_MID)
        text_box(slide, cx + 0.06, Y2 + 0.62, col_w - 0.12, 0.24,
                 tier["label"], font_size=7.5, bold=True, color=WHITE,
                 align=PP_ALIGN.LEFT, wrap=False)

        # WHY band
        rect(slide, cx, Y2 + 0.88, col_w, 0.24, NAVY_LT)
        text_box(slide, cx + 0.06, Y2 + 0.90, col_w - 0.12, 0.22,
                 tier["why"], font_size=7, bold=False, color=WHITE, wrap=False)

        # Body — height capped so it doesn't overlap the adoption-risk bar below
        rect(slide, cx, Y2 + 1.12, col_w, 2.46, GRAY_LT)
        text_box(slide, cx + 0.08, Y2 + 1.16, col_w - 0.16, 2.38,
                 tier["body"], font_size=8, color=NAVY, wrap=True)

    # ── Right side: specialist panels ────────────────────────────────────────
    rx = 7.84
    rw = 5.29

    # TM-40 panel
    rect(slide, rx, Y2, rw, 1.62, NAVY2)
    text_box(slide, rx + 0.08, Y2 + 0.06, rw - 0.16, 0.22,
             "TM-40 SERIES  —  12 Specialist Tracks  (Prereq: TM-30)",
             font_size=8.5, bold=True, color=WHITE, wrap=False)
    text_box(slide, rx + 0.08, Y2 + 0.30, rw - 0.16, 1.28,
             ("WFF Tracks: TM-40A Intelligence · TM-40B Fires · TM-40C Movement & Maneuver\n"
              "TM-40D Sustainment · TM-40E Protection · TM-40F Mission Command\n"
              "Specialist: TM-40G ORSA · TM-40H AI Engineer · TM-40M ML Engineer\n"
              "TM-40J Program Manager · TM-40K Knowledge Manager · TM-40L Software Engineer"),
             font_size=8, color=GRAY_LT, wrap=True)

    # TM-50 panel
    rect(slide, rx, Y2 + 1.66, rw, 1.10, NAVY_MID)
    text_box(slide, rx + 0.08, Y2 + 1.72, rw - 0.16, 0.22,
             "TM-50 SERIES  —  6 Advanced Tracks  (Prereq: TM-40G–M)",
             font_size=8.5, bold=True, color=WHITE, wrap=False)
    text_box(slide, rx + 0.08, Y2 + 1.96, rw - 0.16, 0.76,
             "Expert-level continuation for each TM-40 specialist track.\n"
             "TM-50G ORSA · TM-50H AI Eng · TM-50M MLE · TM-50J PM · TM-50K KM · TM-50L SWE",
             font_size=8, color=GRAY_LT, wrap=True)

    # Training management panel
    Y3 = Y2 + 2.80
    rect(slide, rx, Y3, rw, 0.58, DARK_BG)
    text_box(slide, rx + 0.08, Y3 + 0.04, rw - 0.16, 0.20,
             "TRAINING MANAGEMENT PACKAGE  (READY TO EXECUTE)",
             font_size=8, bold=True, color=WHITE, wrap=False)
    text_box(slide, rx + 0.08, Y3 + 0.26, rw - 0.16, 0.28,
             "MTP · POI · CAD · TEO · Annual Schedule · Enrollment SOP · Lesson Plans · Syllabi",
             font_size=7.5, color=GRAY_LT, wrap=True)

    # ── Adoption risk bar ─────────────────────────────────────────────────────
    Y4 = Y2 + 3.62   # below columns
    rect(slide, 0.10, Y4, 7.60, 0.54, DARK_BG)
    text_box(slide, 0.20, Y4 + 0.04, 7.40, 0.20,
             "ADOPTION RISK WITHOUT THIS TRAINING",
             font_size=8, bold=True, color=GOLD, wrap=False)
    text_box(slide, 0.20, Y4 + 0.26, 7.40, 0.24,
             "Personnel cannot log in or navigate  ·  Staff sections depend on 1–2 specialists  "
             "·  Data products not reproduced  ·  Maven investment wasted",
             font_size=7.5, color=WHITE, wrap=True)

    # ── Stats strip ──────────────────────────────────────────────────────────
    Y5 = Y4 + 0.58
    sw = 7.60 / len(STATS)
    for k, (val, lbl) in enumerate(STATS):
        sx = 0.10 + k * sw
        bg = NAVY if k % 2 == 0 else NAVY_MID
        rect(slide, sx, Y5, sw, 0.52, bg)
        text_box(slide, sx + 0.04, Y5 + 0.02, sw - 0.08, 0.24,
                 val, font_size=14, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, wrap=False)
        text_box(slide, sx + 0.04, Y5 + 0.26, sw - 0.08, 0.24,
                 lbl, font_size=6.5, color=GRAY_LT,
                 align=PP_ALIGN.CENTER, wrap=True)

    # ── Next steps (right, aligned with stats) ───────────────────────────────
    rect(slide, rx, Y4, rw, 1.10, DARK_BG)
    text_box(slide, rx + 0.08, Y4 + 0.04, rw - 0.16, 0.22,
             "NEXT STEPS — COMMAND ACTION REQUIRED",
             font_size=8.5, bold=True, color=GOLD, wrap=False)
    text_box(slide, rx + 0.08, Y4 + 0.28, rw - 0.16, 0.78,
             ("①  Issue training policy letter (template ready)\n"
              "②  Certify instructors per Faculty Development Plan\n"
              "③  Provision Foundry access; schedule first cohort"),
             font_size=8, color=WHITE, wrap=True)

    # Footer
    text_box(slide, 0.10, 6.63, 13.13, 0.20,
             "USAREUR-AF  ·  C2DAO  ·  MSS-POI-001  ·  MAR 2026",
             font_size=7, color=NAVY2, align=PP_ALIGN.CENTER, wrap=False)


def build_project_overview():
    out = OUT_DIR / "MSS_Project_Overview_PROTO.pptx"
    prs = prep_template(out)
    update_title_chrome(prs.slides[0],
                        ["MAVEN SMART SYSTEM",
                         "Command-Wide Adoption Strategy  —  Training Infrastructure"],
                        "USAREUR-AF Operational Data Team  |  C2DAO  |  MAR 2026")
    build_overview_content(prs.slides[1])
    prs.save(str(out))
    print(f"Saved → {out}")


# ══════════════════════════════════════════════════════════════════════════════
# PROTOTYPE 3 — MSS Program Overview (one-pager)
# ══════════════════════════════════════════════════════════════════════════════

def panel(slide, x, y, w, h, title, body_lines,
          header_fill=NAVY, body_fill=GRAY_LT,
          title_size=8, body_size=7.5):
    """Titled info panel: header band + body."""
    rect(slide, x, y, w, 0.26, header_fill)
    text_box(slide, x + 0.06, y + 0.03, w - 0.12, 0.20,
             title, font_size=title_size, bold=True, color=WHITE, wrap=False)
    rect(slide, x, y + 0.26, w, h - 0.26, body_fill)
    body_text = "\n".join(body_lines) if isinstance(body_lines, list) else body_lines
    text_box(slide, x + 0.06, y + 0.30, w - 0.12, h - 0.34,
             body_text, font_size=body_size, color=NAVY, wrap=True)


def build_program_overview_content(slide):
    """6-panel one-pager: left col, center col, right col (2 rows each)."""
    update_template_chrome(slide,
                           "MAVEN SMART SYSTEM (MSS) — TRAINING PROGRAM OVERVIEW")

    # 3-column layout: col widths 4.08, 4.08, 4.10 + margins
    CW   = 4.08     # column width
    GAP  = 0.09
    X0   = 0.10
    Y0   = 0.77
    RH1  = 1.62     # row 1 height (top row)
    RH2  = 4.06     # row 2 height — must not extend past footer at y=6.63

    xs = [X0, X0 + CW + GAP, X0 + 2*(CW + GAP)]

    # ── Row 1: 3 summary panels ───────────────────────────────────────────────
    panel(slide, xs[0], Y0, CW, RH1,
          "WHAT IT IS",
          ["A structured curriculum — three sequential core courses (TM-10/20/30) for "
           "all personnel, followed by 12 functional specialist tracks (TM-40A–M) and "
           "6 advanced tracks (TM-50G–M) for senior practitioners.",
           "",
           "All training delivered in a dedicated MSS Training Environment with "
           "synthetic OPDATA supporting hands-on lab exercises."])

    panel(slide, xs[1], Y0, CW, RH1,
          "TRAINING ENVIRONMENT",
          ["All training conducted in a dedicated MSS Training Environment with "
           "synthetic operational data — realistic scenarios across all WFF and "
           "specialist tracks.",
           "",
           "Students complete lab exercises against live Foundry instances; "
           "no production data is used during training."])

    panel(slide, xs[2], Y0, CW, RH1,
          "SCOPE & SCALE",
          ["▪  3 core courses (TM-10/20/30) + 18 specialty/advanced tracks",
           "▪  160+ published PDFs; Markdown source in version control",
           "▪  Dedicated exercises with hands-on labs per track",
           "▪  PRE + POST assessments for every course",
           "▪  Authority: C2DAO Data Governance Directive (MSS-POI-001)",
           "▪  Review cycle: annual or on major platform update"])

    Y1 = Y0 + RH1 + 0.08   # row 2 start

    # ── Row 2: learning path, TM-40 grid, supporting publications ────────────
    panel(slide, xs[0], Y1, CW, RH2,
          "LEARNING PATH",
          ["TM-10  MAVEN USER  —  All personnel",
           "  Navigate MSS, read dashboards, verify data currency",
           "",
           "TM-20  BUILDER  —  All staff  (Prereq: TM-10)",
           "  No-code pipelines, Workshop dashboards, forms",
           "",
           "TM-30  ADVANCED BUILDER  —  Data specialists  (Prereq: TM-20)",
           "  Python/SQL transforms, ontology design, AIP Logic",
           "",
           "TM-40 SERIES  —  WFF + Specialist tracks  (Prereq: TM-30)",
           "  WFF A–F: Intel, Fires, M&M, Sust, Prot, Msn Cmd",
           "  Specialist G–M: ORSA, AI Eng, MLE, PM, KM, SWE",
           "",
           "TM-50 SERIES  —  Advanced specialist  (Prereq: TM-40 G–M)",
           "  50G–M: expert-level continuation per role"])

    panel(slide, xs[1], Y1, CW, RH2,
          "TM-40 WFF FUNCTIONAL TRACKS  (A–F)",
          ["Prereq: TM-30  ·  No coding required  ·  WFF functional staff",
           "",
           "TM-40A  Intelligence",
           "TM-40B  Fires",
           "TM-40C  Movement & Maneuver",
           "TM-40D  Sustainment",
           "TM-40E  Protection",
           "TM-40F  Mission Command",
           "",
           "TM-40 SPECIALIST TRACKS  (G–M)",
           "Prereq: TM-30  ·  Role-specific technical depth",
           "",
           "TM-40G  ORSA",
           "TM-40H  AI Engineer",
           "TM-40M  ML Engineer",
           "TM-40J  Program Manager",
           "TM-40K  Knowledge Manager",
           "TM-40L  Software Engineer"])

    panel(slide, xs[2], Y1, CW, RH2,
          "SUPPORTING PUBLICATIONS",
          ["DOCTRINE",
           "  Data Literacy — Sr Leaders & Technical Reference",
           "  Foundry Glossary · CDA Constraints & Directives",
           "  EA Series (EA-00 through EA-05)",
           "",
           "SYLLABI",
           "  Per-course objectives, tasks, evaluation criteria",
           "  All 21 courses (TM-10 through TM-50L) published",
           "",
           "EXERCISES",
           "  Hands-on lab guides + environment setup per track",
           "  PRE + POST exams for every course",
           "",
           "TRAINING MANAGEMENT",
           "  POI · MTP · CAD · TEO · Enrollment SOP",
           "  Annual Schedule · Faculty Development Plan",
           "  Lesson Plan Outlines (TM-10 through TM-50)"])

    # Footer
    text_box(slide, 0.10, 6.63, 13.13, 0.20,
             "USAREUR-AF  ·  C2DAO  ·  MSS-POI-001  ·  MAR 2026",
             font_size=7, color=NAVY2, align=PP_ALIGN.CENTER, wrap=False)


def build_program_overview():
    out = OUT_DIR / "MSS_PROGRAM_OVERVIEW_PROTO.pptx"
    prs = prep_template(out)
    update_title_chrome(prs.slides[0],
                        ["MAVEN SMART SYSTEM (MSS)",
                         "Training Program Overview  —  USAREUR-AF C2DAO"],
                        "USAREUR-AF Operational Data Team  |  C2DAO  |  MAR 2026")
    build_program_overview_content(prs.slides[1])
    prs.save(str(out))
    print(f"Saved → {out}")


# ══════════════════════════════════════════════════════════════════════════════
# Entry point
# ══════════════════════════════════════════════════════════════════════════════

if __name__ == "__main__":
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    build_progression()
    build_project_overview()
    build_program_overview()
    print("All 3 prototypes complete.")
