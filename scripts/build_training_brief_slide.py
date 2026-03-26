"""
Single-slide MSS Training Program dep map for CG demo.
Schema mirrors DEPENDENCY_MAP.html:
  - Solid gold arrows  = prereq edges  (SL 1 → SL 2 → SL 3)
  - Dashed border box  = companion/parallel (FBC — not in prereq chain)
  - Companion boxes    = Syllabus / Exercise / Exams per row
Output: maven_training/training_management/MSS_TRAINING_BRIEF.pptx
"""

import copy
import os
from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

REPO_ROOT = os.path.dirname(os.path.dirname(os.path.afbcath(__file__)))

# ── Colors ─────────────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x1F, 0x38, 0x64)
GOLD   = RGBColor(0xC9, 0xA8, 0x4C)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
DGRAY  = RGBColor(0x44, 0x44, 0x44)
LGRAY  = RGBColor(0x88, 0x88, 0x88)
COMP   = RGBColor(0x2E, 0x4D, 0x8A)
TM10C  = RGBColor(0x1F, 0x38, 0x64)
TM20C  = RGBColor(0x2E, 0x4D, 0x8A)
TM30C  = RGBColor(0x3A, 0x62, 0xAA)
FBCC   = RGBColor(0x7C, 0x4D, 0x0A)   # amber — FBC parallel track
CLIGHT = RGBColor(0xCC, 0xDD, 0xFF)

W = Inches(13.33)
H = Inches(7.5)


# ── Shape / XML helpers ────────────────────────────────────────────────────────
def _shape_name(elem):
    for child in elem.iter():
        if child.tag.endswith('}cNvPr'):
            return child.get('name', '')
    return ''


def clear_shapes(slide, preserve=None):
    preserve = preserve or set()
    sp_tree = slide.shapes._spTree
    to_remove = [c for c in list(sp_tree)
                 if c.tag.split('}')[-1] in {'sp', 'pic', 'graphicFrame', 'grpSp', 'cxnSp'}
                 and _shape_name(c) not in preserve]
    for e in to_remove:
        sp_tree.remove(e)


def update_shape_text(slide, name, new_text):
    for shape in slide.shapes:
        if shape.name == name and shape.has_text_frame:
            runs = shape.text_frame.paragraphs[0].runs
            if runs:
                runs[0].text = new_text
                for r in runs[1:]:
                    r.text = ''
            break


def rect(slide, l, t, w, h, fill=None, line_color=None, lw=Pt(0.5), dash=False):
    s = slide.shapes.add_shape(1, l, t, w, h)
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line_color:
        s.line.color.rgb = line_color
        s.line.width = lw
        if dash:
            # Set dashed stroke via XML (prstDash val="dash")
            sp_pr = s._element.find(qn('p:spPr'))
            ln = sp_pr.find(qn('a:ln'))
            if ln is None:
                ln = etree.SubElement(sp_pr, qn('a:ln'))
                ln.set('w', str(int(lw)))
            prstDash = etree.SubElement(ln, qn('a:prstDash'))
            prstDash.set('val', 'dashDot')
    else:
        s.line.fill.background()
    return s


def txt(slide, text, l, t, w, h, sz=Pt(11), bold=False,
        color=DGRAY, align=PP_ALIGN.LEFT, italic=False):
    txb = slide.shapes.add_textbox(l, t, w, h)
    txb.word_wrap = True
    tf = txb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = sz; r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    return txb


def h_arrow(slide, l, t, aw, row_h):
    mid = t + row_h / 2
    rect(slide, l + Inches(0.03), mid - Inches(0.04),
         aw - Inches(0.06), Inches(0.08), fill=GOLD)
    ah = slide.shapes.add_shape(13, l + aw - Inches(0.13),
                                 mid - Inches(0.13), Inches(0.13), Inches(0.26))
    ah.fill.solid(); ah.fill.fore_color.rgb = GOLD
    ah.line.fill.background()


def v_arrow(slide, cx, t, ah):
    rect(slide, cx - Inches(0.05), t, Inches(0.10), ah - Inches(0.14), fill=GOLD)
    a = slide.shapes.add_shape(13, cx - Inches(0.13), t + ah - Inches(0.28),
                                Inches(0.26), Inches(0.28))
    a.fill.solid(); a.fill.fore_color.rgb = GOLD
    a.line.fill.background()


def dashed_connector(slide, x1, y1, x2, y2):
    """Horizontal dashed gold line — companion edge (dep-map schema)."""
    NS_A = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    NS_P = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    from pptx.enum.shapes import MSO_CONNECTOR_TYPE
    cxn = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT, x1, y1, x2, y2)
    sp_pr = cxn._element.find(f'{{{NS_P}}}spPr')
    if sp_pr is None:
        sp_pr = etree.SubElement(cxn._element, f'{{{NS_P}}}spPr')
    ln = sp_pr.find(f'{{{NS_A}}}ln')
    if ln is None:
        ln = etree.SubElement(sp_pr, f'{{{NS_A}}}ln')
    ln.set('w', str(int(Pt(1.5))))
    sf = etree.SubElement(ln, f'{{{NS_A}}}solidFill')
    clr = etree.SubElement(sf, f'{{{NS_A}}}srgbClr')
    clr.set('val', f'{GOLD[0]:02X}{GOLD[1]:02X}{GOLD[2]:02X}')
    pd = etree.SubElement(ln, f'{{{NS_A}}}prstDash')
    pd.set('val', 'dash')
    return cxn


def doc_box(slide, l, t, w, h, label, name, fill):
    rect(slide, l, t, w, h, fill=fill, line_color=RGBColor(0x60, 0x80, 0xC0), lw=Pt(0.5))
    txt(slide, label,
        l + Inches(0.10), t + Inches(0.07),
        w - Inches(0.20), Inches(0.18), sz=Pt(7.5), bold=True, color=GOLD)
    txt(slide, name,
        l + Inches(0.10), t + Inches(0.28),
        w - Inches(0.20), h - Inches(0.34), sz=Pt(9.5), color=WHITE)


def tm_box(slide, l, t, w, h, course, title, hrs, audience, note, fill):
    rect(slide, l, t, w, h, fill=fill, line_color=WHITE, lw=Pt(0.75))
    txt(slide, course,
        l + Inches(0.13), t + Inches(0.07),
        w - Inches(0.26), Inches(0.34), sz=Pt(22), bold=True, color=WHITE)
    txt(slide, title,
        l + Inches(0.13), t + Inches(0.40),
        w - Inches(0.26), Inches(0.22), sz=Pt(10), bold=True, color=GOLD)
    txt(slide, f"{hrs}  ·  {audience}",
        l + Inches(0.13), t + Inches(0.62),
        w - Inches(0.26), Inches(0.20), sz=Pt(9.5), color=WHITE)
    if note:
        txt(slide, note,
            l + Inches(0.13), t + Inches(0.84),
            w - Inches(0.26), Inches(0.20), sz=Pt(8.5), italic=True, color=CLIGHT)


# ── Load PROTO ─────────────────────────────────────────────────────────────────
proto_path = os.path.join(REPO_ROOT, "maven_training", "pdf",
                          "MSS_Project_Overview_PROTO.pptx")
prs = Presentation(proto_path)
PROTO_HEADER = {'hlSlideMaster.BlankHeader', 'TextBox 2', 'TextBox 7'}

s = prs.slides[1]
clear_shapes(s, preserve=PROTO_HEADER)
update_shape_text(s, 'TextBox 2', 'MSS TRAINING PROGRAM — USAREUR-AF')

txt(s,
    "Operator training path  ·  prereq chain: SL 1 → SL 2 → SL 3  "
    "·  FBC: quarterly parallel foundry bootcamp (dashed = not in prereq chain)",
    Inches(1.16), Inches(0.70), W - Inches(1.30), Inches(0.20),
    sz=Pt(8.5), italic=True, color=LGRAY)

# ── Layout geometry ────────────────────────────────────────────────────────────
# Main grid: TM | Syllabus | Exercise | Exams  (left-biased to leave room for FBC)
TM_W   = Inches(2.40)
DOC_W  = Inches(2.20)
ARR_GX = Inches(0.17)
BH     = Inches(1.20)
ROW_GAP = Inches(0.38)

# Left-align main grid
LEFT_X = Inches(0.35)
C_TM  = LEFT_X
C_SYL = C_TM  + TM_W  + ARR_GX
C_EX  = C_SYL + DOC_W + ARR_GX
C_EXM = C_EX  + DOC_W + ARR_GX
MAIN_R = C_EXM + DOC_W   # right edge of main grid

# FBC panel — to the right, spans SL 2 and SL 3 rows
FBC_GAP = Inches(0.30)
FBC_X   = MAIN_R + FBC_GAP
FBC_W   = W - FBC_X - Inches(0.35)

R1_T = Inches(0.98)
R2_T = R1_T + BH + ROW_GAP
R3_T = R2_T + BH + ROW_GAP
TM_CX = LEFT_X + TM_W / 2

FBC_T = R2_T
FBC_H = BH + ROW_GAP + BH   # spans SL 2 and SL 3 rows

# ── Row 1: SL 1 ───────────────────────────────────────────────────────────────
tm_box(s, C_TM, R1_T, TM_W, BH,
       "SL 1", "MAVEN USER", "8 hrs", "All Personnel",
       "No prereq  ·  New soldier in-processing", TM10C)
h_arrow(s, C_SYL - ARR_GX, R1_T, ARR_GX, BH)
doc_box(s, C_SYL, R1_T, DOC_W, BH, "SYLLABUS", "SL 1", COMP)
h_arrow(s, C_EX - ARR_GX, R1_T, ARR_GX, BH)
doc_box(s, C_EX,  R1_T, DOC_W, BH, "EXERCISE", "EX_10  Operator", COMP)
h_arrow(s, C_EXM - ARR_GX, R1_T, ARR_GX, BH)
doc_box(s, C_EXM, R1_T, DOC_W, BH, "EXAMS", "Pre-assessment\nPost-assessment", COMP)

v_arrow(s, TM_CX, R1_T + BH, ROW_GAP)

# ── Row 2: SL 2 ───────────────────────────────────────────────────────────────
tm_box(s, C_TM, R2_T, TM_W, BH,
       "SL 2", "BUILDER", "40 hrs", "Data-Adjacent Staff",
       "Prereq: SL 1", TM20C)
h_arrow(s, C_SYL - ARR_GX, R2_T, ARR_GX, BH)
doc_box(s, C_SYL, R2_T, DOC_W, BH, "SYLLABUS", "SL 2", COMP)
h_arrow(s, C_EX - ARR_GX, R2_T, ARR_GX, BH)
doc_box(s, C_EX,  R2_T, DOC_W, BH, "EXERCISE", "EX_20  No-Code Builder", COMP)
h_arrow(s, C_EXM - ARR_GX, R2_T, ARR_GX, BH)
doc_box(s, C_EXM, R2_T, DOC_W, BH, "EXAMS", "Pre-assessment\nPost-assessment", COMP)

# Dashed companion edge: SL 2 right edge → FBC left edge (dep-map schema)
dashed_connector(s,
    C_TM + TM_W, R2_T + BH / 2,
    FBC_X,        R2_T + BH / 2)

v_arrow(s, TM_CX, R2_T + BH, ROW_GAP)

# ── Row 3: SL 3 ───────────────────────────────────────────────────────────────
tm_box(s, C_TM, R3_T, TM_W, BH,
       "SL 3", "ADVANCED BUILDER", "40 hrs", "Unit Data Leads",
       "Prereq: SL 2", TM30C)
h_arrow(s, C_SYL - ARR_GX, R3_T, ARR_GX, BH)
doc_box(s, C_SYL, R3_T, DOC_W, BH, "SYLLABUS", "SL 3", COMP)
h_arrow(s, C_EX - ARR_GX, R3_T, ARR_GX, BH)
doc_box(s, C_EX,  R3_T, DOC_W, BH, "EXERCISE", "EX_30  Adv Builder", COMP)
h_arrow(s, C_EXM - ARR_GX, R3_T, ARR_GX, BH)
doc_box(s, C_EXM, R3_T, DOC_W, BH, "EXAMS", "Pre-assessment\nPost-assessment", COMP)

# ── FBC sidebar — dashed border, amber, outside prereq chain ──────────────────
rect(s, FBC_X, FBC_T, FBC_W, FBC_H,
     fill=RGBColor(0x3A, 0x22, 0x04),   # very dark amber background
     line_color=GOLD, lw=Pt(1.5), dash=True)

txt(s, "FBC",
    FBC_X + Inches(0.12), FBC_T + Inches(0.10),
    FBC_W - Inches(0.24), Inches(0.30),
    sz=Pt(18), bold=True, color=GOLD)

txt(s, "FOUNDRY BOOTCAMP",
    FBC_X + Inches(0.12), FBC_T + Inches(0.40),
    FBC_W - Inches(0.24), Inches(0.20),
    sz=Pt(9.5), bold=True, color=GOLD)

txt(s,
    "Quarterly · 5-day supervised\nbuild event\n\n"
    "Prereq: SL 2 Go\n+ command-validated project\n\n"
    "Does NOT grant SL 3 credit\nDoes NOT unlock SL 4",
    FBC_X + Inches(0.12), FBC_T + Inches(0.65),
    FBC_W - Inches(0.24), FBC_H - Inches(0.75),
    sz=Pt(9), color=WHITE)

# "parallel track" label on the dashed connector
txt(s, "parallel track",
    C_TM + TM_W + Inches(0.05), R2_T + BH / 2 - Inches(0.22),
    FBC_GAP + Inches(0.10), Inches(0.18),
    sz=Pt(7), italic=True, color=LGRAY)

# ── Depth note ─────────────────────────────────────────────────────────────────
NOTE_T = R3_T + BH + Inches(0.16)
txt(s,
    "SL 4 / SL 5  —  21 specialist & advanced tracks published and ready  "
    "(WFF: Intel · Fires · M&M · Sustainment · Protection · MC  |  "
    "Technical: ORSA · AI · ML · PM · KM · SWE)  —  Execution follows SL 1/2/3 baseline",
    Inches(0.35), NOTE_T, MAIN_R - Inches(0.35), Inches(0.24),
    sz=Pt(8.5), italic=True, color=DGRAY, align=PP_ALIGN.CENTER)

# ── Remove cover (index 0) ────────────────────────────────────────────────────
sldIdLst = prs.slides._sldIdLst
if len(sldIdLst) > 0:
    first = sldIdLst[0]
    rId = first.get(qn('r:id'))
    del sldIdLst[0]
    if rId:
        try:
            prs.part.drop_rel(rId)
        except Exception:
            pass

# ── Save ──────────────────────────────────────────────────────────────────────
# ── Slide 2: Capacity Projection ─────────────────────────────────────────────
# Uses the pre-rendered PNG from scripts/capacity_graph.py
cap_img = os.path.join(REPO_ROOT, "capacity_projection.png")

# Duplicate the layout from slide 1 (same PROTO header)
layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
s2 = prs.slides.add_slide(layout)
clear_shapes(s2, preserve=PROTO_HEADER)
update_shape_text(s2, 'TextBox 2', 'MSS TRAINING PROGRAM — CAPACITY PROJECTION')

if os.path.exists(cap_img):
    # Center the image on the slide with padding for the header
    from PIL import Image
    try:
        with Image.open(cap_img) as im:
            img_w, img_h = im.size
    except ImportError:
        # Fallback aspect ratio from our 15×7.5 figsize @ 180 dpi
        img_w, img_h = 2700, 1350

    aspect = img_w / img_h
    # Available area below header
    avail_w = W - Inches(0.60)
    avail_h = H - Inches(1.10)
    # Fit to available area preserving aspect ratio
    if avail_w / avail_h > aspect:
        pic_h = avail_h
        pic_w = int(pic_h * aspect)
    else:
        pic_w = avail_w
        pic_h = int(pic_w / aspect)

    pic_l = (W - pic_w) // 2
    pic_t = Inches(0.80)
    s2.shapes.add_picture(cap_img, pic_l, pic_t, pic_w, pic_h)
    print("  2 slides: dep map + capacity projection")
else:
    txt(s2, f"[capacity_projection.png not found — run scripts/capacity_graph.py first]",
        Inches(1.0), Inches(3.0), Inches(10.0), Inches(1.0),
        sz=Pt(14), bold=True, color=RGBColor(0x8A, 0x1A, 0x1A), align=PP_ALIGN.CENTER)
    print("  WARNING: capacity_projection.png not found — placeholder added")
    print("  Run: python3 scripts/capacity_graph.py  (then re-run this script)")

# ── Save ──────────────────────────────────────────────────────────────────────
out_path = os.path.join(REPO_ROOT, "maven_training", "training_management",
                        "MSS_TRAINING_BRIEF.pptx")
prs.save(out_path)
print(f"Saved: {out_path}")
