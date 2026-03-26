#!/usr/bin/env python3
"""
Prototype: reformat Maven_Rollout_Plan to USAREUR-AF PPT template.
Run from repo root: python3 scripts/build_proto_rollout.py
Output: maven_training/pdf/Maven_Rollout_Plan_PROTO.pptx
"""

import shutil
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml

# ── Paths ──────────────────────────────────────────────────────────────────────
TEMPLATE = Path("/home/dale/Downloads/USAREUR-AF PPT Template.pptx")
OUTPUT   = Path("maven_training/pdf/Maven_Rollout_Plan_PROTO.pptx")

# ── Palette (extracted from template inspection) ───────────────────────────────
NAVY    = RGBColor(0x17, 0x32, 0x5C)   # primary header / text
NAVY2   = RGBColor(0x18, 0x33, 0x5F)   # phase-column fill (slightly lighter)
PURPLE  = RGBColor(0x70, 0x30, 0xA0)   # classification banner text
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
ROW_ODD  = RGBColor(0xD6, 0xDC, 0xE4)  # alternating row fill A
ROW_EVEN = RGBColor(0xEB, 0xEE, 0xF3)  # alternating row fill B
TIMELINE_BG = RGBColor(0x0D, 0x1F, 0x3C)  # dark navy for timeline bar

# ── Slide metadata ─────────────────────────────────────────────────────────────
CLASSIFICATION = ""
SLIDE_TITLE    = "MAVEN BUILDER TRAINING — FORCE ROLLOUT PLAN"
POC_LINE       = "USAREUR-AF Operational Data Team  |  C2DAO  |  MAR 2026"

TIMELINE = (
    "RECOMMENDED TIMELINE:  Directive signed NLT D+30  |  Cadre designated NLT D+45  |  "
    "First SL 1 cohort NLT D+60  |  Full SL 4 coverage NLT D+180"
)

# ── Rollout content (6 phases) ─────────────────────────────────────────────────
PHASES = [
    (
        "1 | COMMAND\nAUTHORITY",
        "Issue training directive (EXORD or policy memo) designating Maven Builder Training as a command priority.",
        "CDR decision: mandatory vs. encouraged by role.\nG3/J3 staff action.",
        "Signed directive published to all subordinate units.",
    ),
    (
        "2 | CADRE\nSTANDUP",
        "Identify and assign 5 cadre: 1 OIC + 4 SMEs covering SL 1/2/3/4 tracks. Develop instructor SOP.",
        "G6/G2/S6 nominations.\nOIC appointed by CDR.",
        "Cadre certified; SOP signed; instructors completed SL 4 series.",
    ),
    (
        "3 | AUDIENCE\nSEQUENCING",
        "Phase training: SL 1 (all staff) → SL 2/3 (data-adjacent billets) → SL 4 series by role.",
        "G3 decision: phased by echelon or MOS priority.\nUnit S3s build into ATP.",
        "100% SL 1 completion; SL 4 track assignments published.",
    ),
    (
        "4 | RESOURCE\n& ACCESS",
        "Provision Foundry seats, schedule lab time/VTC blocks, integrate into unit training calendar.",
        "G6 action: license provisioning.\nG3: calendar deconfliction.",
        "All students have Foundry access NLT training start date.",
    ),
    (
        "5 | TRACK\n& REPORT",
        "Establish completion tracking in DTMS/ATRRS. Set Go/No-Go checkpoints per MTP. Report monthly.",
        "G1/S1 action: reporting format.\nCDR sets completion deadline (recommended D+180).",
        "Monthly completion report to CDR; Go/No-Go recorded per student.",
    ),
    (
        "6 | FEEDBACK\n& SUSTAIN",
        "Conduct AAR after first cohort; update curriculum and cadre SOP quarterly. Integrate into onboarding.",
        "Cadre OIC owns curriculum updates.\nCDR approves policy changes.",
        "Quarterly AAR complete; curriculum version controlled; onboarding updated.",
    ),
]


# ── Helpers ────────────────────────────────────────────────────────────────────

def set_cell_bg(cell, rgb: RGBColor):
    """Apply solid fill to a table cell via XML."""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    # Remove any existing fill
    for tag in ("a:solidFill", "a:noFill", "a:gradFill", "a:pattFill"):
        existing = tcPr.find(qn(tag))
        if existing is not None:
            tcPr.remove(existing)
    ns = 'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"'
    solid = parse_xml(
        f'<a:solidFill {ns}><a:srgbClr val="{str(rgb)}"/></a:solidFill>'
    )
    tcPr.append(solid)


def style_cell(cell, text: str, *,
               font_size=9, bold=False,
               fg: RGBColor = None, bg: RGBColor = None,
               align=PP_ALIGN.LEFT):
    """Set cell background, text, and font properties.
    Multi-line text delimited by \\n becomes separate paragraphs.
    """
    if bg:
        set_cell_bg(cell, bg)

    tf = cell.text_frame
    tf.word_wrap = True

    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.bold = bold
        if fg:
            run.font.color.rgb = fg


def remove_shape(shape):
    """Remove a shape element from its parent XML node."""
    el = shape._element
    el.getparent().remove(el)


# ── Slide builders ─────────────────────────────────────────────────────────────

def build_title_slide(slide):
    """Update template title slide with rollout plan metadata."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        txt  = shape.text_frame.text.strip()
        name = shape.name

        if name == "Title 1":
            tf = shape.text_frame
            tf.clear()
            p1 = tf.paragraphs[0]
            p1.alignment = PP_ALIGN.CENTER
            r1 = p1.add_run()
            r1.text = "MAVEN BUILDER TRAINING"
            r1.font.size = Pt(28)
            r1.font.bold = True

            p2 = tf.add_paragraph()
            p2.alignment = PP_ALIGN.CENTER
            r2 = p2.add_run()
            r2.text = "FORCE ROLLOUT PLAN"
            r2.font.size = Pt(20)
            r2.font.bold = False

        elif name.startswith("hlSlideMaster") and txt in ("Classification", "CUI", "UNCLASSIFIED"):
            # Classification banners — update text, keep purple color
            for p in shape.text_frame.paragraphs:
                for run in p.runs:
                    run.text = CLASSIFICATION

        elif name == "TextBox 3" and txt == "Classification":
            # Center-of-slide classification line
            for p in shape.text_frame.paragraphs:
                for run in p.runs:
                    run.text = CLASSIFICATION

        elif name == "TextBox 5" and txt == "POC":
            tf = shape.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = POC_LINE
            run.font.size = Pt(12)
            run.font.color.rgb = NAVY


def build_content_slide(slide):
    """Build the 6-phase rollout table on the template's blank content slide."""
    # ── Update existing template text boxes ────────────────────────────────────
    to_remove = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        txt  = shape.text_frame.text.strip()
        name = shape.name

        if name == "TextBox 2" or (txt == "Header"):
            # Slide header line
            tf = shape.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = SLIDE_TITLE
            run.font.size = Pt(16)
            run.font.bold = True
            run.font.color.rgb = NAVY

        elif name.startswith("hlSlideMaster") and txt in ("CUI", "UNCLASSIFIED", "Classification"):
            for p in shape.text_frame.paragraphs:
                for run in p.runs:
                    run.text = CLASSIFICATION

        elif name == "TextBox 8":
            # Large default content area — remove so it doesn't overlap our table
            to_remove.append(shape)

    for shape in to_remove:
        remove_shape(shape)

    # ── Build the rollout table ────────────────────────────────────────────────
    # Geometry: top just below the header (0.92"), leave room for timeline + bottom banner
    N_ROWS  = 7          # 1 header + 6 phase rows
    N_COLS  = 4
    T_LEFT  = Inches(0.10)
    T_TOP   = Inches(0.92)
    T_WIDTH = Inches(13.13)
    T_HEIGHT = Inches(5.55)   # header row ~0.38" + 6 × ~0.862"

    tbl_shape = slide.shapes.add_table(N_ROWS, N_COLS, T_LEFT, T_TOP, T_WIDTH, T_HEIGHT)
    table = tbl_shape.table

    # Column widths: Phase | Action | Decision/Owner | Success Standard
    col_widths = [Inches(1.20), Inches(4.45), Inches(4.10), Inches(3.38)]
    for i, w in enumerate(col_widths):
        table.columns[i].width = w

    # Row heights
    hdr_h  = Emu(int(0.38 * 914400))
    row_h  = Emu(int(0.862 * 914400))
    table.rows[0].height = hdr_h
    for i in range(1, N_ROWS):
        table.rows[i].height = row_h

    # Header row
    HEADERS = ["PHASE", "ACTION", "DECISION / OWNER", "SUCCESS STANDARD"]
    for j, hdr in enumerate(HEADERS):
        style_cell(table.cell(0, j), hdr,
                   font_size=10, bold=True,
                   fg=WHITE, bg=NAVY,
                   align=PP_ALIGN.CENTER)

    # Data rows (alternating fills)
    for i, (phase, action, decision, standard) in enumerate(PHASES):
        row_bg = ROW_ODD if i % 2 == 0 else ROW_EVEN
        style_cell(table.cell(i + 1, 0), phase,
                   font_size=8, bold=True,
                   fg=WHITE, bg=NAVY2,
                   align=PP_ALIGN.CENTER)
        style_cell(table.cell(i + 1, 1), action,
                   font_size=8, fg=NAVY, bg=row_bg)
        style_cell(table.cell(i + 1, 2), decision,
                   font_size=8, fg=NAVY, bg=row_bg)
        style_cell(table.cell(i + 1, 3), standard,
                   font_size=8, fg=NAVY, bg=row_bg)

    # ── Timeline bar below table ───────────────────────────────────────────────
    tl_top  = T_TOP + T_HEIGHT + Inches(0.04)
    tl_left = Inches(0.10)
    tl_w    = Inches(13.13)
    tl_h    = Inches(0.28)

    tl_box = slide.shapes.add_textbox(tl_left, tl_top, tl_w, tl_h)
    # Dark navy background via XML fill on the txBody's parent spPr
    spPr = tl_box._element.spPr
    ns = 'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"'
    solidFill = parse_xml(
        f'<a:solidFill {ns}><a:srgbClr val="{str(TIMELINE_BG)}"/></a:solidFill>'
    )
    spPr.append(solidFill)

    tf = tl_box.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = TIMELINE
    run.font.size = Pt(7.5)
    run.font.bold = False
    run.font.color.rgb = WHITE

    # ── Footer line ────────────────────────────────────────────────────────────
    ft_top = tl_top + tl_h + Inches(0.02)
    ft_box = slide.shapes.add_textbox(Inches(0.10), ft_top, Inches(13.13), Inches(0.20))
    ft = ft_box.text_frame
    ft.word_wrap = False
    p = ft.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "USAREUR-AF  ·  C2DAO  ·  MSS-POI-001  ·  MAR 2026"
    run.font.size = Pt(7)
    run.font.color.rgb = NAVY2


# ── Main ───────────────────────────────────────────────────────────────────────

def main():
    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    shutil.copy(TEMPLATE, OUTPUT)

    prs = Presentation(str(OUTPUT))

    build_title_slide(prs.slides[0])
    build_content_slide(prs.slides[1])

    prs.save(str(OUTPUT))
    print(f"Prototype saved → {OUTPUT}")


if __name__ == "__main__":
    main()
