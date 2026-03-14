#!/usr/bin/env python3
"""
Prototype: Army Data Orientation — reformatted to USAREUR-AF PPT template.
Source: maven_training/source_material/course_portal/assets/army_data_orientation_v1/
Run from repo root: python3 scripts/build_proto_army_data_orientation.py
Output: maven_training/pdf/Army_Data_Orientation_PROTO.pptx
"""

import shutil
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml

# ── Paths ──────────────────────────────────────────────────────────────────────
TEMPLATE = Path("/home/dale/Downloads/USAREUR-AF PPT Template.pptx")
OUTPUT   = Path("maven_training/pdf/Army_Data_Orientation_PROTO.pptx")

# ── Palette ────────────────────────────────────────────────────────────────────
NAVY      = RGBColor(0x17, 0x32, 0x5C)
NAVY2     = RGBColor(0x18, 0x33, 0x5F)
NAVY_MID  = RGBColor(0x26, 0x4D, 0x7E)
DARK_BG   = RGBColor(0x0D, 0x1F, 0x3C)
PURPLE    = RGBColor(0x70, 0x30, 0xA0)   # classification
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE = RGBColor(0xF5, 0xF6, 0xF8)
GRAY_LT   = RGBColor(0xE2, 0xE8, 0xF0)
GOLD      = RGBColor(0xC8, 0x9A, 0x00)
GREEN     = RGBColor(0x2A, 0x6E, 0x3F)   # "new model" / positive
GREEN_LT  = RGBColor(0xD4, 0xED, 0xDA)
RED       = RGBColor(0xBD, 0x2E, 0x2E)   # "problem" / legacy
RED_LT    = RGBColor(0xF8, 0xD7, 0xDA)
BLUE      = RGBColor(0x1A, 0x5C, 0xA8)   # data integration layer
BLUE_LT   = RGBColor(0xD0, 0xE4, 0xF5)
TEAL      = RGBColor(0x1A, 0x78, 0x6E)   # ontology layer
TEAL_LT   = RGBColor(0xCF, 0xE8, 0xE6)
VIOLET    = RGBColor(0x6B, 0x32, 0x8A)   # understanding layer
VIOLET_LT = RGBColor(0xE8, 0xD5, 0xF0)

CLASSIFICATION = "UNCLASSIFIED"

# ── Content geometry ───────────────────────────────────────────────────────────
Y0  = 0.70   # content start (below header + small gap)
YB  = 7.20   # content end (above bottom banner)
CH  = YB - Y0  # 6.50" total content height
X0  = 0.15
XE  = 13.18
CW  = XE - X0  # 13.03" total content width


# ── Low-level helpers ──────────────────────────────────────────────────────────

def solid_fill_xml(rgb: RGBColor) -> str:
    ns = 'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"'
    return f'<a:solidFill {ns}><a:srgbClr val="{str(rgb)}"/></a:solidFill>'


def set_cell_bg(cell, rgb: RGBColor):
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for tag in ("a:solidFill", "a:noFill", "a:gradFill"):
        el = tcPr.find(qn(tag))
        if el is not None:
            tcPr.remove(el)
    tcPr.append(parse_xml(solid_fill_xml(rgb)))


def set_textbox_fill(tb, rgb: RGBColor):
    spPr = tb._element.spPr
    for tag in (qn("a:solidFill"), qn("a:noFill"), qn("a:gradFill")):
        el = spPr.find(tag)
        if el is not None:
            spPr.remove(el)
    spPr.append(parse_xml(solid_fill_xml(rgb)))


def rect(slide, x, y, w, h, fill: RGBColor, no_line=True):
    shp = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if no_line:
        shp.line.fill.background()
    shp.text_frame.clear()
    return shp


def tb(slide, x, y, w, h, lines, *, font_size=9, bold=False,
       color: RGBColor = NAVY, align=PP_ALIGN.LEFT, wrap=True,
       fill: RGBColor = None, italic=False):
    """Add text box. lines can be str (split on \\n) or list[str]."""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    if fill:
        set_textbox_fill(box, fill)
    tf = box.text_frame
    tf.word_wrap = wrap
    if isinstance(lines, str):
        lines = lines.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return box


def panel(slide, x, y, w, h, title, body, *,
          header_fill=NAVY, body_fill=OFF_WHITE,
          title_color=WHITE, body_color=NAVY,
          title_size=9, body_size=8.5, title_h=0.28):
    """Titled content panel with header band + body."""
    rect(slide, x, y, w, title_h, header_fill)
    tb(slide, x + 0.06, y + 0.02, w - 0.12, title_h - 0.04,
       title, font_size=title_size, bold=True, color=title_color, wrap=False)
    rect(slide, x, y + title_h, w, h - title_h, body_fill)
    if isinstance(body, list):
        body_text = "\n".join(body)
    else:
        body_text = body
    tb(slide, x + 0.08, y + title_h + 0.06, w - 0.16, h - title_h - 0.08,
       body_text, font_size=body_size, color=body_color, wrap=True)


def remove_shape(shape):
    el = shape._element
    el.getparent().remove(el)


# ── Chrome helpers ─────────────────────────────────────────────────────────────

def clone_content_slide(prs):
    """Clone template slide 2 (content layout) as a new slide.
    Copies the full shape tree so every content slide inherits the exact
    same chrome as the template — background, classification banner, header.
    """
    import copy
    src = prs.slides[1]
    slide = prs.slides.add_slide(src.slide_layout)
    sp_tree = slide.shapes._spTree
    # Remove everything except the 2 required group elements (nvGrpSpPr, grpSpPr)
    for child in list(sp_tree)[2:]:
        sp_tree.remove(child)
    # Copy all shapes from template slide 2
    for child in list(src.shapes._spTree)[2:]:
        sp_tree.append(copy.deepcopy(child))
    return slide


def setup_slide(slide, header_text, classification=CLASSIFICATION):
    """Update chrome on a cloned content slide: set header text, fix
    classification strings, remove the default content placeholder.
    """
    to_remove = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        txt  = shape.text_frame.text.strip()
        name = shape.name
        if name == "TextBox 2" or txt == "Header":
            tf = shape.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = header_text
            run.font.size = Pt(14)
            run.font.bold = True
            run.font.color.rgb = NAVY
        elif txt in ("CUI", "UNCLASSIFIED", "Classification"):
            for p in shape.text_frame.paragraphs:
                for run in p.runs:
                    run.text = classification
        elif name == "TextBox 8":
            to_remove.append(shape)
    for s in to_remove:
        remove_shape(s)
    # Bottom banner (slide 2 only has top in the template)
    bot = slide.shapes.add_textbox(Inches(0), Inches(7.25), Inches(13.33), Inches(0.25))
    tf = bot.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = classification
    run.font.size = Pt(9)
    run.font.bold = True
    run.font.color.rgb = PURPLE


def add_footer(slide, text="USAREUR-AF  ·  C2DAO  ·  MSS-POI-001  ·  MAR 2026  ·  UNCLASSIFIED"):
    tb(slide, X0, 6.96, CW, 0.20, text,
       font_size=7, color=NAVY2, align=PP_ALIGN.CENTER, wrap=False)


def update_title_chrome(slide, title_lines, poc, classification=CLASSIFICATION):
    """Populate the template title slide."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        txt  = shape.text_frame.text.strip()
        name = shape.name
        if name == "Title 1":
            tf = shape.text_frame
            tf.clear()
            for i, line in enumerate(title_lines):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.alignment = PP_ALIGN.CENTER
                run = p.add_run()
                run.text = line
                run.font.size = Pt(26 if i == 0 else 17)
                run.font.bold = (i == 0)
        elif name.startswith("hlSlideMaster") and txt in ("Classification", "CUI", "UNCLASSIFIED"):
            for p in shape.text_frame.paragraphs:
                for run in p.runs:
                    run.text = classification
        elif name == "TextBox 3" and txt == "Classification":
            for p in shape.text_frame.paragraphs:
                for run in p.runs:
                    run.text = classification
        elif name == "TextBox 5" and txt == "POC":
            tf = shape.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = poc
            run.font.size = Pt(11)
            run.font.color.rgb = NAVY




# ══════════════════════════════════════════════════════════════════════════════
# Slide builders
# ══════════════════════════════════════════════════════════════════════════════

def slide_first_principles(slide):
    """Slide 1: FIRST PRINCIPLES — 2-col layout."""
    mid = X0 + CW / 2

    # Left panel: WHY THE ARMY EXISTS
    panel(slide, X0, Y0, CW/2 - 0.08, CH - 0.28,
          "WHY THE ARMY EXISTS",
          ["The Army exists to fight and win the Nation's wars by providing prompt, "
           "sustained land dominance across the full range of military operations.",
           "",
           "Everything the Army does serves this purpose:",
           "",
           "▸  Organize, train, and equip forces for combat",
           "▸  Deploy land forces anywhere in the world",
           "▸  Conduct sustained operations across all domains",
           "▸  Integrate with joint and multinational partners"],
          header_fill=NAVY, title_color=WHITE, title_size=9)
    tb(slide, X0 + 0.08, Y0 + CH - 0.30, CW/2 - 0.20, 0.22,
       "Title 10, U.S. Code  ·  ADP 1",
       font_size=7.5, color=NAVY_MID, italic=True)

    # Right panel: WHY A STAFF EXISTS
    panel(slide, mid + 0.08, Y0, CW/2 - 0.08, CH - 0.28,
          "WHY A STAFF EXISTS",
          ["A staff exists for one reason: to help the commander make and execute decisions.",
           "",
           "The staff does this by:",
           "",
           "▸  Collecting and analyzing information",
           "▸  Transforming data into understanding",
           "▸  Presenting options to the commander",
           "▸  Synchronizing execution of orders",
           "▸  Assessing whether operations achieve objectives"],
          header_fill=TEAL, title_color=WHITE, title_size=9)
    tb(slide, mid + 0.16, Y0 + CH - 0.30, CW/2 - 0.20, 0.22,
       "Every role in the data space serves this chain.",
       font_size=8, color=TEAL, bold=True, italic=True)

    add_footer(slide)


def slide_operational_environment(slide):
    """Slide 2: THE OPERATIONAL ENVIRONMENT — 3-col + quote bar."""

    tb(slide, X0, Y0, CW, 0.22,
       "ADP 3-13  —  Army forces understand their environment through three dimensions.",
       font_size=8.5, color=NAVY_MID, italic=True)

    cw3 = (CW - 0.04) / 3
    dims = [
        ("HUMAN DIMENSION",
         ["People and the interaction between individuals and groups — "
          "how they understand, decide, generate will, and act."],
         "Values · Fear · Morale · Culture · Will to fight",
         NAVY),
        ("INFORMATION DIMENSION",
         ["The content and data that individuals, groups, and systems "
          "communicate and exchange — and the processes to do so."],
         "Data · Narratives · Intel · Comms · Networks",
         TEAL),
        ("PHYSICAL DIMENSION",
         ["The material characteristics and capabilities — both natural and "
          "manufactured — within the environment."],
         "Terrain · Weather · Weapons · Infrastructure",
         BLUE),
    ]
    col_h = CH - 0.72 - 0.28  # leave room for intro + quote bar

    for i, (title, body, footer, color) in enumerate(dims):
        cx = X0 + i * (cw3 + 0.02)
        panel(slide, cx, Y0 + 0.28, cw3, col_h,
              title, body,
              header_fill=color, title_color=WHITE, body_size=8.5)
        tb(slide, cx + 0.08, Y0 + 0.28 + col_h - 0.24, cw3 - 0.16, 0.22,
           footer, font_size=7.5, color=color, italic=True)

    # Quote bar
    qy = Y0 + 0.28 + col_h + 0.06
    rect(slide, X0, qy, CW, 0.70, DARK_BG)
    tb(slide, X0 + 0.12, qy + 0.06, CW - 0.24, 0.58,
       ['"Information is central to everything we do — it is the basis of intelligence, '
        'a fundamental element of command and control, and the foundation for '
        'communicating thoughts, opinions, and ideas."',
        "ADP 3-13 Foreword  ·  Your work lives in the information dimension."],
       font_size=8, color=WHITE, italic=False)

    add_footer(slide)


def slide_data_decisions_doctrine(slide):
    """Slide 3: DATA · DECISIONS · DOCTRINE — hierarchy stack + role panel."""

    tb(slide, X0, Y0, CW, 0.22,
       "ADP 3-13  —  The cognitive hierarchy: how raw signals become commander action.",
       font_size=8.5, color=NAVY_MID, italic=True)

    # Left: cognitive hierarchy (cascading/indented blocks)
    lw = CW * 0.52
    layers = [
        ("DATA",          "Any signal or observation from the environment.\nA radar sounding. A fuel report. A sensor reading.",          OFF_WHITE,  NAVY,    NAVY),
        ("INFORMATION",   "Data placed in context — who sent it, what it\nmeans, and how it relates to the mission.",                     GRAY_LT,    NAVY,    GREEN),
        ("KNOWLEDGE",     "Analysis refines information. Patterns emerge.\nRelationships between data points connect.",                   BLUE_LT,    NAVY,    BLUE),
        ("UNDERSTANDING", "Judgment transforms knowledge into insight.\nThe commander sees the whole picture.",                           VIOLET_LT,  NAVY,    VIOLET),
        ("DECISION → ACTION", "",                                                                                                         GOLD,       NAVY,    NAVY),
    ]
    indent = 0.0
    ly = Y0 + 0.32
    lh = 0.72
    for label, desc, bg, fg, accent in layers:
        rect(slide, X0 + indent, ly, lw - indent, lh, bg)
        rect(slide, X0 + indent, ly, 0.06, lh, accent)  # left accent bar
        tb(slide, X0 + indent + 0.10, ly + 0.04, lw - indent - 0.14, 0.24,
           label, font_size=9, bold=True, color=accent if label != "DECISION → ACTION" else NAVY)
        if desc:
            tb(slide, X0 + indent + 0.10, ly + 0.28, lw - indent - 0.14, 0.40,
               desc, font_size=8, color=fg)
        # Arrow between layers (except last)
        if label != "DECISION → ACTION":
            tb(slide, X0 + lw/2 - 0.10, ly + lh, 0.20, 0.16,
               "▼", font_size=10, color=NAVY_MID, align=PP_ALIGN.CENTER)
        ly += lh + 0.16
        indent += 0.18

    # Right: YOUR ROLE IN THIS CHAIN
    rx = X0 + lw + 0.20
    rw = CW - lw - 0.20
    panel(slide, rx, Y0 + 0.32, rw, CH - 0.54,
          "YOUR ROLE IN THIS CHAIN",
          ["As a data practitioner, you operate at the foundation of this hierarchy — "
           "and that is where the leverage is.",
           "",
           "If the data is wrong, every layer above it fails.",
           "",
           "If data is unavailable, the commander is blind.",
           "",
           "Foundry is the engine that transforms raw data into information and "
           "knowledge at enterprise scale.",
           "",
           "Your pipelines feed commander decisions."],
          header_fill=NAVY_MID, title_color=WHITE, body_size=9)

    add_footer(slide)


def slide_what_makes_us_different(slide):
    """Slide 4: WHAT MAKES US DIFFERENT — 2-col + full-width bottom bar."""

    panel_h = CH - 1.10

    # Left: THE RAT AND THE MAZE
    panel(slide, X0, Y0, CW/2 - 0.08, panel_h,
          "THE RAT AND THE MAZE",
          ["You can teach a rat to run a maze. Through trial and error, "
           "it will eventually master the route.",
           "",
           "But that rat cannot pass its knowledge to the next rat.",
           "",
           "Every generation starts from zero.",
           "",
           "Humans are different. We codify experience.\nWe write it down. We call it doctrine."],
          header_fill=GOLD, title_color=NAVY, body_color=NAVY)

    # Right: WHAT IS DOCTRINE?
    panel(slide, X0 + CW/2 + 0.08, Y0, CW/2 - 0.08, panel_h,
          "WHAT IS DOCTRINE?",
          ["Fundamental principles, with supporting tactics, techniques, procedures, "
           "and terms and symbols, used for the conduct of operations and as a guide "
           "for actions.",
           "",
           "It is authoritative but requires judgment in application.",
           "",
           "ADP 1-01  ·  Doctrine Primer"],
          header_fill=NAVY, title_color=WHITE)

    # Bottom full-width bar: THE ONLY CONSTANT
    by = Y0 + panel_h + 0.08
    rect(slide, X0, by, CW, 0.90, DARK_BG)
    tb(slide, X0 + 0.12, by + 0.04, 0.22, 0.22,
       "◉", font_size=14, color=GOLD)
    tb(slide, X0 + 0.36, by + 0.04, 1.20, 0.22,
       "THE ONLY CONSTANT", font_size=9, bold=True, color=GOLD)
    tb(slide, X0 + 0.36, by + 0.28, CW - 0.48, 0.58,
       "People rotate. Technology changes. Platforms evolve. The soldiers who sit in your seat "
       "20 years from now will have different tools, different threats, and different missions. "
       "The only thing that connects you to them is doctrine — the accumulated, codified "
       "professional knowledge of an institution that has been learning since 1775.",
       font_size=8, color=WHITE)

    add_footer(slide)


def slide_where_does_data_fit(slide):
    """Slide 5: SO WHERE DOES DATA FIT? — 4-col flow + summary text."""

    flow = [
        ("DOCTRINE",          "Defines what we need to know and do",                NAVY),
        ("ONTOLOGY",          "Models the world as objects and relationships",       TEAL),
        ("DATA ARCHITECTURE", "Connects, cleans, and delivers enterprise data",     BLUE),
        ("DECISIONS",         "Commanders act on trusted information",              NAVY_MID),
    ]

    cw4 = (CW - 0.18) / 4  # 3 arrows between 4 boxes = 3 × 0.06
    bh  = 2.20
    by  = Y0 + 0.40

    for i, (title, body, color) in enumerate(flow):
        cx = X0 + i * (cw4 + 0.06)
        rect(slide, cx, by, cw4, 0.08, color)   # top accent bar
        rect(slide, cx, by + 0.08, cw4, bh - 0.08, OFF_WHITE)
        tb(slide, cx + 0.10, by + 0.14, cw4 - 0.20, 0.30,
           title, font_size=10, bold=True, color=color,
           align=PP_ALIGN.CENTER)
        tb(slide, cx + 0.10, by + 0.50, cw4 - 0.20, bh - 0.62,
           body, font_size=9, color=NAVY,
           align=PP_ALIGN.CENTER)
        # Arrow connector (not after last)
        if i < 3:
            tb(slide, cx + cw4 + 0.01, by + bh/2 - 0.14, 0.08, 0.28,
               "▶", font_size=11, color=NAVY_MID, align=PP_ALIGN.CENTER)

    # Summary text below
    sy = by + bh + 0.40
    rect(slide, X0, sy - 0.04, CW, 0.02, GRAY_LT)  # divider
    tb(slide, X0, sy + 0.08, CW, 0.30,
       "Doctrine tells us what matters.  Data architecture makes it visible.  "
       "You are learning to build the bridge between them.",
       font_size=11, color=NAVY, align=PP_ALIGN.CENTER, bold=False)
    tb(slide, X0, sy + 0.44, CW, 0.24,
       "Next: The operational landscape you're walking into  →",
       font_size=9, color=NAVY_MID, align=PP_ALIGN.CENTER, italic=True)

    add_footer(slide)


def slide_operational_landscape(slide):
    """Slide 6: THE OPERATIONAL LANDSCAPE — section break with stats."""

    # Left: large section title
    lw = CW * 0.52
    tb(slide, X0, Y0 + 0.30, lw, 1.80,
       ["THE OPERATIONAL", "LANDSCAPE OF", "ARMY DATA"],
       font_size=36, bold=True, color=NAVY, align=PP_ALIGN.LEFT)

    rect(slide, X0, Y0 + 2.30, lw * 0.60, 0.04, GOLD)  # gold divider

    tb(slide, X0, Y0 + 2.50, lw, 0.28,
       "Day One Orientation for Data Practitioners",
       font_size=11, color=NAVY_MID, italic=True)
    tb(slide, X0, Y0 + 2.84, lw, 0.28,
       "Understanding Why You're Here — and Where the Army Is Going",
       font_size=10, color=NAVY)
    tb(slide, X0, Y0 + 3.30, lw * 0.80, 0.20,
       "Source: Army CIO Leo Garciga, AFCEA NOVA Army IT Day 2026",
       font_size=7.5, color=NAVY_MID, italic=True)

    # Right: stat boxes
    rx = X0 + lw + 0.30
    rw = CW - lw - 0.30
    sh = (CH - 0.20) / 2 - 0.10

    rect(slide, rx, Y0, rw, sh, NAVY2)
    rect(slide, rx, Y0, rw, 0.06, GREEN)  # top accent
    tb(slide, rx + 0.10, Y0 + 0.20, rw - 0.20, sh * 0.55,
       "~800", font_size=60, bold=True, color=GREEN,
       align=PP_ALIGN.CENTER, wrap=False)
    tb(slide, rx + 0.10, Y0 + sh * 0.65, rw - 0.20, sh * 0.30,
       "legacy business systems\nare being consolidated",
       font_size=10, color=WHITE, align=PP_ALIGN.CENTER)

    rect(slide, rx, Y0 + sh + 0.20, rw, sh, NAVY2)
    rect(slide, rx, Y0 + sh + 0.20, rw, 0.06, GREEN)
    tb(slide, rx + 0.10, Y0 + sh + 0.40, rw - 0.20, sh * 0.55,
       "100", font_size=60, bold=True, color=GREEN,
       align=PP_ALIGN.CENTER, wrap=False)
    tb(slide, rx + 0.10, Y0 + sh + 0.20 + sh * 0.65, rw - 0.20, sh * 0.30,
       "killed in two quarters",
       font_size=10, color=WHITE, align=PP_ALIGN.CENTER)

    add_footer(slide)


def slide_legacy_sprawl(slide):
    """Slide 7: THE PROBLEM: LEGACY SYSTEM SPRAWL."""

    tb(slide, X0, Y0, CW, 0.22,
       'Army CIO Leo Garciga called legacy business systems the "Achilles heel of every enterprise."',
       font_size=8.5, color=RED, italic=True)

    # 3 stat boxes
    stats = [
        ("42",  "Training &\nReadiness Systems"),
        ("58",  "HR Management\nSystems"),
        ("75",  "Logistics\nSystems"),
    ]
    sw = (CW - 0.04) / 3
    sh = 1.60
    for i, (val, lbl) in enumerate(stats):
        cx = X0 + i * (sw + 0.02)
        rect(slide, cx, Y0 + 0.30, sw, 0.06, RED)
        rect(slide, cx, Y0 + 0.36, sw, sh, OFF_WHITE)
        tb(slide, cx + 0.10, Y0 + 0.44, sw - 0.20, 0.90,
           val, font_size=52, bold=True, color=RED, align=PP_ALIGN.CENTER, wrap=False)
        tb(slide, cx + 0.10, Y0 + 1.36, sw - 0.20, 0.52,
           lbl, font_size=9.5, color=NAVY, align=PP_ALIGN.CENTER)

    # What this looked like panel
    py = Y0 + 0.36 + sh + 0.20
    ph = CH - 0.36 - sh - 0.42
    panel(slide, X0, py, CW, ph,
          "WHAT THIS LOOKED LIKE IN PRACTICE",
          ["▸  Every functional area built its own custom software",
           "▸  Customized process after customized process — many never used again",
           "▸  No shared data layer — systems couldn't talk to each other",
           "▸  Result: massive duplication, zero interoperability, enormous cost"],
          header_fill=DARK_BG, title_color=WHITE, body_size=9.5)

    add_footer(slide)


def slide_army_response(slide):
    """Slide 8: THE ARMY'S RESPONSE: CONSOLIDATE OR DIE."""

    tb(slide, X0, Y0, CW, 0.22,
       "The largest enterprise consolidation push the Army has ever attempted.",
       font_size=8.5, color=NAVY_MID)

    ph = CH - 0.52
    # Left: OLD MODEL
    panel(slide, X0, Y0 + 0.28, CW/2 - 0.15, ph - 0.30,
          "THE OLD MODEL",
          ["▸  ~800 separate business systems",
           "▸  Custom-built for each functional area",
           "▸  Maintained forever by dedicated teams",
           "▸  Data locked in silos",
           "▸  Years to deliver a new capability"],
          header_fill=RED, title_color=WHITE, body_size=9.5)

    # Arrow
    tb(slide, X0 + CW/2 - 0.12, Y0 + 0.28 + ph/2 - 0.32, 0.28, 0.50,
       "→", font_size=36, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    # Right: NEW MODEL
    panel(slide, X0 + CW/2 + 0.15, Y0 + 0.28, CW/2 - 0.15, ph - 0.30,
          "THE NEW MODEL",
          ["▸  Enterprise platforms (shared infrastructure)",
           "▸  Reusable data across the force",
           "▸  Composable applications on top",
           "▸  Low-code / no-code development",
           "▸  Weeks to deliver, not years"],
          header_fill=GREEN, title_color=WHITE, body_size=9.5)

    # Quote footer
    qy = Y0 + 0.28 + ph - 0.26
    tb(slide, X0, qy, CW, 0.24,
       '"We\'ve turned off 100 systems in basically two quarters. Killed them."  — CIO Leo Garciga',
       font_size=8.5, color=NAVY_MID, italic=True, align=PP_ALIGN.CENTER)

    add_footer(slide)


def slide_readiness_example(slide):
    """Slide 9: THIS IS ALREADY WORKING: THE READINESS EXAMPLE."""

    def flow_row(slide, y, boxes, colors, row_h=0.52):
        bw = (CW - (len(boxes)-1) * 0.40) / len(boxes)
        for i, (label, color) in enumerate(zip(boxes, colors)):
            cx = X0 + i * (bw + 0.40)
            rect(slide, cx, y, bw, row_h, color)
            tb(slide, cx + 0.10, y + 0.12, bw - 0.20, row_h - 0.16,
               label, font_size=9.5, bold=False, color=WHITE,
               align=PP_ALIGN.CENTER, wrap=False)
            if i < len(boxes) - 1:
                tb(slide, cx + bw + 0.08, y + row_h/2 - 0.12, 0.24, 0.24,
                   "→", font_size=14, color=NAVY, align=PP_ALIGN.CENTER)

    # BEFORE
    tb(slide, X0, Y0, 0.70, 0.24, "BEFORE", font_size=9, bold=True, color=RED)
    flow_row(slide, Y0 + 0.28,
             ["Unit Data", "Manual Excel", "Staff Builds PPT", "USR Brief"],
             [RED, RED, RED, RED])
    tb(slide, X0, Y0 + 0.88, CW, 0.22,
       "Thousands of man-hours. Static snapshots. Stale by the time it's briefed.",
       font_size=8, color=NAVY_MID, italic=True)

    # AFTER
    tb(slide, X0, Y0 + 1.22, 0.60, 0.24, "AFTER", font_size=9, bold=True, color=GREEN)
    flow_row(slide, Y0 + 1.50,
             ["Live Data Feeds", "Enterprise Platform", "Real-Time Readiness"],
             [GREEN, GREEN, GREEN])

    # Quote block
    qy = Y0 + 2.22
    rect(slide, X0, qy, 0.06, 1.40, NAVY_MID)  # left bar
    rect(slide, X0 + 0.06, qy, CW - 0.06, 1.40, GRAY_LT)
    tb(slide, X0 + 0.18, qy + 0.12, CW - 0.30, 1.16,
       "Since November 2025, the Army's Total Army Readiness Review has shifted from "
       "static presentations built through thousands of man-hours to live data reporting. "
       "Senior leaders can now see the status of every truck, tank, and helicopter in real time.",
       font_size=9, color=NAVY)

    add_footer(slide)


def slide_what_is_foundry(slide):
    """Slide 10: WHAT IS PALANTIR FOUNDRY? — 2×2 grid."""

    tb(slide, X0, Y0, CW, 0.22,
       "The enterprise data and AI platform at the center of the Army's technology stack.",
       font_size=8.5, color=NAVY_MID)

    quadrants = [
        ("DATA INTEGRATION",
         "Connect and transform data from dozens of source systems into clean, unified datasets.",
         NAVY),
        ("ONTOLOGY",
         "Model real-world objects — units, equipment, personnel — as a connected knowledge graph.",
         TEAL),
        ("APPLICATIONS",
         "Build operational apps directly on top of live data — no separate development cycle.",
         BLUE),
        ("ANALYTICS & AI",
         "Run analysis, machine learning, and AI workflows against enterprise-scale data.",
         NAVY_MID),
    ]

    qw = CW / 2 - 0.08
    qh = (CH - 0.32) / 2 - 0.08

    positions = [
        (X0,            Y0 + 0.30),
        (X0 + qw + 0.16, Y0 + 0.30),
        (X0,            Y0 + 0.30 + qh + 0.16),
        (X0 + qw + 0.16, Y0 + 0.30 + qh + 0.16),
    ]

    for (qx, qy), (title, body, color) in zip(positions, quadrants):
        rect(slide, qx, qy, qw, 0.06, color)
        rect(slide, qx, qy + 0.06, qw, qh - 0.06, OFF_WHITE)
        tb(slide, qx + 0.10, qy + 0.12, qw - 0.20, 0.28,
           title, font_size=9.5, bold=True, color=color)
        tb(slide, qx + 0.10, qy + 0.44, qw - 0.20, qh - 0.54,
           body, font_size=9, color=NAVY)

    add_footer(slide)


def slide_cultural_shift(slide):
    """Slide 11: THE CULTURAL SHIFT YOU'RE WALKING INTO — 2-col + quote."""

    ph = CH - 0.52

    # Left: LEGACY MINDSET
    panel(slide, X0, Y0, CW/2 - 0.08, ph,
          "LEGACY MINDSET",
          ['"We need our own custom system"',
           "Data belongs to my organization",
           "Build from scratch every time",
           "Manual reporting and analysis",
           "Protect the system we have"],
          header_fill=RED, title_color=WHITE, body_size=10)

    # Right: ENTERPRISE MINDSET
    panel(slide, X0 + CW/2 + 0.08, Y0, CW/2 - 0.08, ph,
          "ENTERPRISE MINDSET",
          ['"Use the enterprise platform first"',
           "Data is a shared Army asset",
           "Go out-of-the-box whenever possible",
           "Live data, real-time decisions",
           "Converge onto enterprise solutions"],
          header_fill=GREEN, title_color=WHITE, body_size=10)

    # Quote footer
    qy = Y0 + ph + 0.08
    tb(slide, X0, qy, CW, 0.22,
       '"Go as much out of the box as you can."  — CIO Garciga on the Army\'s low-code/no-code mandate',
       font_size=8.5, color=NAVY_MID, italic=True, align=PP_ALIGN.CENTER)

    add_footer(slide)


def slide_where_you_fit(slide):
    """Slide 12: WHERE YOU FIT IN — stack diagram."""

    tb(slide, X0, Y0, CW, 0.22,
       "The enterprise platforms are in place. The missing layer is enterprise DATA ARCHITECTURE.",
       font_size=9, color=TEAL, bold=True)

    # Stack layers
    lw  = CW * 0.66
    lh  = 0.72
    gap = 0.12
    stack = [
        ("SOURCE SYSTEMS",    "ERPs, logistics, HR, readiness feeds",      GRAY_LT,  NAVY),
        ("DATA INTEGRATION",  "Pipelines that clean, transform, and connect data", BLUE_LT, BLUE),
        ("ONTOLOGY",          "Object types, link types, properties — the Army's data model", TEAL_LT, TEAL),
        ("APPLICATIONS",      "Dashboards, workflows, decision-support tools", VIOLET_LT, VIOLET),
        ("DECISIONS",         "Commanders act on real-time, trusted data",   NAVY,     WHITE),
    ]
    you_layers = {1, 2}  # DATA INTEGRATION + ONTOLOGY

    ly = Y0 + 0.32
    for i, (label, desc, bg, fg) in enumerate(stack):
        rect(slide, X0, ly, lw, lh, bg)
        rect(slide, X0, ly, 0.06, lh, fg)
        tb(slide, X0 + 0.12, ly + 0.06, lw * 0.40, 0.26,
           label, font_size=9, bold=True, color=fg)
        tb(slide, X0 + 0.12, ly + 0.32, lw * 0.85, 0.36,
           desc, font_size=8.5, color=fg if fg != WHITE else WHITE)

        if i in you_layers:
            rect(slide, X0 + lw + 0.12, ly, 0.90, lh, NAVY)
            tb(slide, X0 + lw + 0.12, ly + 0.20, 0.90, 0.32,
               "YOU", font_size=18, bold=True, color=WHITE,
               align=PP_ALIGN.CENTER, wrap=False)

        if i < len(stack) - 1:
            tb(slide, X0 + lw/2 - 0.10, ly + lh, 0.20, gap,
               "▼", font_size=9, color=NAVY_MID, align=PP_ALIGN.CENTER)
        ly += lh + gap

    # Summary
    tb(slide, X0 + lw + 0.16, Y0 + 0.32 + 2 * (lh + gap) + lh + 0.20,
       CW - lw - 0.16, 0.56,
       "Your work makes the\nentire stack function.",
       font_size=10, color=NAVY)

    add_footer(slide)


def slide_what_you_learn(slide):
    """Slide 13: WHAT YOU'LL LEARN IN THIS COURSE — 3-col."""

    cols = [
        ("DATA PIPELINES",    NAVY,  ["Connect source systems", "Transform raw data",
                                       "Build reliable ingestion", "Automate data flows"]),
        ("ONTOLOGY DESIGN",   TEAL,  ["Define object types", "Model relationships",
                                       "Build the data layer", "Enable interoperability"]),
        ("OPERATIONAL APPS",  BLUE,  ["Build on live data", "Create decision tools",
                                       "Replace manual reports", "Deliver to the force"]),
    ]

    cw3 = (CW - 0.08) / 3
    for i, (title, color, bullets) in enumerate(cols):
        cx = X0 + i * (cw3 + 0.04)
        rect(slide, cx, Y0, cw3, 0.08, color)  # top accent
        rect(slide, cx, Y0 + 0.08, cw3, CH - 0.30, OFF_WHITE)
        tb(slide, cx + 0.10, Y0 + 0.24, cw3 - 0.20, 0.32,
           title, font_size=11, bold=True, color=color,
           align=PP_ALIGN.CENTER)
        rect(slide, cx + 0.20, Y0 + 0.62, cw3 - 0.40, 0.03, color)
        by = Y0 + 0.78
        for bullet in bullets:
            tb(slide, cx + 0.12, by, cw3 - 0.24, 0.34,
               f"▸  {bullet}", font_size=9.5, color=NAVY)
            by += 0.80

    add_footer(slide)


def slide_your_mission(slide):
    """Slide 14: YOUR MISSION — centered closing slide."""

    # Centered mission title
    rect(slide, X0, Y0 + 0.30, CW, 0.04, GOLD)  # top rule
    tb(slide, X0, Y0 + 0.60, CW, 0.80,
       "YOUR MISSION", font_size=44, bold=True, color=NAVY,
       align=PP_ALIGN.CENTER, wrap=False)
    rect(slide, X0 + CW/4, Y0 + 1.46, CW/2, 0.04, NAVY)  # divider

    tb(slide, X0, Y0 + 1.66, CW, 0.52,
       ["Build the data architecture that turns",
        "enterprise platforms into enterprise capability."],
       font_size=14, color=NAVY, align=PP_ALIGN.CENTER)

    # 3 checkmarks
    checks = [
        "The Army is consolidating hundreds of systems onto enterprise platforms — this is irreversible.",
        "Palantir Foundry is the data and AI backbone. Your work here is foundational.",
        "Without strong data architecture, none of the enterprise vision works. That's your job.",
    ]
    cy = Y0 + 2.50
    for chk in checks:
        rect(slide, X0 + 1.0, cy, 0.42, 0.42, GREEN)
        tb(slide, X0 + 1.0, cy + 0.06, 0.42, 0.30,
           "✓", font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        tb(slide, X0 + 1.54, cy + 0.04, CW - 2.54, 0.36,
           chk, font_size=10, color=NAVY)
        cy += 0.68

    add_footer(slide)


# ══════════════════════════════════════════════════════════════════════════════
# Main
# ══════════════════════════════════════════════════════════════════════════════

# (header text, builder function) — one entry per content slide
SLIDES = [
    ("FIRST PRINCIPLES",                          slide_first_principles),
    ("THE OPERATIONAL ENVIRONMENT",               slide_operational_environment),
    ("DATA  ·  DECISIONS  ·  DOCTRINE",           slide_data_decisions_doctrine),
    ("WHAT MAKES US DIFFERENT",                   slide_what_makes_us_different),
    ("SO WHERE DOES DATA FIT?",                   slide_where_does_data_fit),
    ("THE OPERATIONAL LANDSCAPE OF ARMY DATA",    slide_operational_landscape),
    ("THE PROBLEM: LEGACY SYSTEM SPRAWL",         slide_legacy_sprawl),
    ("THE ARMY'S RESPONSE: CONSOLIDATE OR DIE",   slide_army_response),
    ("THIS IS ALREADY WORKING: THE READINESS EXAMPLE", slide_readiness_example),
    ("WHAT IS PALANTIR FOUNDRY?",                 slide_what_is_foundry),
    ("THE CULTURAL SHIFT YOU'RE WALKING INTO",    slide_cultural_shift),
    ("WHERE YOU FIT IN",                          slide_where_you_fit),
    ("WHAT YOU'LL LEARN IN THIS COURSE",          slide_what_you_learn),
    ("YOUR MISSION",                              slide_your_mission),
]


def main():
    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    shutil.copy(TEMPLATE, OUTPUT)
    prs = Presentation(str(OUTPUT))

    # Title slide (template slide 1 — title layout only)
    update_title_chrome(
        prs.slides[0],
        ["ARMY DATA ORIENTATION",
         "Day One Orientation for Data Practitioners"],
        "USAREUR-AF C2DAO  |  Understanding Why You're Here — and Where the Army Is Going",
    )

    # Clone template slide 2 (content layout) for all 14 content slides.
    # Clone BEFORE any modification so every copy starts from the clean template.
    for _ in range(len(SLIDES) - 1):   # -1 because slide 2 already exists
        clone_content_slide(prs)

    # Setup chrome + build content on each content slide
    for i, (header, builder) in enumerate(SLIDES):
        slide = prs.slides[i + 1]      # slides[1] … slides[14]
        setup_slide(slide, header)
        builder(slide)

    prs.save(str(OUTPUT))
    print(f"Saved → {OUTPUT}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
