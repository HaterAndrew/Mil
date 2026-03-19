"""
Build a single-slide PowerPoint summary of the USAREUR-AF Operational Data Team project.
Focus: TM-10/20/30 baseline tiers and their criticality to command-wide Maven adoption.
Output: maven_training/pdf/MSS_Project_Overview.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# ── Palette ───────────────────────────────────────────────────────────────────
OD_GREEN    = RGBColor(0x4B, 0x5B, 0x2E)   # dark OD green (background)
DARK_GREEN  = RGBColor(0x32, 0x3D, 0x1C)   # darker panel
MID_GREEN   = RGBColor(0x6B, 0x7F, 0x3E)   # box fill
GOLD        = RGBColor(0xC8, 0xA0, 0x32)   # Army gold
GOLD_DARK   = RGBColor(0x9A, 0x78, 0x22)   # darker gold for borders
OFF_WHITE   = RGBColor(0xF0, 0xED, 0xE0)   # parchment text
RED_ALERT   = RGBColor(0x8B, 0x1A, 0x1A)   # risk/alert color
BLUE_ACCENT = RGBColor(0x2A, 0x4A, 0x6B)   # specialist tier accent


def add_rect(slide, left, top, width, height, fill_color,
             line_color=None, line_width_pt=0):
    shape = slide.shapes.add_shape(
        1,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width_pt)
    else:
        shape.line.fill.background()
    return shape


def tb(slide, text, left, top, width, height,
       size=10, bold=False, color=OFF_WHITE,
       align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Arial"
    return txb


def tier_box(slide, level, audience, headline, bullets,
             left, top, width, height,
             header_color, header_text_color=None):
    """Full tier box: colored header strip + dark body."""
    if header_text_color is None:
        header_text_color = OD_GREEN

    HEADER_H = 0.62

    # Header fill
    add_rect(slide, left, top, width, HEADER_H, header_color)
    # Body fill
    add_rect(slide, left, top + HEADER_H, width, height - HEADER_H,
             DARK_GREEN, header_color, 1.2)

    # Level label (large)
    tb(slide, level,
       left + 0.08, top + 0.04, width * 0.35, 0.32,
       size=18, bold=True, color=header_text_color)

    # Audience
    tb(slide, audience,
       left + 0.08, top + 0.36, width - 0.16, 0.22,
       size=7.5, bold=False, color=header_text_color)

    # Headline (gold, inside body)
    tb(slide, headline,
       left + 0.1, top + HEADER_H + 0.06, width - 0.2, 0.25,
       size=8.5, bold=True, color=GOLD)

    # Bullets
    body = "\n".join(bullets)
    tb(slide, body,
       left + 0.1, top + HEADER_H + 0.32, width - 0.2, height - HEADER_H - 0.36,
       size=7.5, bold=False, color=OFF_WHITE)


# ── Slide setup ───────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

slide = prs.slides.add_slide(prs.slide_layouts[6])   # blank

# Background
add_rect(slide, 0, 0, 13.33, 7.5, OD_GREEN)

# ── Top bar ───────────────────────────────────────────────────────────────────
add_rect(slide, 0, 0, 13.33, 0.68, GOLD)

tb(slide, "USAREUR-AF  |  Maven Smart System — Command-Wide Adoption Strategy",
   0.15, 0.05, 10.0, 0.34,
   size=13, bold=True, color=OD_GREEN, align=PP_ALIGN.LEFT)

tb(slide, "Training Infrastructure: TM-10 through TM-30  |  The Baseline That Makes Maven Work",
   0.15, 0.38, 10.0, 0.26,
   size=9, bold=False, color=DARK_GREEN, align=PP_ALIGN.LEFT)

tb(slide, "UNCLASSIFIED  |  March 2026",
   10.0, 0.12, 3.15, 0.44,
   size=8.5, bold=False, color=OD_GREEN, align=PP_ALIGN.RIGHT)

# ── Bottom bar ────────────────────────────────────────────────────────────────
add_rect(slide, 0, 7.17, 13.33, 0.33, GOLD)
tb(slide, "USAREUR-AF Operational Data Team  ·  Wiesbaden, Germany  ·  UNCLASSIFIED",
   0.15, 7.19, 13.0, 0.26,
   size=7.5, bold=False, color=OD_GREEN, align=PP_ALIGN.CENTER)

# ── BLUF ──────────────────────────────────────────────────────────────────────
add_rect(slide, 0.2, 0.76, 12.93, 0.38, DARK_GREEN, GOLD, 0.75)
tb(slide,
   "BLUF:  Maven adoption fails at the unit level when personnel cannot access, navigate, or trust the platform. "
   "TM-10/20/30 are the command's solution — doctrine-based, zero-prerequisite training that builds "
   "the data-literate formation Maven requires to deliver operational value.",
   0.36, 0.78, 12.6, 0.34,
   size=8.5, bold=False, color=GOLD, align=PP_ALIGN.LEFT)

# ── Section label ─────────────────────────────────────────────────────────────
tb(slide, "THE BASELINE TIERS  (mandatory for all personnel)",
   0.2, 1.22, 7.2, 0.24,
   size=8, bold=True, color=GOLD, align=PP_ALIGN.LEFT)

tb(slide, "SPECIALIST / ADVANCED",
   7.8, 1.22, 5.3, 0.24,
   size=8, bold=True, color=RGBColor(0xA0, 0xB8, 0xD0), align=PP_ALIGN.LEFT)

# ── Divider line between sections ─────────────────────────────────────────────
add_rect(slide, 7.7, 1.18, 0.04, 4.72, GOLD)

# ── TM-10 tier box ────────────────────────────────────────────────────────────
tier_box(slide,
    level="TM-10",
    audience="ALL personnel — every Soldier and civilian",
    headline="WHY IT'S CRITICAL:  No access = no adoption.",
    bullets=[
        "Covers: CAC login, project navigation, Workshop apps,",
        "  authorized Actions, Contour/Quiver, AIP interfaces",
        "  classification markings, export procedures",
        "",
        "Without TM-10, personnel cannot use Maven at all.",
        "Every unit-level use case begins here.",
        "Prerequisite to all other levels.",
        "",
        "Go/No-Go: 10-task practical evaluation",
        "Duration: 1 day (8 hours) | No prerequisites",
    ],
    left=0.2, top=1.46, width=2.44, height=3.84,
    header_color=MID_GREEN, header_text_color=OFF_WHITE)

# ── TM-20 tier box ────────────────────────────────────────────────────────────
tier_box(slide,
    level="TM-20",
    audience="ALL staff — light builders",
    headline="WHY IT'S CRITICAL:  Units must own their data products.",
    bullets=[
        "Covers: Pipeline Builder (visual, no-code), Workshop",
        "  dashboards, forms, Object Types, Actions,",
        "  project organization, access management,",
        "  branching, USAREUR-AF naming standards",
        "",
        "Without TM-20, the command depends on a small",
        "specialist cadre for every dashboard and form.",
        "TM-20 distributes capability to every staff section.",
        "",
        "Go/No-Go: Build a functional pipeline + dashboard",
        "Duration: 5 days (40 hours) | Prereq: TM-10",
    ],
    left=2.74, top=1.46, width=2.44, height=3.84,
    header_color=MID_GREEN, header_text_color=OFF_WHITE)

# ── TM-30 tier box ────────────────────────────────────────────────────────────
tier_box(slide,
    level="TM-30",
    audience="Data-adjacent specialists (17/25-series, S6, G2)",
    headline="WHY IT'S CRITICAL:  Maven's power is in complex, connected products.",
    bullets=[
        "Covers: Multi-source pipelines, complex Workshop apps,",
        "  Ontology architecture (no-code), advanced Contour/Quiver,",
        "  AIP Logic configuration, data lineage review,",
        "  governance + data steward workflows, C2DAO standards",
        "",
        "Without TM-30, the command cannot build the",
        "operational products that Maven was procured for:",
        "readiness tracking, logistics dashboards, ISR integration.",
        "",
        "Go/No-Go: Multi-source pipeline + integrated app",
        "Duration: 5 days (40 hours) | Prereq: TM-10, TM-20",
    ],
    left=5.28, top=1.46, width=2.32, height=3.84,
    header_color=MID_GREEN, header_text_color=OFF_WHITE)

# ── Specialist panel (right side) ─────────────────────────────────────────────
add_rect(slide, 7.74, 1.46, 5.39, 1.72, BLUE_ACCENT, GOLD_DARK, 0.75)
tb(slide, "TM-40 SERIES  —  12 Tracks (WFF A–F + Specialist G–M)",
   7.86, 1.50, 5.1, 0.22,
   size=8.5, bold=True, color=GOLD)
tb(slide,
   "WFF: 40A Intel  ·  40B Fires  ·  40C M&M  ·  40D Sust  ·  40E Prot  ·  40F MC\n"
   "Specialist: 40G ORSA  ·  40H AI Eng  ·  40M MLE  ·  40J PM  ·  40K KM  ·  40L SWE\n\n"
   "Prereq: TM-30. Role-assigned. Each track: Technical Manual + Concepts Guide + Syllabus.\n"
   "Produces the analysts, engineers, and PMs who build the platform's most capable products.",
   7.86, 1.74, 5.1, 1.38,
   size=7.5, bold=False, color=OFF_WHITE)

add_rect(slide, 7.74, 3.26, 5.39, 1.30, BLUE_ACCENT, GOLD_DARK, 0.75)
tb(slide, "TM-50 SERIES  —  6 Advanced Tracks",
   7.86, 3.30, 5.1, 0.22,
   size=8.5, bold=True, color=GOLD)
tb(slide,
   "Advanced mastery post-TM-40, per role. 6 tracks (ORSA, AI Eng, MLE, PM, KM, SWE).\n\n"
   "Produces the command's subject matter experts — the personnel who sustain, extend,\n"
   "and govern the platform long-term.",
   7.86, 3.54, 5.1, 0.96,
   size=7.5, bold=False, color=OFF_WHITE)

add_rect(slide, 7.74, 4.64, 5.39, 0.66, DARK_GREEN, GOLD_DARK, 0.75)
tb(slide, "TRAINING MANAGEMENT PACKAGE (READY TO EXECUTE)",
   7.86, 4.67, 5.1, 0.20,
   size=7.5, bold=True, color=GOLD)
tb(slide,
   "MTP · POI · CAD · TEO · Annual Schedule · Enrollment SOP · Lesson Plans · Syllabi · Policy Letter",
   7.86, 4.88, 5.1, 0.36,
   size=7.5, bold=False, color=OFF_WHITE)

# ── Adoption risk bar ─────────────────────────────────────────────────────────
add_rect(slide, 0.2, 5.38, 7.44, 0.60, RED_ALERT, RGBColor(0xFF, 0x60, 0x60), 0.75)
tb(slide, "ADOPTION RISK WITHOUT THIS TRAINING",
   0.32, 5.40, 7.2, 0.20,
   size=8, bold=True, color=RGBColor(0xFF, 0xC0, 0xC0))
tb(slide,
   "Personnel cannot log in or navigate  ·  Staff sections depend on 1–2 specialists for every product  "
   "·  Commanders cannot evaluate or trust data outputs  ·  Platform investment goes unrealized",
   0.32, 5.59, 7.2, 0.34,
   size=7.5, bold=False, color=OFF_WHITE)

# ── Metrics strip ─────────────────────────────────────────────────────────────
add_rect(slide, 0.2, 6.06, 7.44, 0.46, DARK_GREEN, GOLD, 0.5)

metrics = [
    ("3", "Mandatory\nBaseline TMs"),
    ("100%", "MOS\nCoverage"),
    ("88 hrs", "TM-10→30\nTotal Instruction"),
    ("3-tier", "Progression\n(No-Code → Code)"),
    ("9", "Syllabi\nReady"),
    ("33", "PDFs\nPublished"),
]
col_w = 7.44 / len(metrics)
for i, (num, label) in enumerate(metrics):
    x = 0.2 + i * col_w
    tb(slide, num,
       x + 0.04, 6.08, col_w - 0.08, 0.20,
       size=11, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    tb(slide, label,
       x + 0.04, 6.26, col_w - 0.08, 0.22,
       size=6, bold=False, color=OFF_WHITE, align=PP_ALIGN.CENTER)

# ── Next steps ────────────────────────────────────────────────────────────────
add_rect(slide, 7.74, 5.38, 5.39, 1.14, DARK_GREEN, GOLD, 0.75)
tb(slide, "NEXT STEPS — COMMAND ACTION REQUIRED",
   7.86, 5.41, 5.1, 0.20,
   size=8, bold=True, color=GOLD)
tb(slide,
   "① Issue training policy letter (template ready)\n"
   "② Certify instructors per Faculty Development Plan\n"
   "③ Provision MSS accounts — prerequisite to TM-10 execution\n"
   "④ Execute first TM-10 cohort — all personnel, no exceptions",
   7.86, 5.62, 5.1, 0.84,
   size=8, bold=False, color=OFF_WHITE)

# ── Save ──────────────────────────────────────────────────────────────────────
out_dir = os.path.join(os.path.dirname(__file__), "..", "maven_training", "pdf")
os.makedirs(out_dir, exist_ok=True)
out_path = os.path.join(out_dir, "MSS_Project_Overview.pptx")
prs.save(out_path)
print(f"Saved: {out_path}")
