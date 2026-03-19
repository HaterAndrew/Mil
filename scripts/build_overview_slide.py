"""
Build a single-slide PowerPoint overview for the MSS Training Program.
Output: maven_training/MSS_PROGRAM_OVERVIEW.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Army color palette ────────────────────────────────────────────────────────
ARMY_GREEN   = RGBColor(0x4B, 0x5A, 0x2A)   # OD green
ARMY_TAN     = RGBColor(0xC8, 0xB5, 0x60)   # desert tan
ARMY_BLACK   = RGBColor(0x1A, 0x1A, 0x1A)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY   = RGBColor(0xF2, 0xF2, 0xF2)
MID_GRAY     = RGBColor(0xCC, 0xCC, 0xCC)
GOLD         = RGBColor(0xC8, 0xA0, 0x00)
DARK_GREEN   = RGBColor(0x2E, 0x38, 0x18)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def rgb(r, g, b):
    return RGBColor(r, g, b)


def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=Pt(0)):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.line.width = line_width
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, text, left, top, width, height,
                font_size=Pt(11), bold=False, color=ARMY_BLACK,
                align=PP_ALIGN.LEFT, italic=False, word_wrap=True):
    txb = slide.shapes.add_textbox(left, top, width, height)
    txb.word_wrap = word_wrap
    tf = txb.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def add_text_in_shape(shape, lines, font_size=Pt(10), bold=False, color=WHITE, align=PP_ALIGN.LEFT):
    """Set text in a shape's text frame, one paragraph per line."""
    tf = shape.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = font_size
        run.font.bold = bold
        run.font.color.rgb = color


# ── Build presentation ────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]  # completely blank
slide = prs.slides.add_slide(blank_layout)

# ── Background ────────────────────────────────────────────────────────────────
bg = add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_color=LIGHT_GRAY)

# ── Header banner ─────────────────────────────────────────────────────────────
HEADER_H = Inches(1.05)
hdr = add_rect(slide, 0, 0, SLIDE_W, HEADER_H, fill_color=ARMY_GREEN)

# Gold accent stripe under header
add_rect(slide, 0, HEADER_H - Inches(0.04), SLIDE_W, Inches(0.04), fill_color=ARMY_TAN)

# Title
add_textbox(slide,
    "MAVEN SMART SYSTEM (MSS) — TRAINING PROGRAM",
    Inches(0.25), Inches(0.08),
    Inches(9), Inches(0.55),
    font_size=Pt(26), bold=True, color=WHITE, align=PP_ALIGN.LEFT)

# Subtitle
add_textbox(slide,
    "USAREUR-AF C2DAO  ·  Wiesbaden, Germany",
    Inches(0.25), Inches(0.62),
    Inches(9), Inches(0.35),
    font_size=Pt(13), bold=False, color=ARMY_TAN, align=PP_ALIGN.LEFT)



# ── BODY — three columns ──────────────────────────────────────────────────────
BODY_TOP  = HEADER_H + Inches(0.18)
BODY_BOT  = SLIDE_H  - Inches(0.55)   # leave room for footer
BODY_H    = BODY_BOT - BODY_TOP
COL_GAP   = Inches(0.18)
COL_W     = (SLIDE_W - Inches(0.35)*2 - COL_GAP*2) / 3
COL1_L    = Inches(0.35)
COL2_L    = COL1_L + COL_W + COL_GAP
COL3_L    = COL2_L + COL_W + COL_GAP

CARD_RADIUS = 0    # python-pptx rectangles only; keep square

def add_card(slide, left, top, width, height, header_text="", header_color=ARMY_GREEN):
    """Draw a white card with a colored header bar."""
    # Card body
    card = add_rect(slide, left, top, width, height,
                    fill_color=WHITE, line_color=MID_GRAY, line_width=Pt(0.75))
    # Header strip
    hdr_h = Inches(0.32)
    h = add_rect(slide, left, top, width, hdr_h, fill_color=header_color)
    return card, hdr_h

# ─── COLUMN 1: What It Is + Learning Path ────────────────────────────────────
c1_top = BODY_TOP

# WHAT IT IS card
card1a_h = Inches(1.65)
add_card(slide, COL1_L, c1_top, COL_W, card1a_h)
add_textbox(slide, "WHAT IT IS",
    COL1_L + Inches(0.08), c1_top + Inches(0.05),
    COL_W - Inches(0.1), Inches(0.26),
    font_size=Pt(10), bold=True, color=WHITE)
add_textbox(slide,
    "A structured curriculum — three sequential core courses for all personnel, "
    "followed by role-specific specialty tracks — that builds every USAREUR-AF "
    "soldier, officer, and civilian into a data-capable user of the Maven Smart "
    "System (Palantir Foundry as the Army's operational data platform).",
    COL1_L + Inches(0.1), c1_top + Inches(0.36),
    COL_W - Inches(0.15), Inches(1.22),
    font_size=Pt(10.5), color=ARMY_BLACK)

# LEARNING PATH card
c1_path_top = c1_top + card1a_h + Inches(0.15)
card1b_h = BODY_BOT - c1_path_top
add_card(slide, COL1_L, c1_path_top, COL_W, card1b_h, header_color=DARK_GREEN)
add_textbox(slide, "LEARNING PATH",
    COL1_L + Inches(0.08), c1_path_top + Inches(0.05),
    COL_W - Inches(0.1), Inches(0.26),
    font_size=Pt(10), bold=True, color=WHITE)

path_text = (
    "TM-10  Maven User\n"
    "All personnel — navigate MSS, read\n"
    "dashboards, verify data currency\n"
    "\n"
    "      ↓  prereq: TM-10\n"
    "\n"
    "TM-20  Builder\n"
    "All staff — dashboards, forms,\n"
    "basic pipelines (no-code / low-code)\n"
    "\n"
    "      ↓  prereq: TM-20\n"
    "\n"
    "TM-30  Advanced Builder\n"
    "Data-adjacent specialists, unit data\n"
    "leads — pipeline authoring, governance,\n"
    "production promotion\n"
    "\n"
    "      ↓  prereq: TM-30\n"
    "\n"
    "TM-40  Role-Specific Tracks\n"
    "WFF Functional (A–F): INT · FIRES\n"
    "M&M · SUST · PROT · Mission Cmd\n"
    "Technical Specialist (G–M)\n"
    "\n"
    "      ↓  (specialists only)\n"
    "\n"
    "TM-50  Advanced / Expert"
)
add_textbox(slide, path_text,
    COL1_L + Inches(0.1), c1_path_top + Inches(0.38),
    COL_W - Inches(0.15), card1b_h - Inches(0.45),
    font_size=Pt(9.5), color=ARMY_BLACK)

# ─── COLUMN 2: WFF Tracks + Training Environment ────────────────────────────
c2_top = BODY_TOP

# WFF TRACKS card
card2a_h = Inches(2.55)
add_card(slide, COL2_L, c2_top, COL_W, card2a_h, header_color=ARMY_GREEN)
add_textbox(slide, "TM-40 WFF FUNCTIONAL TRACKS (A–F)",
    COL2_L + Inches(0.08), c2_top + Inches(0.05),
    COL_W - Inches(0.1), Inches(0.26),
    font_size=Pt(10), bold=True, color=WHITE)

wff_rows = [
    ("TM-40A", "Intelligence"),
    ("TM-40B", "Fires"),
    ("TM-40C", "Movement & Maneuver"),
    ("TM-40D", "Sustainment"),
    ("TM-40E", "Protection"),
    ("TM-40F", "Mission Command"),
]
row_top = c2_top + Inches(0.38)
row_h   = Inches(0.29)
for i, (code, name) in enumerate(wff_rows):
    bg_color = LIGHT_GRAY if i % 2 == 0 else WHITE
    add_rect(slide,
        COL2_L + Inches(0.08), row_top + i * row_h,
        COL_W - Inches(0.15), row_h,
        fill_color=bg_color)
    add_textbox(slide, code,
        COL2_L + Inches(0.12), row_top + i * row_h + Inches(0.04),
        Inches(0.75), row_h,
        font_size=Pt(9.5), bold=True, color=ARMY_GREEN)
    add_textbox(slide, name,
        COL2_L + Inches(0.9), row_top + i * row_h + Inches(0.04),
        COL_W - Inches(1.0), row_h,
        font_size=Pt(9.5), color=ARMY_BLACK)

prereq_note_top = row_top + len(wff_rows) * row_h + Inches(0.06)
add_textbox(slide, "Prereq: TM-30  ·  No coding required  ·  Audience: functional WFF staff",
    COL2_L + Inches(0.1), prereq_note_top,
    COL_W - Inches(0.15), Inches(0.32),
    font_size=Pt(8.5), italic=True, color=RGBColor(0x55, 0x55, 0x55))

# TRAINING ENVIRONMENT card
c2_env_top = c2_top + card2a_h + Inches(0.15)
card2b_h = Inches(1.55)
add_card(slide, COL2_L, c2_env_top, COL_W, card2b_h, header_color=DARK_GREEN)
add_textbox(slide, "TRAINING ENVIRONMENT",
    COL2_L + Inches(0.08), c2_env_top + Inches(0.05),
    COL_W - Inches(0.1), Inches(0.26),
    font_size=Pt(10), bold=True, color=WHITE)
add_textbox(slide,
    "All training conducted in a dedicated MSS Training Environment "
    "with synthetic operational data. No training in the production "
    "MSS instance. Trainees never ingest, modify, or export production "
    "operational data.",
    COL2_L + Inches(0.1), c2_env_top + Inches(0.38),
    COL_W - Inches(0.15), Inches(1.1),
    font_size=Pt(10), color=ARMY_BLACK)

# SCOPE & SCALE card
c2_scope_top = c2_env_top + card2b_h + Inches(0.15)
card2c_h = BODY_BOT - c2_scope_top
add_card(slide, COL2_L, c2_scope_top, COL_W, card2c_h, header_color=ARMY_GREEN)
add_textbox(slide, "SCOPE & SCALE",
    COL2_L + Inches(0.08), c2_scope_top + Inches(0.05),
    COL_W - Inches(0.1), Inches(0.26),
    font_size=Pt(10), bold=True, color=WHITE)

bullets = [
    "3 core courses (TM-10/20/30) + specialty tracks for all roles",
    "160+ published PDFs; Markdown source in version control",
    "Dedicated exercises with hands-on labs per track",
    "PRE + POST assessments for every course",
    "Authority: C2DAO Data Governance Directive (MSS-POI-001)",
    "Review cycle: annual or on major platform update",
]
for i, b in enumerate(bullets):
    add_textbox(slide, f"▪  {b}",
        COL2_L + Inches(0.1), c2_scope_top + Inches(0.34) + i * Inches(0.155),
        COL_W - Inches(0.15), Inches(0.18),
        font_size=Pt(8), color=ARMY_BLACK)

# ─── COLUMN 3: Supporting Publications + Doctrine ────────────────────────────
c3_top = BODY_TOP

# SUPPORTING PUBLICATIONS card
card3a_h = Inches(3.35)
add_card(slide, COL3_L, c3_top, COL_W, card3a_h, header_color=ARMY_GREEN)
add_textbox(slide, "SUPPORTING PUBLICATIONS",
    COL3_L + Inches(0.08), c3_top + Inches(0.05),
    COL_W - Inches(0.1), Inches(0.26),
    font_size=Pt(10), bold=True, color=WHITE)

pub_rows = [
    ("Doctrine",        "Data Literacy (Sr Leaders + Technical), Foundry Glossary, CDA Constraints & Directives"),
    ("Syllabi",         "Per-course learning objectives, task lists, and evaluation criteria (all 18 courses)"),
    ("Exercises",       "Hands-on lab guides + environment setup per track"),
    ("Assessments",     "PRE + POST exams for every course"),
    ("Training Mgmt",   "POI · MTP · CAD · TEO · Enrollment SOP · Annual Schedule · Faculty Dev Plan"),
    ("Standards",       "Naming & governance standards; dependency map (cross-reference index)"),
]

pub_top = c3_top + Inches(0.38)
pub_row_h = Inches(0.47)
for i, (cat, desc) in enumerate(pub_rows):
    bg_color = LIGHT_GRAY if i % 2 == 0 else WHITE
    add_rect(slide,
        COL3_L + Inches(0.08), pub_top + i * pub_row_h,
        COL_W - Inches(0.15), pub_row_h,
        fill_color=bg_color)
    add_textbox(slide, cat,
        COL3_L + Inches(0.12), pub_top + i * pub_row_h + Inches(0.03),
        Inches(1.0), Inches(0.22),
        font_size=Pt(9), bold=True, color=ARMY_GREEN)
    add_textbox(slide, desc,
        COL3_L + Inches(0.12), pub_top + i * pub_row_h + Inches(0.22),
        COL_W - Inches(0.22), Inches(0.24),
        font_size=Pt(8.5), color=ARMY_BLACK)

# DOCTRINE card
c3_doc_top = c3_top + card3a_h + Inches(0.15)
card3b_h = BODY_BOT - c3_doc_top
add_card(slide, COL3_L, c3_doc_top, COL_W, card3b_h, header_color=DARK_GREEN)
add_textbox(slide, "DOCTRINE PUBLICATIONS",
    COL3_L + Inches(0.08), c3_doc_top + Inches(0.05),
    COL_W - Inches(0.1), Inches(0.26),
    font_size=Pt(10), bold=True, color=WHITE)

doc_items = [
    ("Data Literacy — Senior Leaders",
     "Strategic framing for O-5+, CSM+, and senior civilians. Principles, not procedures."),
    ("Data Literacy — Technical Reference",
     "Foundation data concepts for all personnel. Recommended pre-reading before TM-10."),
    ("Foundry Glossary",
     "Army data terms mapped to Palantir Foundry equivalents. All personnel."),
    ("CDA Constraints & Directives",
     "Authoritative constraints on data architecture decisions. Required for TM-30+ and all specialists."),
]
di_top = c3_doc_top + Inches(0.38)
di_h   = Inches(0.52)
for i, (title, desc) in enumerate(doc_items):
    add_textbox(slide, title,
        COL3_L + Inches(0.12), di_top + i * di_h,
        COL_W - Inches(0.2), Inches(0.22),
        font_size=Pt(9.5), bold=True, color=ARMY_GREEN)
    add_textbox(slide, desc,
        COL3_L + Inches(0.12), di_top + i * di_h + Inches(0.21),
        COL_W - Inches(0.2), Inches(0.28),
        font_size=Pt(8.8), color=ARMY_BLACK)

# ── Footer ────────────────────────────────────────────────────────────────────
FOOTER_TOP = SLIDE_H - Inches(0.45)
add_rect(slide, 0, FOOTER_TOP, SLIDE_W, Inches(0.45), fill_color=ARMY_GREEN)
add_rect(slide, 0, FOOTER_TOP, SLIDE_W, Inches(0.04), fill_color=ARMY_TAN)



# Center: org line (width capped to avoid overlapping right-aligned distribution text)
add_textbox(slide, "USAREUR-AF  ·  C2DAO  ·  MSS-POI-001  ·  March 2026",
    Inches(2.2), FOOTER_TOP + Inches(0.08),
    Inches(6.1), Inches(0.32),
    font_size=Pt(9), color=WHITE, align=PP_ALIGN.CENTER)

# Right: distribution note
add_textbox(slide, "Distribution: U.S. Government agencies and contractors",
    SLIDE_W - Inches(4.8), FOOTER_TOP + Inches(0.08),
    Inches(4.6), Inches(0.32),
    font_size=Pt(7), italic=True, color=ARMY_TAN, align=PP_ALIGN.RIGHT)

# ── Save ──────────────────────────────────────────────────────────────────────
out = "maven_training/MSS_PROGRAM_OVERVIEW.pptx"
prs.save(out)
print(f"Saved: {out}")
