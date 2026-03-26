#!/usr/bin/env python3
"""
Convert all PNG slide decks to PPTX using USAREUR-AF template.

Each deck directory under source_material/course_portal/assets/ that contains
slide-NNN.png files gets converted to a .pptx:
  - Slide 1: title card with deck name, track, and TM level from manifest
  - Slides 2+: content slides (cloned from template slide 2) with PNG embedded

Metadata sourced from: source_material/course_portal/manifest.json

Run from repo root: python3 scripts/build_png_decks.py
Output: maven_training/pdf/{DECK_NAME}.pptx
"""

import copy
import json
import shutil
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Paths ──────────────────────────────────────────────────────────────────────
TEMPLATE    = Path("/home/dale/Downloads/USAREUR-AF PPT Template.pptx")
SOURCE_DIR  = Path("maven_training/source_material/course_portal/assets")
MANIFEST      = Path("maven_training/source_material/course_portal/manifest.json")
ALIGNMENT_MAP = Path("maven_training/source_material/course_portal/alignment_map.json")
OUTPUT_DIR  = Path("maven_training/pdf")

# army_data_orientation already has a hand-crafted script; skip it here
SKIP_DIRS = {"army_data_orientation_v1"}


def load_manifest() -> dict:
    """Load manifest.json and return a dict keyed by normalised id (underscores)."""
    entries = json.loads(MANIFEST.read_text())
    return {e["id"].replace(" ", "_"): e for e in entries}


def load_alignment() -> dict:
    """Load alignment_map.json and merge into manifest metadata.
    Alignment map carries enriched titles, correct track labels, sequence order,
    and also_supports cross-refs produced by the portal structure agent.
    """
    base     = load_manifest()
    enriched = json.loads(ALIGNMENT_MAP.read_text())
    for entry in enriched:
        key = entry["id"].replace(" ", "_")
        if key in base:
            base[key].update(entry)          # alignment map wins on overlapping keys
        else:
            base[key] = entry
    return base

# ── Palette ────────────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x17, 0x32, 0x5C)
PURPLE = RGBColor(0x70, 0x30, 0xA0)

CLASSIFICATION = ""

# ── Content area geometry ──────────────────────────────────────────────────────
IMG_X = Inches(0.15)
IMG_Y = Inches(0.70)   # below header strip
IMG_W = Inches(13.03)
IMG_H = Inches(6.50)   # above bottom banner


# ── Helpers ────────────────────────────────────────────────────────────────────

def deck_title(dir_name: str) -> str:
    """Convert directory name to a readable all-caps title (fallback only)."""
    return dir_name.replace("_", " ").replace("-", " ").upper()


def remove_shape(shape):
    el = shape._element
    el.getparent().remove(el)


def clone_content_slide(prs: Presentation):
    """Clone template slide 2 (content layout) as a new slide.
    Must be called BEFORE slide 2 is modified so every copy is a clean clone.
    """
    src = prs.slides[1]
    slide = prs.slides.add_slide(src.slide_layout)
    sp_tree = slide.shapes._spTree
    # Strip bare stub shapes (keep only nvGrpSpPr + grpSpPr at indices 0,1)
    for child in list(sp_tree)[2:]:
        sp_tree.remove(child)
    # Deep-copy all shapes from template slide 2
    for child in list(src.shapes._spTree)[2:]:
        sp_tree.append(copy.deepcopy(child))
    return slide


def setup_slide(slide, header_text: str):
    """Update header and classification text on a cloned content slide.
    Also removes the default large content placeholder (TextBox 8).
    """
    to_remove = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        txt  = shape.text_frame.text.strip()
        name = shape.name
        if name == "TextBox 2" or txt == "Header":
            tf = shape.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = header_text
            run.font.size = Pt(14)
            run.font.bold = True
            run.font.color.rgb = NAVY
        elif txt in ("CUI", "UNCLASSIFIED", "Classification"):
            for p in shape.text_frame.paragraphs:
                for run in p.runs:
                    run.text = CLASSIFICATION
        elif name == "TextBox 8":
            to_remove.append(shape)
    for s in to_remove:
        remove_shape(s)
    # Bottom classification banner
    bot = slide.shapes.add_textbox(Inches(0), Inches(7.25), Inches(13.33), Inches(0.25))
    tf = bot.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = CLASSIFICATION
    run.font.size = Pt(9)
    run.font.bold = True
    run.font.color.rgb = PURPLE


def update_title_chrome(slide, title_lines: list[str], poc: str):
    """Populate the template title slide (slide 1)."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        txt  = shape.text_frame.text.strip()
        name = shape.name
        if name == "Title 1":
            tf = shape.text_frame
            tf.clear()
            for i, line in enumerate(title_lines):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.alignment = PP_ALIGN.CENTER
                run = p.add_run()
                run.text = line
                run.font.size = Pt(26 if i == 0 else 17)
                run.font.bold = (i == 0)
        elif name.startswith("hlSlideMaster") and txt in ("Classification", "CUI", "UNCLASSIFIED"):
            for p in shape.text_frame.paragraphs:
                for run in p.runs:
                    run.text = CLASSIFICATION
        elif name == "TextBox 3" and txt == "Classification":
            for p in shape.text_frame.paragraphs:
                for run in p.runs:
                    run.text = CLASSIFICATION
        elif name == "TextBox 5" and txt == "POC":
            tf = shape.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = poc
            run.font.size = Pt(11)
            run.font.color.rgb = NAVY


def build_deck(deck_dir: Path, manifest: dict):
    """Build one PPTX from a directory of slide-NNN.png files."""
    png_files = sorted(
        [f for f in deck_dir.iterdir()
         if f.name.startswith("slide-") and f.suffix == ".png"],
        key=lambda p: int(p.stem.split("-")[1])
    )
    if not png_files:
        print(f"  SKIP  {deck_dir.name} — no slides found")
        return

    # Pull metadata from alignment map (enriched) + manifest (fallback)
    meta         = manifest.get(deck_dir.name, {})
    title        = meta.get("title") or deck_title(deck_dir.name)
    track        = meta.get("track", "")
    tm_level     = meta.get("maven_tm_level", "")
    also_supports = meta.get("maven_tm_also_supports", [])
    seq_order    = meta.get("sequence_order")

    # Subtitle line: "SL 2  ·  Primer" + optional also-supports
    parts = [p for p in [tm_level, track] if p]
    if also_supports:
        parts.append("also: " + ", ".join(also_supports))
    subtitle = "  ·  ".join(parts)

    out    = OUTPUT_DIR / f"{deck_dir.name}.pptx"
    n_pngs = len(png_files)

    shutil.copy(TEMPLATE, out)
    prs = Presentation(str(out))

    # ── Title slide (template slide 0) ────────────────────────────────────────
    title_lines = [title.upper(), subtitle] if subtitle else [title.upper()]
    update_title_chrome(
        prs.slides[0],
        title_lines=title_lines,
        poc="USAREUR-AF  ·  C2DAO  ·  MSS",
    )

    # ── Clone content slides BEFORE modifying slide 1 ─────────────────────────
    # We need n_pngs content slides. Slide index 1 is the unmodified template;
    # clone n_pngs - 1 more so we have slides[1 .. n_pngs].
    for _ in range(n_pngs - 1):
        clone_content_slide(prs)

    # ── Populate each content slide with its PNG ───────────────────────────────
    # Header: "TITLE  ·  TM-XX  ·  Track" — full context visible on every slide
    header_parts = [title.upper()] + [p for p in [tm_level, track] if p]
    header = "  ·  ".join(header_parts)
    for i, png_path in enumerate(png_files):
        slide = prs.slides[i + 1]          # slides[1] … slides[n_pngs]
        setup_slide(slide, header)
        slide.shapes.add_picture(str(png_path), IMG_X, IMG_Y, IMG_W, IMG_H)

    seq_tag = f"  #{seq_order}" if seq_order else ""
    prs.save(str(out))
    print(f"  OK    {out.name}  ({n_pngs + 1} slides)  [{tm_level}  {track}{seq_tag}]")


def main():
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    manifest = load_alignment()

    deck_dirs = sorted(
        d for d in SOURCE_DIR.iterdir()
        if d.is_dir() and d.name not in SKIP_DIRS
    )

    # Warn about any dirs missing from manifest
    for d in deck_dirs:
        if d.name not in manifest:
            print(f"  WARN  {d.name} — not in manifest, using dir name")

    print(f"Building {len(deck_dirs)} decks → {OUTPUT_DIR}/\n")
    for deck_dir in deck_dirs:
        build_deck(deck_dir, manifest)
    print("\nDone.")


if __name__ == "__main__":
    main()
