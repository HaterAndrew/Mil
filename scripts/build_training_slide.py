"""
Build a single-slide PPT: Maven Training Progression (TM-10 → TM-30 focus).
Output: maven_training/pdf/MSS_Training_Progression.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt
import os

# ── Army colors ──────────────────────────────────────────────────────────────
ARMY_GREEN   = RGBColor(0x4B, 0x5F, 0x27)   # OD green
ARMY_GOLD    = RGBColor(0xFF, 0xC3, 0x25)   # Army gold
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY   = RGBColor(0xF2, 0xF2, 0xF2)
DARK_GRAY    = RGBColor(0x40, 0x40, 0x40)
MED_GRAY     = RGBColor(0xA0, 0xA0, 0xA0)
DARK_GREEN   = RGBColor(0x2E, 0x3A, 0x18)   # darker for header bg

# ── Slide dimensions: widescreen 16:9 ────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

slide_layout = prs.slide_layouts[6]   # blank
slide = prs.slides.add_slide(slide_layout)
shapes = slide.shapes


def add_rect(slide, l, t, w, h, fill_rgb=None, line_rgb=None, line_width_pt=0):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
    else:
        shape.fill.background()
    if line_rgb:
        shape.line.color.rgb = line_rgb
        shape.line.width = Pt(line_width_pt)
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, l, t, w, h, text, font_size, bold=False,
                color=WHITE, align=PP_ALIGN.CENTER, wrap=True):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txb


# ── Background ────────────────────────────────────────────────────────────────
add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=LIGHT_GRAY)

# ── Header bar ────────────────────────────────────────────────────────────────
add_rect(slide, 0, 0, 13.33, 1.0, fill_rgb=DARK_GREEN)
add_textbox(slide, 0.15, 0.08, 9.0, 0.55,
            "MSS TRAINING PROGRESSION", 26, bold=True, color=ARMY_GOLD)
add_textbox(slide, 0.15, 0.58, 9.0, 0.35,
            "USAREUR-AF Operational Data Team  |  Maven Smart System (MSS)",
            10, bold=False, color=WHITE)


# ── Gold accent rule ──────────────────────────────────────────────────────────
add_rect(slide, 0, 1.0, 13.33, 0.06, fill_rgb=ARMY_GOLD)

# ─────────────────────────────────────────────────────────────────────────────
# MAIN CONTENT: three primary blocks (TM-10, TM-20, TM-30) + side note column
# Layout:  [TM-10] → [TM-20] → [TM-30]  |  [TM-40/50 side note]
# ─────────────────────────────────────────────────────────────────────────────

BLOCK_TOP   = 1.25
BLOCK_H     = 5.5
BLOCK_W     = 3.0
GAP         = 0.18
START_L     = 0.25
ARROW_Y     = BLOCK_TOP + BLOCK_H / 2 - 0.05   # vertical center of blocks

blocks = [
    {
        "label": "TM-10",
        "title": "MAVEN USER",
        "audience": "ALL PERSONNEL",
        "audience_detail": "Every Soldier, Officer & Civilian",
        "bullets": [
            "Navigate the MSS portal",
            "Access datasets & dashboards",
            "Run saved queries / filters",
            "Export and share products",
            "Understand data provenance",
        ],
        "prereq": "None",
    },
    {
        "label": "TM-20",
        "title": "BUILDER",
        "audience": "ALL STAFF",
        "audience_detail": "Light Builders (no coding required)",
        "bullets": [
            "Build Workshop dashboards",
            "Create no-code transforms",
            "Design basic data pipelines",
            "Configure object views",
            "Publish products to unit portals",
        ],
        "prereq": "Prereq: TM-10",
    },
    {
        "label": "TM-30",
        "title": "ADVANCED BUILDER",
        "audience": "DATA-ADJACENT SPECIALISTS",
        "audience_detail": "17/25-series, S6/G6, G2",
        "bullets": [
            "Author Python/SQL transforms",
            "Model ontology object types",
            "Build AIP Logic workflows",
            "Performance-tune pipelines",
            "Govern datasets & lineage",
        ],
        "prereq": "Prereq: TM-20",
    },
]

for i, blk in enumerate(blocks):
    left = START_L + i * (BLOCK_W + GAP + 0.35)   # 0.35 for arrow space

    # Card shadow (offset rect)
    add_rect(slide, left + 0.06, BLOCK_TOP + 0.06, BLOCK_W, BLOCK_H,
             fill_rgb=MED_GRAY)

    # Main card
    card = add_rect(slide, left, BLOCK_TOP, BLOCK_W, BLOCK_H,
                    fill_rgb=WHITE, line_rgb=ARMY_GREEN, line_width_pt=2)

    # Header strip
    add_rect(slide, left, BLOCK_TOP, BLOCK_W, 0.85, fill_rgb=ARMY_GREEN)

    # TM label
    add_textbox(slide, left, BLOCK_TOP + 0.04, BLOCK_W, 0.38,
                blk["label"], 22, bold=True, color=ARMY_GOLD)

    # Title
    add_textbox(slide, left, BLOCK_TOP + 0.42, BLOCK_W, 0.35,
                blk["title"], 11, bold=True, color=WHITE)

    # Audience band
    add_rect(slide, left, BLOCK_TOP + 0.85, BLOCK_W, 0.52, fill_rgb=DARK_GREEN)
    add_textbox(slide, left + 0.05, BLOCK_TOP + 0.87, BLOCK_W - 0.1, 0.26,
                blk["audience"], 9, bold=True, color=ARMY_GOLD)
    add_textbox(slide, left + 0.05, BLOCK_TOP + 1.11, BLOCK_W - 0.1, 0.22,
                blk["audience_detail"], 8, bold=False, color=WHITE)

    # Prereq badge
    add_textbox(slide, left + 0.1, BLOCK_TOP + 1.4, BLOCK_W - 0.2, 0.22,
                blk["prereq"], 8, bold=False, color=MED_GRAY, align=PP_ALIGN.LEFT)

    # Bullet points
    bullet_top = BLOCK_TOP + 1.65
    for j, bullet in enumerate(blk["bullets"]):
        add_textbox(slide, left + 0.15, bullet_top + j * 0.68,
                    BLOCK_W - 0.25, 0.62,
                    f"▸  {bullet}", 9, bold=False,
                    color=DARK_GRAY, align=PP_ALIGN.LEFT)

    # Arrow between cards (not after last)
    if i < len(blocks) - 1:
        arrow_l = left + BLOCK_W + 0.05
        add_textbox(slide, arrow_l, ARROW_Y - 0.22, 0.32, 0.44,
                    "▶", 20, bold=True, color=ARMY_GREEN)


# ─────────────────────────────────────────────────────────────────────────────
# TM-40 / TM-50 SIDE NOTE PANEL
# ─────────────────────────────────────────────────────────────────────────────
NOTE_L = START_L + 3 * (BLOCK_W + GAP + 0.35) - 0.1
NOTE_W = 13.33 - NOTE_L - 0.2
NOTE_TOP = BLOCK_TOP
NOTE_H = BLOCK_H

# Background panel (muted)
add_rect(slide, NOTE_L, NOTE_TOP, NOTE_W, NOTE_H,
         fill_rgb=RGBColor(0xE8, 0xEC, 0xDF),   # very light green-gray
         line_rgb=MED_GRAY, line_width_pt=1)

# Header
add_rect(slide, NOTE_L, NOTE_TOP, NOTE_W, 0.55,
         fill_rgb=MED_GRAY)
add_textbox(slide, NOTE_L, NOTE_TOP + 0.06, NOTE_W, 0.42,
            "SPECIALIST TRACKS", 9, bold=True, color=WHITE)

# TM-40 block
add_textbox(slide, NOTE_L + 0.1, NOTE_TOP + 0.65, NOTE_W - 0.15, 0.28,
            "TM-40 SERIES  (Prereq: TM-30)", 8, bold=True,
            color=DARK_GREEN, align=PP_ALIGN.LEFT)

tracks_40 = [
    "WFF Tracks (A–F):",
    "  A Intel · B Fires · C M&M",
    "  D Sust · E Prot · F MC",
    "Specialist Tracks (G–M):",
    "  G ORSA · H AI Eng · I MLE",
    "  J PM · K KM · L SWE",
]
for k, t in enumerate(tracks_40):
    is_header = t.endswith(":")
    add_textbox(slide, NOTE_L + 0.10, NOTE_TOP + 0.95 + k * 0.36,
                NOTE_W - 0.15, 0.32,
                t, 7.5, bold=is_header,
                color=DARK_GREEN if is_header else DARK_GRAY,
                align=PP_ALIGN.LEFT)

# Divider
add_rect(slide, NOTE_L + 0.1, NOTE_TOP + 3.20, NOTE_W - 0.2, 0.04,
         fill_rgb=MED_GRAY)

# TM-50 block
add_textbox(slide, NOTE_L + 0.1, NOTE_TOP + 3.32, NOTE_W - 0.15, 0.28,
            "TM-50 SERIES  (Prereq: TM-40)", 8, bold=True,
            color=DARK_GREEN, align=PP_ALIGN.LEFT)
add_textbox(slide, NOTE_L + 0.1, NOTE_TOP + 3.62, NOTE_W - 0.15, 1.6,
            "Advanced / expert-level continuation\nfor each TM-40G–M track.\n\nFocus: production pipelines,\nMLOps, AIP Agents, advanced\nOntology patterns.",
            7.5, bold=False, color=DARK_GRAY, align=PP_ALIGN.LEFT)

# ── Footer ────────────────────────────────────────────────────────────────────
add_rect(slide, 0, 7.18, 13.33, 0.32, fill_rgb=DARK_GREEN)
add_textbox(slide, 0.2, 7.19, 9.0, 0.28,
            "USAREUR-AF OD Team  ·  MSS Training Program  ·  Wiesbaden, Germany",
            7, bold=False, color=MED_GRAY)


# ── Save ──────────────────────────────────────────────────────────────────────
out_dir = os.path.join(os.path.dirname(__file__), "..", "maven_training", "pdf")
os.makedirs(out_dir, exist_ok=True)
out_path = os.path.join(out_dir, "MSS_Training_Progression.pptx")
prs.save(out_path)
print(f"Saved: {os.path.afbcath(out_path)}")
