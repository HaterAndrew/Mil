#!/usr/bin/env python3
"""
Build new PPTX decks for curriculum gaps identified in content audit:
  - WFF Overview          (SL 4A–F)  → WFF_Overview.pptx
  - ML Engineering        (SL 4M)    → TM_40M_ML_Engineering.pptx
  - Program Manager       (SL 4J)    → TM_40J_Program_Manager.pptx
  - Advanced MSS Overview (SL 5)     → TM_50_Advanced_Overview.pptx
  - Advanced ORSA         (SL 5G)    → TM_50G_Advanced_ORSA.pptx
  - Advanced AI Eng       (SL 5H)    → TM_50H_Advanced_AI_Engineer.pptx
  - Advanced ML Eng       (SL 5M)    → TM_50M_Advanced_ML_Engineer.pptx
  - Advanced PM           (SL 5J)    → TM_50J_Advanced_Program_Manager.pptx
  - Advanced KM           (SL 5K)    → TM_50K_Advanced_Knowledge_Manager.pptx
  - Advanced SWE          (SL 5L)    → TM_50L_Advanced_Software_Engineer.pptx

Run from repo root: python3 scripts/build_new_decks.py
Output: maven_training/pdf/
"""

import copy
import os
import shutil
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Paths ──────────────────────────────────────────────────────────────────────
TEMPLATE   = Path(os.environ.get(
    "PPTX_TEMPLATE",
    Path(__file__).parent.parent / "maven_training" / "source_material" / "USAREUR-AF PPT Template.pptx"
))
if not TEMPLATE.exists():
    raise FileNotFoundError(f"Template not found: {TEMPLATE}  — set PPTX_TEMPLATE env var")
OUTPUT_DIR = Path("maven_training/pdf")

# ── Palette ────────────────────────────────────────────────────────────────────
NAVY      = RGBColor(0x17, 0x32, 0x5C)
NAVY2     = RGBColor(0x18, 0x33, 0x5F)
NAVY_MID  = RGBColor(0x26, 0x4D, 0x7E)
PURPLE    = RGBColor(0x70, 0x30, 0xA0)
GOLD      = RGBColor(0xC8, 0x9A, 0x00)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_LT   = RGBColor(0xE2, 0xE8, 0xF0)

CLASSIFICATION = ""

# ── Content geometry ──────────────────────────────────────────────────────────
X0  = 0.15
Y0  = 0.70
YB  = 7.20
XE  = 13.18
CW  = XE - X0


# ── Chrome helpers (same pattern as build_proto_army_data_orientation) ────────

def remove_shape(shape):
    el = shape._element
    el.getparent().remove(el)


def clone_content_slide(prs):
    src = prs.slides[1]
    slide = prs.slides.add_slide(src.slide_layout)
    sp_tree = slide.shapes._spTree
    for child in list(sp_tree)[2:]:
        sp_tree.remove(child)
    for child in list(src.shapes._spTree)[2:]:
        sp_tree.append(copy.deepcopy(child))
    return slide


def setup_slide(slide, header_text):
    to_remove = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        txt  = shape.text_frame.text.strip()
        name = shape.name
        if name == "TextBox 2" or txt == "Header":
            tf = shape.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = header_text
            run.font.size = Pt(14)
            run.font.bold = True
            run.font.color.rgb = NAVY
        elif txt in ("CUI", "UNCLASSIFIED", "Classification"):
            for p in shape.text_frame.paragraphs:
                for run in p.runs:
                    run.text = CLASSIFICATION
        elif name == "TextBox 8":
            to_remove.append(shape)
    for s in to_remove:
        remove_shape(s)
    bot = slide.shapes.add_textbox(Inches(0), Inches(7.25), Inches(13.33), Inches(0.25))
    tf = bot.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = CLASSIFICATION
    run.font.size = Pt(9)
    run.font.bold = True
    run.font.color.rgb = PURPLE


def update_title_chrome(slide, title_lines, poc):
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        txt  = shape.text_frame.text.strip()
        name = shape.name
        if name == "Title 1":
            tf = shape.text_frame
            tf.clear()
            for i, line in enumerate(title_lines):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.alignment = PP_ALIGN.CENTER
                run = p.add_run()
                run.text = line
                run.font.size = Pt(26 if i == 0 else 17)
                run.font.bold = (i == 0)
        elif name.startswith("hlSlideMaster") and txt in ("Classification", "CUI", "UNCLASSIFIED"):
            for p in shape.text_frame.paragraphs:
                for run in p.runs:
                    run.text = CLASSIFICATION
        elif name == "TextBox 3" and txt == "Classification":
            for p in shape.text_frame.paragraphs:
                for run in p.runs:
                    run.text = CLASSIFICATION
        elif name == "TextBox 5" and txt == "POC":
            tf = shape.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = poc
            run.font.size = Pt(11)
            run.font.color.rgb = NAVY


# ── Slide content helpers ─────────────────────────────────────────────────────

def add_section_title(slide, text, y=0.80, size=20):
    """Bold navy section title below the chrome header."""
    tb = slide.shapes.add_textbox(Inches(X0), Inches(y), Inches(CW), Inches(0.50))
    tf = tb.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = NAVY


def add_rule(slide, y=1.32):
    """Thin gold rule under section title."""
    from pptx.oxml.ns import qn
    from pptx.oxml import parse_xml
    tb = slide.shapes.add_textbox(Inches(X0), Inches(y), Inches(CW), Inches(0.02))
    sp = tb._element
    spPr = sp.spPr
    ns = 'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"'
    fill_xml = f'<a:solidFill {ns}><a:srgbClr val="{str(GOLD)}"/></a:solidFill>'
    spPr.append(parse_xml(fill_xml))


def add_bullets(slide, items, y_start=1.40, font_size=12.5, col_split=None):
    """
    Render a list of bullet items. Items are strings.
    Leading '▶ ' marks a sub-section label (bold navy).
    Leading '  ' marks indented sub-bullet.
    col_split: if int, split items into two columns at that index.
    """
    def render_column(slide, items, x, y, w, h):
        tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = True
        first = True
        for item in items:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            if item.startswith("▶ "):
                p.space_before = Pt(8)
                run = p.add_run()
                run.text = item[2:]
                run.font.size = Pt(font_size)
                run.font.bold = True
                run.font.color.rgb = NAVY
            elif item.startswith("  "):
                p.space_before = Pt(1)
                run = p.add_run()
                run.text = "    –  " + item.strip()
                run.font.size = Pt(font_size - 1)
                run.font.color.rgb = NAVY_MID
            else:
                p.space_before = Pt(3)
                run = p.add_run()
                run.text = "▪  " + item
                run.font.size = Pt(font_size)
                run.font.color.rgb = NAVY2

    avail_h = YB - y_start - 0.15
    if col_split:
        left  = items[:col_split]
        right = items[col_split:]
        col_w = CW / 2 - 0.10
        render_column(slide, left,  X0,             y_start, col_w, avail_h)
        render_column(slide, right, X0 + col_w + 0.20, y_start, col_w, avail_h)
    else:
        render_column(slide, items, X0, y_start, CW, avail_h)


def add_two_column_table(slide, left_header, right_header, rows, y_start=1.40):
    """Render a simple two-column comparison."""
    col_w  = CW / 2 - 0.10
    h_size = 13

    def header_box(x, text):
        tb = slide.shapes.add_textbox(Inches(x), Inches(y_start), Inches(col_w), Inches(0.30))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        run.font.size = Pt(h_size)
        run.font.bold = True
        run.font.color.rgb = GOLD

    header_box(X0, left_header)
    header_box(X0 + col_w + 0.20, right_header)
    add_bullets(slide, [r[0] for r in rows], y_start=y_start + 0.35,
                font_size=11.5, col_split=None)

    tb = slide.shapes.add_textbox(
        Inches(X0 + col_w + 0.20), Inches(y_start + 0.35),
        Inches(col_w), Inches(YB - y_start - 0.50))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for row in rows:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(3)
        run = p.add_run()
        run.text = "▪  " + row[1]
        run.font.size = Pt(11.5)
        run.font.color.rgb = NAVY2


# ── Deck builder ──────────────────────────────────────────────────────────────

def build_deck(filename, title_lines, poc, header_text, slides_content):
    """
    slides_content: list of (section_title, bullet_items) tuples.
    Each tuple becomes one content slide.
    """
    out = OUTPUT_DIR / filename
    shutil.copy(TEMPLATE, out)
    prs = Presentation(str(out))

    update_title_chrome(prs.slides[0], title_lines, poc)

    n = len(slides_content)
    for _ in range(n - 1):
        clone_content_slide(prs)

    for i, (sec_title, bullets) in enumerate(slides_content):
        slide = prs.slides[i + 1]
        setup_slide(slide, header_text)
        add_section_title(slide, sec_title)
        add_rule(slide)
        if bullets:
            add_bullets(slide, bullets)

    prs.save(str(out))
    print(f"  OK    {out.name}  ({n + 1} slides)")


# ══════════════════════════════════════════════════════════════════════════════
# DECK DEFINITIONS
# ══════════════════════════════════════════════════════════════════════════════

def deck_wff_overview():
    build_deck(
        filename="WFF_Overview.pptx",
        title_lines=["HOW MSS SUPPORTS WARFIGHTING FUNCTIONS",
                     "SL 4A · SL 4B · SL 4C · SL 4D · SL 4E · SL 4F"],
        poc="USAREUR-AF  ·  C2DAO  ·  MSS",
        header_text="HOW MSS SUPPORTS WARFIGHTING FUNCTIONS  ·  SL 4 WFF",
        slides_content=[
            ("THE SIX WARFIGHTING FUNCTIONS AND MSS", [
                "MSS is not a single-function tool — it is the common data layer for all six WFFs",
                "Each WFF has its own SL 4 track (A–F) with doctrine-aligned workflows",
                "▶ Warfighting Function Tracks",
                "  SL 4A — Intelligence (IPOE, all-source analysis, targeting)",
                "  SL 4B — Fires (D3A, FSC, AMD, C-RAM, joint fires)",
                "  SL 4C — Movement & Maneuver (MDMP, routes, force tracking)",
                "  SL 4D — Sustainment (supply chain, maintenance, transport)",
                "  SL 4E — Protection (CRM, CBRN, AT/FP, survivability)",
                "  SL 4F — Mission Command (battle rhythm, COP, CCIR, staff integration)",
                "▶ Prerequisite: SL 3 (Advanced Builder) required for all WFF tracks",
            ]),
            ("INTELLIGENCE (SL 4A) — IPOE AND ANALYSIS IN MSS", [
                "▶ Core Functions",
                "  Intelligence Preparation of the Battlefield (IPOE) — all four steps in MSS",
                "  All-source fusion: GEOINT, SIGINT, HUMINT, OSINT, MASINT, TECHINT",
                "  Pattern of Life, SIGACT analysis, link analysis and network mapping",
                "▶ Collection & Targeting",
                "  PIR/IR management, RFI workflow, collection synchronization matrix",
                "  Intelligence support to targeting: HVT/HPT lists, D3A, BDA tracking",
                "▶ Products & Dissemination",
                "  INTSUM standards, MSS dashboard as intelligence product",
                "  Classification handling within MSS workspaces",
                "▶ Legacy Tool Migration: DCGS-A → MSS",
            ]),
            ("FIRES (SL 4B) — TARGETING AND FIRE SUPPORT IN MSS", [
                "▶ Core Functions",
                "  D3A process: Decide (HPTL/AGM), Detect (sensor feeds), Deliver, Assess (BDA)",
                "  Fire support plan data management and FSCM visualization",
                "  13F FSO workflow — fire support execution matrix",
                "▶ Air Defense and C-RAM",
                "  Air picture visualization, AMD system employment status",
                "  C-RAM tracker, Q-36/Q-37 radar data integration",
                "  Critical Asset List / Defended Asset List management",
                "▶ Joint Fires",
                "  JTAC/JFO workflows, Air Tasking Order data, airspace deconfliction",
                "▶ Legacy Tool Migration: AFATDS → MSS",
            ]),
            ("MOVEMENT & MANEUVER (SL 4C) — OPERATIONS IN MSS", [
                "▶ Operational Planning",
                "  Full MDMP support: COA development, OPORD visualization, task organization",
                "  FRAGO management and version control",
                "▶ Route Planning and Mobility",
                "  MSR/ASR management, obstacle reports, bridge classification",
                "  Route clearance tracking, breach planning data",
                "▶ Force Tracking",
                "  Blue force tracking data integration, unit location reporting standards",
                "  Equipment tracking and status, combat power reporting",
                "▶ Aviation Integration",
                "  Air corridor management, UAS deconfliction, air assault data",
                "▶ Legacy Tool Migration: CPCE / FBCB2 → MSS",
            ]),
            ("SUSTAINMENT (SL 4D) — LOGISTICS IN MSS", [
                "▶ Supply Chain",
                "  All supply classes tracked in MSS (I–X); push vs. pull visualization",
                "  92A workflow, distribution management dashboard, threshold alerts",
                "▶ Maintenance",
                "  Equipment readiness and C-rating tracking (AR 700-138)",
                "  PMCS status, work order management, GCSS-Army integration",
                "▶ Transportation and Ammunition",
                "  Convoy planning, movement requests, vehicle readiness",
                "  Ammunition basic load, ASR/CSR tracking, DODIC management",
                "▶ Personnel and HR",
                "  PERSTAT management, casualty reporting, theater strength management",
                "▶ Legacy Tool Migration: GCSS-Army / ULLS → MSS",
            ]),
            ("PROTECTION (SL 4E) — CRM AND FORCE PROTECTION IN MSS", [
                "▶ Composite Risk Management",
                "  Hazard identification, risk scoring, control measures tracking",
                "  Residual risk acceptance and approving authority documentation",
                "▶ CBRN Defense",
                "  CBRN threat data, NBC 1–6 reports, contamination control points",
                "  Detection equipment status, decontamination site management",
                "▶ Antiterrorism and Force Protection",
                "  FPCON tracking, RAM management, threat/vulnerability assessments",
                "▶ Physical Security and Survivability",
                "  Entry control point data, base perimeter visualization",
                "  Fighting position tracking, hardening status for key facilities",
                "▶ Legacy Tool Migration: APERP / manual tracking → MSS",
            ]),
            ("MISSION COMMAND (SL 4F) — STAFF INTEGRATION IN MSS", [
                "▶ Staff Section Integration",
                "  S1–S6/G1–G6 role-specific MSS workflows defined in SL 4F",
                "  XO as MSS governance authority; battle captain as shift operator",
                "▶ Operations Process Support",
                "  Full MDMP planning, pre-execution checks, battle tracking",
                "  Continuous assessment: MOE/MOP framework on MSS",
                "▶ Battle Rhythm Management",
                "  BUA, CUB, targeting meetings, logistics sync, readiness reviews",
                "  Commander's Update Brief — decision product standards",
                "▶ Common Operating Picture",
                "  Unit status, logistics, intelligence, and threat layers",
                "  COP currency standards (data freshness) and subordinate access",
                "▶ CCIR Framework: PIR, FFIR, EEFI — loaded and monitored in MSS",
            ]),
            ("CROSS-WFF DATA FLOWS AND DECISION ADVANTAGE", [
                "MSS creates a unified data environment — no WFF operates in isolation",
                "▶ Intelligence → Fires Integration",
                "  S2 provides targeting data; S3/FSO coordinates delivery and BDA",
                "  Intelligence-fires integration checklist governs workspace access",
                "▶ Operations → Sustainment Integration",
                "  S3 OPORD drives logistics synchronization (S4 distribution sync)",
                "  Readiness data (S4/S1) feeds combat power reporting (S3)",
                "▶ Protection → All WFFs",
                "  CRM and AT/FP feeds into all operational planning products",
                "  CBRN data informs M&M route planning and staging area selection",
                "▶ Mission Command — the integrator",
                "  SL 4F ties all five WFFs into a coherent COP and battle rhythm",
                "  Decision advantage emerges from timely, accurate cross-WFF data",
            ]),
            ("LEGACY TOOL MIGRATION — WHY MSS REPLACES LEGACY TOOLS", [
                "▶ What legacy tools cannot do",
                "  Single-purpose: AFATDS (fires only), GCSS-A (logistics only), DCGS-A (intel only)",
                "  No cross-domain data sharing without manual hand-jamming",
                "  No common operating picture — each staff section has its own picture",
                "▶ What MSS adds",
                "  Single platform: all six WFFs share one data environment",
                "  Real-time cross-WFF visibility: sustainment status feeds maneuver planning",
                "  Ontology layer: relationships between data are explicit, not implicit",
                "  AIP layer: AI-assisted analysis available to all WFFs simultaneously",
                "▶ Migration principle: MSS does not delete legacy workflows",
                "  It provides a data integration layer that cross-references legacy outputs",
                "  GCSS-Army reconciliation, AFATDS handoff, DCGS-A data import all documented",
            ]),
            ("WFF TRACK ENROLLMENT AND PREREQUISITES", [
                "▶ Hard Prerequisite for ALL WFF tracks",
                "  SL 1 → SL 2 → SL 3 → SL 4 (WFF track)",
                "  No direct enrollment into SL 4A–F without SL 3 completion",
                "▶ WFF track selection: based on duty position and MOS",
                "  S2/G2 / 35-series → SL 4A (Intelligence)",
                "  FA / 13-series / 14-series → SL 4B (Fires)",
                "  S3/G3 / 11-series / 19-series / 12-series / 15-series → SL 4C (M&M)",
                "  S4/G4 / 90-series / 91-series / 92-series / 88-series → SL 4D (Sustainment)",
                "  Protection Officer / 31-series / CBRN / AT → SL 4E (Protection)",
                "  CG / DCG / CoS / XO / Battle Captain / staff officers → SL 4F (Mission Command)",
                "▶ Multiple track enrollment: authorized for multi-functional officers",
                "▶ Specialist tracks (SL 4G–M) are separate from WFF tracks",
            ]),
        ]
    )


def deck_tm40m():
    build_deck(
        filename="TM_40M_ML_Engineering.pptx",
        title_lines=["ML ENGINEERING ON MSS",
                     "SL 4M  ·  Specialist Track  ·  Prereq: SL 3"],
        poc="USAREUR-AF  ·  C2DAO  ·  MSS",
        header_text="ML ENGINEERING ON MSS  ·  SL 4M  ·  Specialist",
        slides_content=[
            ("THE MLE ROLE IN USAREUR-AF", [
                "ML Engineers build, evaluate, deploy, and sustain predictive models on MSS",
                "▶ What MLE Does — Not What It Does Not",
                "  MLE owns the model lifecycle from feature engineering through production",
                "  MLE does NOT make operational decisions — models support human judgment",
                "  MLE does NOT author AIP Logic workflows — that is SL 4H (AI Engineer)",
                "▶ Where MLE Sits in the Data Chain",
                "  Downstream of data engineers (Pipeline Builder, SL 2/30)",
                "  Upstream of analysts and operators who consume model outputs",
                "  Parallel to AI Engineers (AIP) and ORSA (statistical analysis)",
                "▶ Model Governance Requirement",
                "  Every model in production requires a signed Model Card (six-gate process)",
                "  No model may reach production without commander-level risk acceptance",
            ]),
            ("CODE WORKSPACES — THE MLE ENVIRONMENT", [
                "▶ Code Workspace Overview",
                "  Jupyter-compatible environment inside Foundry — no local setup required",
                "  Connects directly to Foundry datasets, ontology objects, and model registry",
                "▶ Key Setup Tasks (SL 4M Chapter 2)",
                "  Create and configure workspace; select correct compute tier",
                "  Connect to Foundry datasets in a notebook (no file downloads)",
                "  Version control via Foundry's branch mechanism — not git-external",
                "▶ Compute Management",
                "  Compute costs real capacity — select minimum tier for task",
                "  Terminate workloads that complete; do not leave large kernels running",
                "▶ Package Management",
                "  Pre-approved package list governs available libraries",
                "  scikit-learn, PyTorch, XGBoost, Optuna, MLflow approved for MSS use",
            ]),
            ("FEATURE ENGINEERING AND FEATURE STORE", [
                "▶ Feature Engineering Transform (Chapter 3)",
                "  Built as a Foundry Transform — versioned, auditable, reproducible",
                "  Input: raw Foundry dataset; Output: feature table registered in Feature Store",
                "▶ Common MSS Feature Patterns",
                "  Unit readiness: C-rating history, equipment status, personnel fill rate",
                "  Temporal features: rolling averages, lag features, trend indicators",
                "  Categorical encoding: MOS, echelon, equipment type",
                "▶ Feature Store",
                "  Central registry — prevents re-computation of shared features",
                "  Feature versioning enables reproducibility across model iterations",
                "  Feature lineage tracked: every feature traces back to a source dataset",
                "▶ Temporal Split Requirement",
                "  ALL training splits must respect time — never random shuffle operational data",
                "  Leakage from future data invalidates model validity for operational use",
            ]),
            ("MODEL TRAINING AND EVALUATION", [
                "▶ Training on MSS (Chapter 4)",
                "  scikit-learn models: standard fit/predict with MLflow experiment tracking",
                "  Hyperparameter tuning: Optuna (approved); grid search acceptable for small models",
                "  PyTorch: available for deep learning use cases with GPU compute tier",
                "▶ Evaluation Standards (Chapter 5)",
                "  Every model requires a documented Evaluation Report before deployment",
                "  Metrics selected based on use case: F1 for classification, RMSE for regression",
                "  Bias and Fairness Assessment required — especially for personnel-affecting models",
                "▶ Evaluation Metrics Reference by Use Case",
                "  C-rating prediction: classification (Precision/Recall/F1)",
                "  Logistics demand: regression (MAE, RMSE, MAPE)",
                "  Anomaly detection: precision at k, area under PR curve",
                "▶ Peer review required before any evaluation report is finalized",
            ]),
            ("MODEL DEPLOYMENT AND MLOPS", [
                "▶ Deployment Patterns (Chapter 6)",
                "  Batch inference: deployed as a Foundry Transform — runs on schedule",
                "  Online inference: available for low-latency use cases (requires architecture review)",
                "  Predictions published as Ontology Object Properties — visible to analysts/operators",
                "▶ MLOps on Foundry (Chapter 7)",
                "  Model versioning: every deployment tagged with version, training date, metrics",
                "  Drift detection: Production Software Index (PSI) monitored continuously",
                "  Retraining triggers: PSI threshold breach OR schedule (whichever comes first)",
                "▶ Inference Caching",
                "  Freshness indicators attached to all predictions",
                "  Stale predictions flagged for operators — no silent staleness allowed",
                "▶ Model Versioning Requirement",
                "  All prior model versions retained for 180 days for audit purposes",
            ]),
            ("APPROVED OPERATIONAL USE CASES", [
                "▶ Three Primary Use Cases (Chapter 8)",
                "  Unit Readiness C-Rating Prediction",
                "    Input: equipment status, personnel fill, maintenance backlog, historical C-ratings",
                "    Output: predicted C-rating (30-day horizon), confidence interval",
                "    Consumer: S3/G3, readiness officers, sustainment S4",
                "  Logistics Demand Forecasting (Class IX)",
                "    Input: usage rates, operational tempo, historical consumption, deployment data",
                "    Output: projected demand by supply class, reorder triggers",
                "    Consumer: S4/G4, BSB/FSB supply officers",
                "  OPDATA Anomaly Detection",
                "    Input: time-series operational data streams",
                "    Output: anomaly flags and severity scores",
                "    Consumer: intelligence analysts, operations center watch officers",
                "▶ All use cases require authorization before model development begins (Gate 1)",
            ]),
            ("MODEL GOVERNANCE — THE SIX-GATE PROCESS", [
                "Every production model in MSS must pass six governance gates:",
                "▶ Gate 1 — Use Case Authorization",
                "  Commander-level approval for the use case before any development",
                "▶ Gate 2 — Design Review",
                "  Architecture, data sources, and risk profile reviewed by peers",
                "▶ Gate 3 — Evaluation Acceptance",
                "  Metrics meet threshold; bias/fairness assessment completed",
                "▶ Gate 4 — Pre-Production Review",
                "  Deployment architecture, monitoring plan, and rollback procedure approved",
                "▶ Gate 5 — Production Deployment Approval",
                "  Final commander sign-off; Model Card completed and signed",
                "▶ Gate 6 — Post-Deployment Monitoring Confirmation",
                "  Drift monitoring active; retraining triggers configured; on-call defined",
                "▶ Model Card: 10-section document (purpose, data, evaluation, deployment, monitoring)",
                "  Model Card is the authoritative artifact — survives personnel turnover",
            ]),
            ("MLE ROLE BOUNDARIES AND CROSS-TRACK COORDINATION", [
                "▶ MLE vs. AI Engineer (SL 4H)",
                "  MLE: statistical/ML models, feature engineering, model lifecycle",
                "  AI Engineer: LLM workflows, AIP Logic, Agent Studio, RAG pipelines",
                "  Overlap: code workspaces, Foundry data access, MLflow; coordinate on shared infra",
                "▶ MLE vs. ORSA (SL 4G)",
                "  ORSA: statistical modeling, simulation, decision analysis, OR products",
                "  MLE: predictive models in production, automated outputs at scale",
                "  ORSA validates MLE model assumptions; MLE operationalizes ORSA research",
                "▶ MLE vs. Software Engineer (SL 4L)",
                "  SWE builds the platform and APIs; MLE consumes them",
                "  SWE designs Feature Store schema; MLE populates it",
                "  Coordinate on inference endpoint architecture and data contracts",
                "▶ MLE Governing References",
                "  DoD RAIMTF, Army AI Ethics Policy, MSS ML Governance SOP, SL 4M",
            ]),
            ("SL 4M CURRICULUM OVERVIEW AND NEXT STEPS", [
                "▶ SL 4M Chapter Structure",
                "  Ch 1: MLE Role and Scope",
                "  Ch 2: Code Workspaces",
                "  Ch 3: Feature Engineering and Feature Store",
                "  Ch 4: Model Training (scikit-learn, Optuna, PyTorch)",
                "  Ch 5: Model Evaluation and Bias Assessment",
                "  Ch 6: Model Deployment (batch, online, ontology integration)",
                "  Ch 7: MLOps (versioning, drift detection, retraining)",
                "  Ch 8: Operational Use Cases (readiness, logistics, anomaly)",
                "  Ch 9: Model Governance (six-gate process, Model Card)",
                "▶ Appendices",
                "  Model Governance Checklist (all six gates)",
                "  Approved Model Use Cases (USAREUR-AF)",
                "  ML Quick Reference (imports, split patterns, drift detection, MLflow)",
                "▶ Next Steps: Complete SL 4M → eligible for SL 5M (Advanced ML Engineer)",
            ]),
        ]
    )


def deck_tm40j():
    build_deck(
        filename="TM_40J_Program_Manager.pptx",
        title_lines=["DATA PROGRAM MANAGER ESSENTIALS",
                     "SL 4J  ·  Specialist Track  ·  Prereq: SL 3"],
        poc="USAREUR-AF  ·  C2DAO  ·  MSS",
        header_text="DATA PROGRAM MANAGER ESSENTIALS  ·  SL 4J  ·  Specialist",
        slides_content=[
            ("THE TECHNICAL PM ROLE IN THE MSS ECOSYSTEM", [
                "The SL 4J PM translates between technical teams and operational stakeholders",
                "▶ What Technical PMs Do",
                "  Scope and sequence data/AI projects from requirement to production",
                "  Manage backlogs, sprints, dependencies, risks, and delivery timelines",
                "  Build visibility: commander-facing dashboards and project tracking on MSS",
                "▶ What Technical PMs Do NOT Do",
                "  Do not make technical architecture decisions (that is SL 4L/H/M)",
                "  Do not interpret analytical outputs (that is ORSA/analyst function)",
                "  Do not approve models for production (commander retains risk authority)",
                "▶ PM Authority vs. Technical Authority",
                "  PM owns schedule, scope, and stakeholder communication",
                "  Technical lead owns design decisions within approved constraints",
                "  Conflict resolution: PM escalates scope disputes; tech lead escalates design disputes",
                "▶ Curriculum Position: SL 4J after SL 3; advanced track is SL 5J",
            ]),
            ("AGILE FOR DATA AND AI PROJECTS (CHAPTER 2)", [
                "▶ Why Agile — Not Waterfall — for Data Projects",
                "  Data requirements are discovered, not specified upfront",
                "  Model performance is unknown until training; iteration is essential",
                "  Operational feedback loops require frequent delivery, not big-bang releases",
                "▶ Scrum Framework for MSS Projects",
                "  Sprint: 2-week default; daily standup; sprint review with stakeholder",
                "  Backlog groomed weekly; stories estimated in story points",
                "  Definition of Done enforced before any story is accepted",
                "▶ Kanban for Operational Support",
                "  Board structure: Backlog → In Progress → Review → Done",
                "  WIP limits prevent overloading the team during high-tempo periods",
                "▶ Sprint Ceremonies: Planning, Standup, Review, Retrospective",
                "  Review: demo working software/data product — no slide decks as proof of work",
                "  Retrospective: identify one improvement per sprint, not a complaint session",
            ]),
            ("ML/AI PROJECT LIFECYCLE (CHAPTER 3)", [
                "The ML/AI project lifecycle has six phases — PM has specific actions in each:",
                "▶ Phase 1 — Problem Definition",
                "  PM translates commander's question into a machine-learnable problem statement",
                "  Success criteria defined in measurable terms before development begins",
                "▶ Phase 2 — Data Audit",
                "  PM coordinates access; data engineers assess quality and availability",
                "  PM documents risk if data is insufficient — does not proceed without acknowledgment",
                "▶ Phase 3 — Prototype",
                "  PM scopes the prototype (not full production); time-boxes to 2–3 sprints",
                "▶ Phase 4 — Evaluation",
                "  PM chairs evaluation review; ensures non-technical stakeholders understand output",
                "▶ Phase 5 — Production",
                "  PM coordinates Model Card completion, governance gates, deployment approval",
                "▶ Phase 6 — Sustainment",
                "  PM owns monitoring cadence; escalates model performance degradation",
            ]),
            ("STAKEHOLDER MANAGEMENT AND REQUIREMENTS (CHAPTER 4)", [
                "▶ The Translation Problem",
                "  Commanders ask questions; ML Engineers build models; PM bridges the gap",
                "  'We need better readiness forecasting' ≠ a model specification",
                "  PM extracts: What decision will this inform? What timeframe? What accuracy?",
                "▶ Requirements Elicitation",
                "  Structured interviews: ask for examples of decisions, not feature lists",
                "  Document as user stories: 'As a [role], I need [output] so that [decision]'",
                "▶ Managing Expectations",
                "  Communicate velocity honestly — resist pressure to commit to unrealistic dates",
                "  'Ready in 2 sprints' is a forecast, not a guarantee; qualify uncertainty",
                "▶ Cross-Track Coordination",
                "  PM maintains dependency log across SL 4H (AI), SL 4M (ML), SL 4L (SWE)",
                "  Blocked dependencies escalated to PM — not resolved ad hoc by engineers",
                "▶ Protecting the Team: PM absorbs stakeholder pressure; engineers execute",
            ]),
            ("PROJECT TRACKING SYSTEMS ON MSS (CHAPTER 5)", [
                "▶ Why Track on MSS — Not Spreadsheets",
                "  Single source of truth: same platform as delivery artifacts",
                "  Stakeholder visibility without meetings: dashboard replaces status briefs",
                "  Audit trail: all changes timestamped and attributed",
                "▶ Sprint Board in Workshop",
                "  Object Types: Epic, Story, Task, Blocker",
                "  Link Types: Story-to-Epic, Task-to-Story, Blocker-to-Story",
                "  Views: Sprint Board (Kanban), Backlog (table), Burndown (chart)",
                "▶ Commander-Facing Dashboard",
                "  One page: status by project, milestones hit/missed, blockers, risk flag",
                "  No jargon: 'Green/Amber/Red' with plain-language description",
                "  Design principle: GO/SES reads in 90 seconds — not a data dump",
                "▶ Automated Status Alerts",
                "  AIP-driven alerts when story count drops or blockers age past threshold",
            ]),
            ("RISK AND DEPENDENCY MANAGEMENT (CHAPTER 6)", [
                "▶ Risk Types Specific to Data Products",
                "  Data quality risk: source data is incomplete, stale, or inconsistent",
                "  Model risk: model performs in training but fails on operational data",
                "  Infrastructure risk: compute, storage, or access not available when needed",
                "  Governance risk: model card not completed; commander approval not obtained",
                "  Personnel risk: key technical staff rotate before knowledge transfer",
                "▶ Risk Register in MSS",
                "  Risk object: description, probability, impact, mitigation, owner, due date",
                "  Sprint risk review: PM reviews risk register every sprint",
                "  Severity matrix: probability × impact → Red/Amber/Green",
                "▶ Dependency Management",
                "  Common patterns: data pipeline must complete before feature store update",
                "  Model training depends on feature engineering completion",
                "  Deployment depends on governance gate approval",
                "▶ Technical Debt: track as a risk; allocate 20% of sprint capacity to debt reduction",
            ]),
            ("DELIVERY PLANNING AND PRODUCTION READINESS (CHAPTER 7)", [
                "▶ Scope / Timeline / Quality Tradeoffs",
                "  PM manages the tradeoff triangle — cannot optimize all three simultaneously",
                "  Default position: do not cut quality; adjust scope or timeline",
                "▶ Definition of Done — Data Product",
                "  Code reviewed and merged; tests passing; documentation current",
                "  Monitoring configured; on-call procedure documented; stakeholder accepted",
                "▶ What 'Production Ready' Means for a Data Product",
                "  Model Card signed; all six governance gates passed",
                "  Rollback procedure tested; drift detection active",
                "  User guide published; training delivered to operational consumers",
                "▶ Milestone Types",
                "  Internal: sprint demo, evaluation complete, staging deployment",
                "  External: commander review, governance gate, production deployment",
                "▶ Post-Release Review: PM chairs at T+14 — measures adoption, not just delivery",
            ]),
            ("CHANGE MANAGEMENT AND USER ADOPTION (CHAPTER 8)", [
                "▶ Why Operational Users Resist New Data Products",
                "  Not laziness — trained to distrust unfamiliar data sources",
                "  Legacy tool habits are deeply embedded (AFATDS, GCSS-Army, Excel)",
                "  Fear: wrong model output leads to bad operational decision",
                "▶ Change Management Planning",
                "  Stakeholder map: identify champions, skeptics, and blockers by name",
                "  Champion first: get one influential user to adopt and endorse",
                "  Train-the-trainer: embed capability with unit S6 or data champion",
                "▶ Managing Resistance at Rollout",
                "  Low adoption at T+14: diagnose cause before adding features",
                "  Common cause: product answers the wrong question — go back to requirements",
                "▶ Platform Governance — PM Perspective",
                "  PM participates in C2DAO governance reviews for own products",
                "  Data product retirement: PM owns decommission plan and user notification",
            ]),
            ("SL 4J CURRICULUM OVERVIEW AND NEXT STEPS", [
                "▶ SL 4J Chapter Structure",
                "  Ch 1: Technical PM Role and Scope",
                "  Ch 2: Agile Project Management (Scrum/Kanban for data projects)",
                "  Ch 3: ML/AI Project Lifecycle (six phases)",
                "  Ch 4: Stakeholder Management and Requirements Translation",
                "  Ch 5: Project Tracking Systems on MSS",
                "  Ch 6: Risk and Dependency Management",
                "  Ch 7: Delivery Planning and Production Readiness",
                "  Ch 8: Change Management and User Adoption",
                "▶ Appendices",
                "  Project Kickoff Checklist",
                "  Definition of Done — Data Product Standards",
                "  Glossary",
                "▶ SL 4J is the only management-track specialist publication",
                "  Cross-track coordination: works with SL 4G, H, M, K, L daily",
                "▶ Next Steps: Complete SL 4J → eligible for SL 5J (Advanced Program Manager)",
            ]),
        ]
    )


def deck_tm50_overview():
    build_deck(
        filename="TM_50_Advanced_Overview.pptx",
        title_lines=["ADVANCED MSS LEADERSHIP",
                     "SL 5 Series  ·  Platform Leadership  ·  Prereq: SL 4 (track-specific)"],
        poc="USAREUR-AF  ·  C2DAO  ·  MSS",
        header_text="ADVANCED MSS LEADERSHIP  ·  SL 5 OVERVIEW",
        slides_content=[
            ("WHAT CHANGES AT SL 5", [
                "SL 5 is not a continuation of SL 4 tasks — it is a change of role",
                "▶ SL 4: Practitioner",
                "  You execute within an established architecture",
                "  You build models, write transforms, manage sprints, maintain knowledge",
                "  Your scope is a project or a function",
                "▶ SL 5: Platform Leader",
                "  You design the architecture others work within",
                "  You set standards, review work, and grow junior practitioners",
                "  Your scope is a program, a capability, or an enterprise",
                "▶ The Core Shift: From Doing to Governing",
                "  SL 5 practitioners spend less time writing code and more time reviewing it",
                "  Less feature engineering; more feature store architecture",
                "  Less sprint planning; more program increment planning",
                "▶ Prerequisite: each SL 5 track requires its corresponding SL 4 track",
            ]),
            ("THE SIX SL 5 SPECIALIST TRACKS", [
                "▶ SL 5G — Advanced ORSA",
                "  Advanced optimization, agent-based modeling, Bayesian analysis",
                "  Wargame and campaign analysis; theater-level OR products",
                "▶ SL 5H — Advanced AI Engineer",
                "  Multi-agent orchestration, LLM fine-tuning, advanced RAG, AI red-teaming",
                "  Enterprise AI architecture and governance",
                "▶ SL 5M — Advanced ML Engineer",
                "  Advanced MLOps, federated learning, neural architectures, adversarial robustness",
                "  ML platform architecture and leadership",
                "▶ SL 5J — Advanced Program Manager",
                "  Scaled Agile, portfolio-level tracking, theater visibility, vendor management",
                "▶ SL 5K — Advanced Knowledge Manager",
                "  Enterprise KM architecture at corps/theater level, coalition knowledge sharing",
                "▶ SL 5L — Advanced Software Engineer",
                "  Platform SDK at scale, high-performance development, DevSecOps, security review",
            ]),
            ("PLATFORM LEADERSHIP RESPONSIBILITIES", [
                "All SL 5 practitioners share a common set of platform leadership duties:",
                "▶ Architecture Governance",
                "  Review designs proposed by SL 4 practitioners",
                "  Enforce standards; escalate architectural risks to C2DAO",
                "▶ Standards Setting",
                "  Define coding standards, documentation requirements, and review criteria",
                "  Own the Definition of Done for your domain",
                "▶ Mentorship and Team Development",
                "  Structured code review and design review for SL 4 practitioners",
                "  Succession planning: document institutional knowledge before rotation",
                "▶ Cross-Track Integration",
                "  SL 5 practitioners coordinate across all SL 5 tracks",
                "  No domain operates in isolation at enterprise scale",
                "▶ Governance Reporting",
                "  Produce domain status assessments for GO/SES leadership",
                "  Communicate technical debt, risk, and capacity constraints clearly",
            ]),
            ("ENTERPRISE RISK AND COMPLIANCE AT SL 5", [
                "▶ Policy Framework (common across all SL 5 tracks)",
                "  DoD Instruction 5000.89 — Test and Evaluation of AI/ML Systems",
                "  DoD RAIMTF — Responsible AI Implementation",
                "  Army AI Ethics Policy and USAREUR-AF implementing instructions",
                "  AR 25-2 — Army Cybersecurity (governs all platform code)",
                "▶ ATO (Authority to Operate) Responsibilities",
                "  SL 5L owns ATO support documentation for platform components",
                "  SL 5H/M own model-level risk documentation (RAIMTF)",
                "  SL 5J owns program-level governance and policy compliance framework",
                "▶ Institutional Memory Risk",
                "  At SL 5, you hold disproportionate institutional knowledge",
                "  Knowledge documentation is not optional — it is a platform resilience requirement",
                "  Model cards, architecture diagrams, decision records must survive your rotation",
            ]),
            ("SL 5 TRACK SELECTION AND ENROLLMENT", [
                "▶ Enrollment Prerequisite: Corresponding SL 4 Track Required",
                "  SL 5G: requires SL 4G (ORSA)",
                "  SL 5H: requires SL 4H (AI Engineer)",
                "  SL 5M: requires SL 4M (ML Engineer)",
                "  SL 5J: requires SL 4J (Program Manager)",
                "  SL 5K: requires SL 4K (Knowledge Manager)",
                "  SL 5L: requires SL 4L (Software Engineer)",
                "▶ Who Should Enroll in SL 5",
                "  Personnel who will lead a platform domain (not just contribute to it)",
                "  Personnel who supervise SL 4 practitioners",
                "  Senior civilian or warrant officer personnel in technical leadership roles",
                "▶ Enrollment is coordinator-approved — not self-service",
                "  Nomination by direct supervisor required; MSS Training Coordinator approval",
                "▶ SL 5 is not a competition — it is a responsibility assignment",
            ]),
            ("CROSS-TRACK SL 5 INTEGRATION", [
                "At enterprise scale, SL 5 tracks must coordinate as a leadership team:",
                "▶ SL 5L + SL 5H/M — Platform and AI/ML Alignment",
                "  SWE provides the platform; AI/ML operates on it",
                "  Shared Feature Store design, inference endpoint architecture, CI/CD pipeline",
                "▶ SL 5G + SL 5H/M — ORSA and AI/ML Integration",
                "  ORSA validates AI/ML model assumptions and interprets outputs",
                "  AI/ML operationalizes ORSA research findings at production scale",
                "▶ SL 5K + All Tracks — Knowledge Architecture",
                "  KM owns the ontology and knowledge graph that all tracks depend on",
                "  Ontology changes require SL 5K review before implementation",
                "▶ SL 5J — Program Integration Authority",
                "  PM is the integration point for all SL 5 domains",
                "  Dependency management, portfolio visibility, and commander reporting are PM-owned",
                "▶ Monthly SL 5 leadership sync recommended: dependency review + risk escalation",
            ]),
            ("FROM SL 4 TO SL 5 — THE READINESS CHECKLIST", [
                "Before enrolling in SL 5, practitioners should be able to:",
                "▶ Technical Readiness",
                "  Independently complete all SL 4 tasks without reference to the manual",
                "  Have reviewed the work of at least one other SL 4 practitioner",
                "  Have experienced at least one production incident in your domain",
                "▶ Leadership Readiness",
                "  Have given structured feedback on a design or code review",
                "  Have documented an architectural decision (ADR format preferred)",
                "  Have communicated technical status to a non-technical senior leader",
                "▶ Governance Readiness",
                "  Understand your domain's policy framework (DoD RAIMTF, AR 25-2, etc.)",
                "  Have completed at least one Model Card or architecture review artifact",
                "▶ Talk to your supervisor and MSS Training Coordinator before enrolling",
                "  SL 5 enrollment without readiness wastes institutional capacity",
            ]),
        ]
    )


def deck_tm50g():
    build_deck(
        filename="TM_50G_Advanced_ORSA.pptx",
        title_lines=["ADVANCED ORSA: THEATER-LEVEL ANALYSIS",
                     "SL 5G  ·  Advanced Specialist  ·  Prereq: SL 4G"],
        poc="USAREUR-AF  ·  C2DAO  ·  MSS",
        header_text="ADVANCED ORSA  ·  SL 5G  ·  Advanced Specialist",
        slides_content=[
            ("FROM SL 4G TO SL 5G — SCOPE CHANGE", [
                "▶ SL 4G (ORSA): Practitioner",
                "  Statistical modeling, time series, Monte Carlo, code workspaces",
                "  Products for brigade/division-level commanders",
                "▶ SL 5G (Advanced ORSA): Platform Leader",
                "  Advanced optimization, agent-based modeling, Bayesian/causal inference",
                "  Wargame support and campaign-level assessment architecture",
                "  Theater-level OR products for GO/SES audiences",
                "▶ Key Shifts at SL 5G",
                "  Analysis under deep uncertainty — not just risk quantification",
                "  Building persistent analytical environments, not one-off analyses",
                "  Leading OR teams and setting analytical standards",
                "  Integrating ML/AI outputs with ORSA products",
            ]),
            ("ADVANCED OPTIMIZATION (CHAPTER 2)", [
                "▶ When Linear Programming Is Insufficient",
                "  LP assumes linearity, certainty, and single objective — operational reality has none",
                "  Move to: Nonlinear, Multi-Objective, and Stochastic programming",
                "▶ Multi-Objective Optimization",
                "  COA comparison: optimize simultaneously for speed, risk, and logistics cost",
                "  Pareto frontier: surface the tradeoff, let the commander choose",
                "▶ Stochastic Programming",
                "  Scenarios drawn from probability distributions — not single-point estimates",
                "  Two-stage stochastic models: decide now, respond to uncertainty later",
                "▶ Metaheuristic Methods",
                "  Genetic algorithms, simulated annealing for large combinatorial problems",
                "  Use when exact optimization is computationally intractable",
                "▶ Key Task: Multi-Objective COA Comparison using Pareto analysis in MSS",
            ]),
            ("AGENT-BASED MODELING AND SIMULATION (CHAPTER 3)", [
                "▶ When to Use ABMS",
                "  Emergent behavior from individual agent interactions — not solvable analytically",
                "  Examples: adversary adaptation, convoy security, crowd dynamics",
                "▶ ABMS Architecture",
                "  Agents: each unit/entity has rules for perception, decision, and action",
                "  Environment: terrain, supply levels, threat data — loaded from MSS datasets",
                "  Rules: doctrine-aligned — agents follow FM-consistent decision logic",
                "▶ Mesa Framework (Python) on MSS",
                "  Mesa: approved ABMS library for MSS code workspaces",
                "  Calibration: tune agent parameters against historical operational data",
                "▶ Advanced ABMS: Adversary Adaptation",
                "  RED force agents adapt to BLUE actions — not static threat modeling",
                "  Enables robust COA testing against adaptive, not scripted, opposition",
            ]),
            ("BAYESIAN ANALYSIS AND CAUSAL INFERENCE (CHAPTER 4)", [
                "▶ Bayesian Framework for Operational Analysis",
                "  Prior belief + new evidence → updated belief",
                "  Explicitly handles small data and expert knowledge incorporation",
                "  Produces probability distributions, not point estimates",
                "▶ Bayesian Networks",
                "  Model causal relationships between operational variables",
                "  Use case: readiness cascade — how maintenance backlog affects C-rating",
                "▶ Causal Inference for Operational Analysis",
                "  Correlation ≠ causation — especially in operational data",
                "  Counterfactual analysis: 'What would readiness have been without intervention X?'",
                "  DoWhy and CausalML frameworks approved for MSS use",
                "▶ Key Distinction: Bayesian inference (updating beliefs) vs. causal inference (mechanism)",
                "  Both required for rigorous theater-level analysis",
            ]),
            ("WARGAME AND CAMPAIGN ANALYSIS (CHAPTER 5)", [
                "▶ OR Support to MDMP — The Analytical Cell",
                "  ORSA embedded in planning team: quantifies COA comparison criteria",
                "  Products at each MDMP step: from mission analysis through COA approval",
                "▶ Designing OR Support to Multi-Echelon Wargames",
                "  Data architecture for wargame: inputs, adjudication data, outputs",
                "  Real-time analysis during wargame execution — not post-hoc",
                "▶ Campaign-Level Assessment Architecture",
                "  Persistent model library supports theater campaign assessment",
                "  MOE/MOP framework: what we measure and why, defined before operations begin",
                "▶ RED Force Modeling",
                "  Adversary decision modeling — not just capability inventory",
                "  Threat COA probability estimates fed from intelligence data in MSS",
            ]),
            ("DECISION ANALYSIS UNDER DEEP UNCERTAINTY (CHAPTER 6)", [
                "▶ The Deep Uncertainty Problem",
                "  Some operational uncertainties cannot be quantified with probabilities",
                "  Deep uncertainty: we do not know the probability distribution itself",
                "▶ XLRM Framework",
                "  X: Exogenous uncertainties (what we cannot control)",
                "  L: Levers (COA options)",
                "  R: Relationships (models connecting X and L to outcomes)",
                "  M: Measures of performance",
                "▶ Robustness Analysis",
                "  Identify COAs that perform acceptably across ALL plausible futures",
                "  Not 'best expected outcome' — 'least regret under worst case'",
                "▶ Scenario-Based Planning and Adaptive Strategy",
                "  Decision trees with branch conditions — when to switch COA based on indicators",
                "▶ ORSA communicates uncertainty honestly to GO/SES — not false precision",
            ]),
            ("SENIOR-LEVEL OR PRODUCTS AND TEAM STANDARDS (CHAPTERS 7–8)", [
                "▶ Understanding the GO/SES Audience",
                "  Senior leaders need: What is the answer? Why should I trust it? What is the risk?",
                "  They do not need: methodology detail, code output, literature citations",
                "▶ OR Product Structure for Senior Audiences",
                "  BLUF: one sentence answer",
                "  Key assumptions: what had to be true for this answer to hold",
                "  Sensitivity: how much does the answer change if assumptions shift",
                "  Uncertainty communication standard: provide ranges, not false precision",
                "▶ Building Persistent OR Capability on MSS",
                "  Model library: documented, versioned, reproducible analyses in Foundry",
                "  Model Cards required for all production OR models",
                "  Team standards: peer review before any product reaches a senior leader",
                "▶ Succession Planning",
                "  Knowledge does not leave with the ORSA — model library and documentation persist",
                "▶ Integrating ML/AI with ORSA: ORSA validates; ML operationalizes",
            ]),
        ]
    )


def deck_tm50h():
    build_deck(
        filename="TM_50H_Advanced_AI_Engineer.pptx",
        title_lines=["ADVANCED AI ENGINEERING",
                     "SL 5H  ·  Advanced Specialist  ·  Prereq: SL 4H"],
        poc="USAREUR-AF  ·  C2DAO  ·  MSS",
        header_text="ADVANCED AI ENGINEERING  ·  SL 5H  ·  Advanced Specialist",
        slides_content=[
            ("FROM SL 4H TO SL 5H — SCOPE CHANGE", [
                "▶ SL 4H (AI Engineer): Practitioner",
                "  Build AIP Logic workflows, Agent Studio agents, single-system RAG pipelines",
                "  Scope: individual AI applications",
                "▶ SL 5H (Advanced AI Engineer): Platform Leader",
                "  Design multi-agent orchestration systems and enterprise AI architecture",
                "  LLM fine-tuning, advanced RAG, AI red-teaming, production observability",
                "  Scope: the AI platform that others build on",
                "▶ The Senior AI Engineer Responsibility",
                "  Set engineering standards for all AI systems on MSS",
                "  Conduct production readiness reviews and architecture reviews",
                "  Own DoD RAIMTF and Army AI Ethics compliance for the platform",
            ]),
            ("MULTI-AGENT ORCHESTRATION (CHAPTER 2)", [
                "▶ Why Multi-Agent Architecture",
                "  Complex operational workflows exceed single-agent context and capability",
                "  Parallel execution: multiple agents working concurrently on sub-tasks",
                "▶ Orchestration Patterns",
                "  Sequential: agent A outputs feed agent B — strict dependency chain",
                "  Parallel: agents A, B, C run simultaneously; orchestrator aggregates",
                "  Hierarchical: orchestrator agent delegates to specialist sub-agents",
                "▶ Shared State Management",
                "  Agents write to and read from shared state object in Foundry",
                "  Conflict resolution: last-writer-wins or priority-based, must be explicit",
                "▶ Circuit Breakers and Failure Isolation",
                "  Single agent failure must not cascade to entire orchestration",
                "  Circuit breaker pattern: fail fast, fallback to degraded mode, alert",
                "▶ Key Task: Design a multi-agent orchestration system with failure isolation",
            ]),
            ("LLM FINE-TUNING AND DOMAIN ADAPTATION (CHAPTER 3)", [
                "▶ When to Fine-Tune vs. When to Prompt",
                "  Fine-tune: persistent domain knowledge, style, or format requirements",
                "  Prompt: task-specific instructions that change frequently",
                "  Rule: exhaust prompt engineering and RAG before committing to fine-tuning",
                "▶ Policy and Legal Requirements",
                "  Training data must be authorized for use — no scraping operational data",
                "  DoD data rights and classification implications reviewed before fine-tuning begins",
                "▶ Training Data Preparation",
                "  Instruction-tuning format: (prompt, response) pairs from Army writing samples",
                "  Data quality gate: reviewed by subject matter expert before training",
                "▶ Fine-Tuning Methods: LoRA/QLoRA preferred (parameter-efficient)",
                "▶ Key Task: Domain adaptation of an LLM for Army writing style (BLUF, Army Active Voice)",
            ]),
            ("ADVANCED RAG ARCHITECTURE (CHAPTER 4)", [
                "▶ Limitations of Baseline RAG",
                "  Single retrieval step misses context for multi-hop questions",
                "  Keyword similarity fails for semantic queries in military domain",
                "▶ Hybrid Retrieval: Dense + Sparse",
                "  BM25 (keyword) + embedding (semantic) combined for better recall",
                "  Critical for military acronyms and technical terminology",
                "▶ Re-Ranking",
                "  Cross-encoder re-ranks retrieved chunks before passing to LLM",
                "  Improves precision without sacrificing recall",
                "▶ Query Transformation",
                "  HyDE: generate hypothetical answer, retrieve against it",
                "  Step-back prompting: abstract the query before retrieval",
                "▶ Corpus Quality Management",
                "  Stale documents degrade RAG performance — corpus freshness is operational",
                "▶ RAG Evaluation Pipeline: RAGAS or equivalent — automated + human spot-check",
            ]),
            ("AI RED-TEAMING AND ADVERSARIAL TESTING (CHAPTER 5)", [
                "▶ Why Red-Team AI Systems",
                "  Operational AI systems are targets — adversaries will attempt manipulation",
                "  Internal testing identifies failure modes before deployment",
                "▶ Threat Model for USAREUR-AF AI Systems",
                "  Prompt injection: adversarial inputs that hijack agent behavior",
                "  Data poisoning: corrupting training data to degrade model performance",
                "  Model extraction: inferring model capabilities through repeated queries",
                "▶ Red-Team Execution Process",
                "  Team: AI engineer + domain expert (ORSA/S2) + security officer",
                "  Scope: define what is in/out of scope before red team begins",
                "  Output: Red-Team Report with findings categorized by severity and mitigations",
                "▶ Mitigations and Defensive Patterns",
                "  Input validation, output filtering, rate limiting, anomaly detection on queries",
                "▶ Red-team is required before production deployment of any agent system",
            ]),
            ("PRODUCTION AI OBSERVABILITY (CHAPTER 6)", [
                "▶ Why AI Observability Is Different from Software Monitoring",
                "  Software metrics: latency, error rate, uptime",
                "  AI metrics: output quality, hallucination rate, retrieval relevance, input distribution",
                "▶ Four Observability Dimensions",
                "  Service health: latency, error rate, availability (standard)",
                "  Input distribution: are queries shifting from what the system was trained on?",
                "  Retrieval quality: are RAG retrievals relevant? (automated scoring)",
                "  Output quality: hallucination detection, factual grounding checks",
                "▶ AI Observability Dashboard",
                "  Built on MSS Workshop: one page, four panels, daily cadence",
                "  Alerts: output quality below threshold triggers SL 5H review",
                "▶ Hallucination Detection",
                "  Faithfulness score: does output match retrieved context?",
                "  Answer relevance score: does output address the query?",
                "▶ Key Task: Implement AI observability for a production system end-to-end",
            ]),
            ("ENTERPRISE AI ARCHITECTURE AND GOVERNANCE (CHAPTER 8)", [
                "▶ Enterprise AI Architecture Principles",
                "  Shared infrastructure: embedding models, vector stores, LLM endpoints shared",
                "  No siloed AI systems — all AI builds on common platform components",
                "  Every AI component documented in the enterprise AI architecture registry",
                "▶ Architecture Review Process",
                "  Any new AI system requires SL 5H architecture review before build",
                "  Review criteria: alignment to platform, security, governance, observability",
                "▶ DoD RAIMTF and Army AI Ethics Compliance",
                "  Six principles: Responsible, Equitable, Traceable, Reliable, Governable",
                "  SL 5H is the compliance point for AI systems — not legal, not compliance office",
                "▶ AI Governance Documentation Requirements",
                "  Model/system card, Red-Team Report, Observability Plan, Incident Response Plan",
                "▶ AI Production Readiness Review (PRR)",
                "  SL 5H chairs PRR before any AI system goes to production",
                "  Checklist covers: security, observability, governance, degraded mode behavior",
            ]),
        ]
    )


def deck_tm50m():
    build_deck(
        filename="TM_50M_Advanced_ML_Engineer.pptx",
        title_lines=["ADVANCED ML ENGINEERING",
                     "SL 5M  ·  Advanced Specialist  ·  Prereq: SL 4M"],
        poc="USAREUR-AF  ·  C2DAO  ·  MSS",
        header_text="ADVANCED ML ENGINEERING  ·  SL 5M  ·  Advanced Specialist",
        slides_content=[
            ("FROM SL 4M TO SL 5M — SCOPE CHANGE", [
                "▶ SL 4M: Build and deploy individual models following governance process",
                "▶ SL 5M: Own the ML platform that all models operate within",
                "  Automated retraining pipelines; advanced neural architectures",
                "  Federated learning for multi-classification environments",
                "  Interpretability and bias auditing at enterprise scale",
                "  ML security, adversarial robustness, and supply chain security",
                "▶ Key Responsibility Shift",
                "  From: submit models through governance gates",
                "  To: design the governance gate process and review others' submissions",
                "  From: implement drift detection for your model",
                "  To: build the shared drift detection infrastructure for all models",
            ]),
            ("ADVANCED MLOPS AND AUTOMATED RETRAINING (CHAPTER 2)", [
                "▶ Automated Retraining Pipeline Architecture",
                "  Trigger sources: PSI drift threshold, scheduled cadence, data volume threshold",
                "  Pipeline stages: data validation → feature refresh → train → evaluate → gate → deploy",
                "  Gate: automated evaluation must pass before promotion — no human bypass",
                "▶ A/B Model Deployment",
                "  Traffic split: new model serves X% of requests alongside production model",
                "  Evaluation window: define in advance — not 'until it looks good'",
                "  Rollback trigger: automated if new model performance drops below threshold",
                "▶ Canary Releases",
                "  Canary: 5% traffic, extended monitoring before full rollout",
                "  Canary failure: automatic rollback, incident report, SL 5M review",
                "▶ Production Readiness Gates — Advanced Version",
                "  All SL 4M gates plus: automated retraining test, adversarial robustness check",
            ]),
            ("ADVANCED NEURAL ARCHITECTURES (CHAPTER 3)", [
                "▶ Transformers for Operational Text",
                "  SITREP classification, OPORD parsing, intelligence report summarization",
                "  Fine-tuning: start from pre-trained checkpoint, not from scratch",
                "  Key Task: Fine-tune transformer for SITREP priority classification",
                "▶ Graph Neural Networks for C2 Network Analysis",
                "  C2 network as graph: units = nodes, command relationships = edges",
                "  GNN predicts readiness propagation across echelons",
                "  Key Task: Build GNN for unit readiness propagation (degraded unit → impact on BCT)",
                "▶ When NOT to Use Neural Architectures",
                "  Small datasets (<10k samples): simpler models outperform deep learning",
                "  Interpretability required: tree-based models preferred for commander-facing outputs",
                "  Compute cost: neural inference at scale requires architecture review",
            ]),
            ("FEDERATED AND PRIVACY-PRESERVING ML (CHAPTER 4)", [
                "▶ Why Federated Learning on MSS",
                "  Multi-classification environment: SECRET and UNCLASSIFIED data cannot be pooled",
                "  Unit-level data sensitivity: some readiness data cannot leave unit workspace",
                "▶ Federated Learning Architecture on MSS",
                "  Central coordinator: aggregates model updates, not raw data",
                "  Participating units: train locally on local data, share only gradient updates",
                "  Privacy preservation: differential privacy noise added before aggregation",
                "▶ Federated Learning Scenarios",
                "  Cross-unit readiness modeling: BCTs train locally, aggregate at division",
                "  Classification-level models: parallel architectures, no cross-domain data flow",
                "▶ Differential Privacy for MSS",
                "  Epsilon-delta privacy bounds required for any federated model",
                "  Privacy budget tracked across queries — no unlimited querying",
            ]),
            ("INTERPRETABILITY, BIAS AUDITING, AND ML SECURITY (CHAPTERS 5, 7)", [
                "▶ Platform-Level Interpretability Infrastructure",
                "  Shared SHAP infrastructure in MSS — not rebuilt per model",
                "  Explanation artifacts stored alongside model predictions in Foundry",
                "  LIME for local explanations: explain individual predictions to operators",
                "▶ Bias Auditing for Operational Models",
                "  Personnel-affecting models: bias audit required across MOS, echelon, gender",
                "  Bias audit report produced by SL 5M; reviewed by independent reviewer",
                "▶ ML Security Threat Model",
                "  Adversarial inputs: inputs crafted to cause misclassification",
                "  Data poisoning: corrupt training data to degrade production model",
                "  Model extraction: infer model via repeated API queries",
                "▶ Adversarial Robustness Evaluation",
                "  FGSM and PGD attacks tested on all production classification models",
                "  Robustness threshold defined and enforced at production gate",
                "▶ Supply chain security: verify provenance of all pre-trained model weights",
            ]),
            ("ML PLATFORM ARCHITECTURE AND LEADERSHIP (CHAPTER 8)", [
                "▶ Shared Feature Store Design",
                "  Central registry: prevents duplicate feature computation across teams",
                "  Feature versioning: reproducibility guarantee for all production models",
                "  Feature Store governance: SL 5M approves new features before registration",
                "▶ Model Registry Architecture",
                "  Centralized: all models in one registry with metadata, lineage, and status",
                "  Status tiers: Experimental → Staging → Production → Retired",
                "▶ Experiment Tracking at Scale",
                "  MLflow at enterprise scale: shared server, team-level namespacing",
                "  Experiment hygiene: stale experiments archived quarterly",
                "▶ Cross-Track Integration",
                "  SL 5H (AI Engineer): share code workspace infrastructure and CI/CD",
                "  SL 4G (ORSA): validate model assumptions; share feature store",
                "▶ Code Review Standards for Production ML",
                "  Every production model code reviewed by SL 5M before deployment",
                "  Review criteria: correctness, efficiency, governance compliance, documentation",
            ]),
            ("LEADING ML CAPABILITY — PRODUCTION READINESS AND TEAM DEVELOPMENT", [
                "▶ Production Readiness Gates (Advanced) — SL 5M Owns the Process",
                "  Automated retraining pipeline tested and verified",
                "  Real-time inference endpoint load-tested (if applicable)",
                "  Advanced neural architecture: interpretability artifacts generated",
                "  Federated learning: privacy budget documented and approved",
                "  Interpretability and bias assessment completed and signed",
                "  Security: adversarial robustness evaluation passed",
                "  Platform governance: model in registry, lineage documented",
                "▶ Developing Junior ML Engineers",
                "  Structured code review: explain why, not just what to fix",
                "  Pair programming on first model deployment: reduces errors, builds standards",
                "  Succession planning: every critical ML system has a documented backup person",
                "▶ ML Capability Roadmap: SL 5M contributes to the enterprise ML architecture",
                "  Annual capability review; identify gaps in platform before they become incidents",
            ]),
        ]
    )


def deck_tm50j():
    build_deck(
        filename="TM_50J_Advanced_Program_Manager.pptx",
        title_lines=["ADVANCED PROGRAM MANAGEMENT",
                     "SL 5J  ·  Advanced Specialist  ·  Prereq: SL 4J"],
        poc="USAREUR-AF  ·  C2DAO  ·  MSS",
        header_text="ADVANCED PROGRAM MANAGEMENT  ·  SL 5J  ·  Advanced Specialist",
        slides_content=[
            ("FROM SL 4J TO SL 5J — SCOPE CHANGE", [
                "▶ SL 4J (PM): Manage one data/AI project team",
                "▶ SL 5J (Advanced PM): Manage a portfolio of concurrent data/AI programs",
                "  Multiple teams, multiple dependencies, theater-level visibility",
                "  GO/SES-facing reporting; scaled Agile across teams",
                "  Vendor management (Palantir), contractor management, budget stewardship",
                "▶ The SL 5J PM Profile",
                "  Works at Division/Corps/Theater level — not BCT-level",
                "  Manages program managers, not just engineers",
                "  Serves as primary interface between technical program and senior leadership",
            ]),
            ("SCALED AGILE — MANAGING MULTIPLE CONCURRENT TEAMS (CHAPTER 2)", [
                "▶ From Project to Program",
                "  Project: one team, one product, one PM",
                "  Program: multiple teams, shared dependencies, portfolio PM",
                "▶ Program Increment (PI) Planning",
                "  Quarterly planning event: all teams define objectives and identify dependencies",
                "  PI objectives reviewed by senior leader before sprint execution begins",
                "  Risk escalation: blockers that cross team boundaries escalated in PI planning",
                "▶ Managing Dependencies Across Teams",
                "  Dependency log maintained in MSS — not in each team's backlog",
                "  Dependency review: weekly at program level, not ad hoc between engineers",
                "▶ Conflict Resolution Between Teams",
                "  Shared resource conflicts escalated to PM — not resolved engineer-to-engineer",
                "  PM arbitrates based on portfolio priority, not loudest voice",
            ]),
            ("PORTFOLIO-LEVEL TRACKING AND THEATER VISIBILITY (CHAPTER 3)", [
                "▶ Why Portfolio Tracking at SL 5J Level",
                "  Theater commander needs: what is working, what is not, what is at risk",
                "  Not sprint velocity — strategic delivery against mission requirements",
                "▶ Program Tracking Ontology Design",
                "  Object Types: Program, Team, Milestone, Risk, Dependency, Resource",
                "  Link Types: Team-to-Program, Milestone-to-Program, Risk-to-Program",
                "▶ Designing for GO/SES Readability",
                "  One-page portfolio view: program status, milestone health, risk heatmap",
                "  Traffic light status with one-line description — no jargon",
                "  Data quality: stale status worse than no status — enforce update cadence",
                "▶ Program Review Briefing Package",
                "  Three slides: what we planned, what we delivered, what is at risk next quarter",
                "  Attach supporting data in MSS — brief talks to the dashboard, not to slides",
            ]),
            ("STAKEHOLDER MANAGEMENT AND DATA CULTURE (CHAPTER 4)", [
                "▶ Managing Up: Briefing GO/SES Leadership",
                "  Bottom Line Up Front — always",
                "  Anticipate questions: prepare backup slides for likely concerns",
                "  Uncertainty communication: 'This forecast has ±15% confidence' not 'it depends'",
                "▶ Building Organizational Buy-In for Data Culture Change",
                "  Champion strategy: identify early adopters in each directorate",
                "  Demonstrate value before expecting behavior change",
                "  'Show me on MSS' replaces 'send me a spreadsheet'",
                "▶ Navigating Command Climate",
                "  Some commanders resist data-driven decision-making — manage the narrative",
                "  Vignette: Commander wants a specific number, not a range — PM must not fabricate precision",
                "▶ Stakeholder Mapping at Program Level",
                "  Power/interest matrix: who can block the program vs. who is just informed",
            ]),
            ("QUANTITATIVE DELIVERY MANAGEMENT (CHAPTER 5)", [
                "▶ Core Delivery Metrics at Program Level",
                "  Velocity trend (not snapshot): is the team accelerating or decelerating?",
                "  Predictability: % of sprint commitments delivered as planned",
                "  Dependency resolution rate: how quickly are cross-team blockers resolved?",
                "  Cycle time: idea to production — the true delivery clock",
                "▶ Interpreting Metrics Correctly",
                "  Velocity drop: diagnose before reacting — could be planned debt reduction",
                "  High predictability + low velocity: team is sandbagging estimates",
                "  Low predictability + high velocity: team is overcommitting",
                "▶ Program-Level Predictability Targets",
                "  70% of PI objectives delivered as committed — Army standard for mature programs",
                "▶ Communicating Metrics to Leadership",
                "  Metrics tell a story — PM frames the story, not the data analyst",
                "  Never present metrics without interpretation and recommended action",
            ]),
            ("VENDOR AND CONTRACTOR MANAGEMENT (CHAPTER 7)", [
                "▶ The Palantir Partnership Model",
                "  Palantir FDEs: technical resources embedded with the MSS program",
                "  FDE scope: technical implementation only — not requirements definition",
                "  PM owns requirements; Palantir implements; government retains decision authority",
                "▶ Task Order Management",
                "  SL 5J PM manages task order scope, deliverables, and acceptance criteria",
                "  Change requests to task order go through PM — not direct FDE-to-user negotiation",
                "▶ Managing Palantir Relationship Day-to-Day",
                "  Weekly sync: FDE lead + PM — scope, blockers, velocity",
                "  Do not allow FDEs to independently prioritize backlog — PM owns priority",
                "▶ Vendor Performance Documentation",
                "  Delivery log: what was committed vs. delivered, by sprint",
                "  Performance documentation supports future contract negotiations",
            ]),
            ("TEAM DESIGN, GOVERNANCE, AND PROGRAM LEADERSHIP (CHAPTERS 8–9)", [
                "▶ Designing the Data Program Team Structure",
                "  Stream-aligned teams: one team per bounded capability domain",
                "  Platform team: shared infrastructure supporting stream teams",
                "  Enabling team: governance, standards, security — support function, not gatekeeper",
                "▶ Skill Mix Planning and Personnel Turbulence",
                "  At Army rotation rates, the program loses key personnel every 12–18 months",
                "  Cross-training and documentation are not optional — they are succession planning",
                "▶ Program Governance Framework",
                "  Decision Authority Matrix: who approves what — documented, not implied",
                "  Change control for production systems: all changes go through PM approval",
                "▶ Program Communication Architecture",
                "  Weekly: team syncs, dependency review",
                "  Monthly: program review with senior stakeholder",
                "  Quarterly: PI planning with full leadership attendance",
                "▶ Writing the Program SITREP: status + forecast + risk + ask",
            ]),
        ]
    )


def deck_tm50k():
    build_deck(
        filename="TM_50K_Advanced_Knowledge_Manager.pptx",
        title_lines=["ADVANCED KNOWLEDGE MANAGEMENT",
                     "SL 5K  ·  Advanced Specialist  ·  Prereq: SL 4K"],
        poc="USAREUR-AF  ·  C2DAO  ·  MSS",
        header_text="ADVANCED KNOWLEDGE MANAGEMENT  ·  SL 5K  ·  Advanced Specialist",
        slides_content=[
            ("FROM SL 4K TO SL 5K — ARCHITECT, NOT OPERATOR", [
                "▶ SL 4K: Operate and maintain knowledge architecture within a unit",
                "▶ SL 5K: Design the architecture that units operate within",
                "  Corps/theater-level ontology design",
                "  Federated architecture across classification levels",
                "  Coalition knowledge sharing under NATO/partner constraints",
                "  Knowledge risk management and institutional memory preservation",
                "▶ The Advanced KM Profile",
                "  Less: day-to-day content entry and taxonomy maintenance",
                "  More: architecture decisions, governance structures, standards enforcement",
                "  Audience: corps G2/G6 leadership, theater J6, NATO KM officers",
            ]),
            ("ENTERPRISE KM ARCHITECTURE — CORPS/THEATER LEVEL (CHAPTER 2)", [
                "▶ Federated Architecture Principles",
                "  Each echelon maintains its own KM domain — not one monolithic ontology",
                "  Federation: lower echelons share up; higher echelons provide reference architecture",
                "  No single point of failure: federated design survives degraded network",
                "▶ Theater Ontology Design",
                "  Core object types shared across all echelons: Unit, Person, Equipment, Location, Event",
                "  Domain-specific extensions: each WFF adds functional object types",
                "  Versioning: ontology changes require architecture review and controlled rollout",
                "▶ Federated Repository Design Patterns",
                "  Canonical repository: theater-level authoritative source",
                "  Shadow repositories: unit-level with sync to theater on defined schedule",
                "▶ Governance Structures",
                "  KM Architecture Review Board: corps G6 + SL 5K leads + C2DAO representative",
            ]),
            ("CROSS-DOMAIN KNOWLEDGE FEDERATION (CHAPTER 3)", [
                "▶ Classification-Level Parallel Architecture",
                "  Separate MSS instances at UNCLASSIFIED and SECRET — no cross-domain flow",
                "  Knowledge bridge: manual transfer only, with classification review",
                "  Taxonomy alignment: both instances use same object type names, different content",
                "▶ Multi-Unit Knowledge Federation",
                "  Coalition partner units: NATO interoperability via STANAG-aligned object types",
                "  Adjacent units: shared link types connect friendly force knowledge domains",
                "▶ Functional Directorate Integration",
                "  Each directorate (G2, G3, G4, G6) has a KM domain — SL 5K coordinates",
                "  Cross-directorate links: SL 5K designs the bridge between domains",
                "▶ Theater-Level Exercise Knowledge Federation",
                "  Exercise-specific KM environment: separate from operational KM",
                "  Post-exercise: lessons harvested and migrated to operational lessons learned system",
            ]),
            ("AI-ASSISTED DOCTRINE AND LESSONS LEARNED AT SCALE (CHAPTER 4)", [
                "▶ Scale Problem: Lessons Learned",
                "  Theater-level exercises generate thousands of AAR entries",
                "  Manual synthesis takes months — value decays before it is actionable",
                "▶ AIP Capabilities for Advanced KM",
                "  Automated tagging: AIP Logic assigns taxonomy tags to unstructured AAR text",
                "  Semantic clustering: group similar lessons without manual reading",
                "  Synthesis: AI-generated summary of theme clusters — human reviews, not trusts",
                "▶ Large-Scale Lessons Synthesis Workflow",
                "  Ingest → Tag → Cluster → Synthesize → Human review → Publish",
                "  Human is the final reviewer — AIP accelerates, not replaces",
                "▶ AI-Assisted Doctrine Development",
                "  Pattern extraction from lessons: 'These three lessons all point to the same gap'",
                "  Doctrine draft supported by evidence from lessons — not from anecdote",
                "▶ Bulk Tagging Quality Control: random sample audit at 10% for accuracy verification",
            ]),
            ("COALITION KNOWLEDGE SHARING (CHAPTER 5)", [
                "▶ Legal and Policy Framework",
                "  SOFA, bilateral agreements, NATO information sharing agreements govern scope",
                "  No sharing without written authorization — SL 5K knows the boundary",
                "▶ NATO Knowledge Management Integration",
                "  STANAG 4778: NATO standard for lessons learned systems",
                "  SL 5K designs MSS-STANAG 4778 conformance mapping",
                "▶ Coalition Knowledge Sharing Architecture",
                "  Partner-accessible workspace: sanitized object types visible to coalition",
                "  Releasability markings enforced at the ontology level, not by user discretion",
                "▶ Partner Nation Knowledge Integration Patterns",
                "  Pull pattern: partner uploads to shared workspace, SL 5K reviews and integrates",
                "  Push pattern: US shares selected products to coalition workspace after review",
                "▶ Federated Mission Networking (FMN): technical reference for NATO interop",
            ]),
            ("KNOWLEDGE GRAPH DESIGN AND MAINTENANCE (CHAPTER 6)", [
                "▶ Knowledge Graph for Military KM",
                "  Nodes: entities (units, personnel, doctrine, equipment, events)",
                "  Edges: relationships (commanded-by, supports, supersedes, references)",
                "  Properties: temporal (effective dates), classification, source",
                "▶ Knowledge Graph Architecture Design",
                "  Modular: each domain is a subgraph — composable, not monolithic",
                "  Temporal: all edges have effective-date range — no timeless facts",
                "  Classification-aware: edge visibility filtered by user clearance",
                "▶ Knowledge Graph Maintenance at Scale",
                "  Automated quality checks: orphaned nodes, broken links, stale entries",
                "  Merge resolution: when two units report conflicting facts about same entity",
                "▶ Long-Term Graph Evolution",
                "  Schema changes: backward-compatible by default; breaking changes require migration plan",
                "  Version control: graph schema versioned same as software code",
            ]),
            ("INSTITUTIONAL MEMORY AND KM RISK MANAGEMENT (CHAPTER 8)", [
                "▶ The Institutional Memory Problem in USAREUR-AF",
                "  Key personnel rotate every 12–24 months; institutional knowledge leaves with them",
                "  Consequence: same mistakes repeated across rotations; no compounding learning",
                "▶ Knowledge Risk Assessment",
                "  Identify critical knowledge dependencies: what breaks if person X leaves?",
                "  Rate by impact and replaceability — not all knowledge is equally critical",
                "▶ Knowledge Elicitation Methodology",
                "  Structured interviews with departing personnel: focus on undocumented practice",
                "  Decision journals: record why decisions were made, not just what was decided",
                "▶ Designing for Personnel Continuity",
                "  No single-person-dependent workflows in production KM systems",
                "  Runbooks: step-by-step procedures for every critical KM operation",
                "▶ Leading the KM Community of Practice",
                "  Monthly CoP: share lessons learned, surface standards questions, coordinate governance",
                "▶ KM Architecture Assessment and Audit: annual review of conformance and effectiveness",
            ]),
        ]
    )


def deck_tm50l():
    build_deck(
        filename="TM_50L_Advanced_Software_Engineer.pptx",
        title_lines=["ADVANCED SOFTWARE ENGINEERING",
                     "SL 5L  ·  Advanced Specialist  ·  Prereq: SL 4L"],
        poc="USAREUR-AF  ·  C2DAO  ·  MSS",
        header_text="ADVANCED SOFTWARE ENGINEERING  ·  SL 5L  ·  Advanced Specialist",
        slides_content=[
            ("FROM SL 4L TO SL 5L — SCOPE CHANGE", [
                "▶ SL 4L: Build OSDK applications and Foundry integrations as a practitioner",
                "▶ SL 5L: Own the platform architecture that practitioners build on",
                "  Platform SDK: infrastructure-level integration at scale",
                "  High-performance development: query optimization, caching, indexing",
                "  Multi-tenant architecture and classification-aware design",
                "  DevSecOps: CI/CD pipeline architecture, automated testing, ATO support",
                "  Security assessment and secure code review leadership",
                "▶ SL 5L in USAREUR-AF",
                "  Primary technical authority for MSS platform code",
                "  ATO documentation owner for platform components",
                "  Code review authority for all production deployments",
            ]),
            ("PLATFORM SDK — INFRASTRUCTURE-LEVEL INTEGRATION (CHAPTER 2)", [
                "▶ Platform SDK vs. OSDK",
                "  OSDK (SL 4L): query ontology objects and execute actions via API",
                "  Platform SDK (SL 5L): manage datasets, branches, file resources at infrastructure level",
                "▶ Platform SDK Authentication at Scale",
                "  Service account management: least-privilege, rotated quarterly",
                "  Token management: short-lived tokens, no hardcoded credentials",
                "▶ Dataset Operations at Scale",
                "  Bulk read/write patterns: pagination, streaming, parallel processing",
                "  Schema evolution: backward-compatible changes only; breaking = migration plan",
                "▶ Branch Management Automation",
                "  Automated branch creation, merge, and promotion pipelines",
                "  Branch naming conventions enforced programmatically",
                "▶ File Resource Management at Scale",
                "  Versioned file resources: media, documents, large binaries in Foundry",
                "  Lifecycle management: archive vs. delete vs. retain policies",
            ]),
            ("HIGH-PERFORMANCE FOUNDRY DEVELOPMENT (CHAPTER 3)", [
                "▶ Query Optimization for OSDK",
                "  Filter early, not late: push filters to the data source, not the application",
                "  Index-aware queries: know which object properties are indexed before querying",
                "  N+1 query problem: batch related object fetches, not loop-per-object",
                "▶ Caching Strategies",
                "  Object-level cache: TTL-based; invalidated on object update event",
                "  Query-level cache: hash query parameters; appropriate for low-change data",
                "  Never cache classification-sensitive data across security boundaries",
                "▶ Indexing and Search Optimization",
                "  Full-text search: Foundry Search API over indexed properties",
                "  Composite indexes for multi-property filter patterns",
                "▶ Compute Cost Management",
                "  Profile before optimizing — measurement precedes optimization",
                "  Compute cost reported monthly; SL 5L reviews unexplained cost increases",
            ]),
            ("MULTI-TENANT ARCHITECTURE AND CLASSIFICATION-AWARE DESIGN (CHAPTER 4)", [
                "▶ Tenant Architecture in MSS",
                "  Tenant: a logically isolated workspace (unit, classification level, coalition partner)",
                "  Data isolation: tenant data never crosses tenant boundary without explicit authorization",
                "▶ Implementing Tenant Isolation in Code",
                "  Tenant context propagated through all API calls — never inferred",
                "  Row-level security: implemented at data layer, not in application logic",
                "  Test isolation: each tenant's test suite runs against isolated test data",
                "▶ Classification-Aware Design",
                "  Classification embedded in data schema: every record has a classification property",
                "  Enforcement at read time: user clearance checked before data returned",
                "  Aggregation risk: SL 5L reviews any feature that combines data across classifications",
                "▶ Coalition Interoperability Patterns",
                "  Read-only API for partner-facing endpoints — no write access to US data",
                "  Releasability enforcement: API filters out non-releasable records automatically",
            ]),
            ("DEVSECOPS FOR FOUNDRY ENVIRONMENTS (CHAPTER 6)", [
                "▶ CI/CD Pipeline Architecture for MSS",
                "  Stages: lint → unit test → integration test → security scan → staging → production",
                "  No production deployment without passing all stages — no exceptions",
                "  Pipeline-as-code: Jenkinsfile or GitLab CI in version control",
                "▶ Automated Testing Standards",
                "  Unit test coverage: minimum 80% for production code; enforced by pipeline",
                "  Integration tests: run against Foundry staging environment — not mocks",
                "  Contract tests: API contracts between services tested automatically",
                "▶ Ontology CI",
                "  Ontology changes submitted as code PRs — reviewed same as application code",
                "  Schema compatibility check automated: breaking change = pipeline fail",
                "▶ Container Deployment Pattern",
                "  Docker containers for all application deployments — no manual server configuration",
                "  Container image scanning: vulnerability check before deployment",
                "▶ Pre-Commit Hooks: lint, format, secret detection — fail fast before push",
            ]),
            ("SECURITY ASSESSMENT AND SECURE CODE REVIEW (CHAPTER 7)", [
                "▶ OWASP Top 10 in the Foundry Context",
                "  Injection: parameterized OSDK queries — never string-concatenated filters",
                "  Broken access control: tenant isolation checks must be explicit, not assumed",
                "  Cryptographic failures: no plaintext sensitive data in logs, environment variables",
                "  Security misconfiguration: default configurations reviewed before deployment",
                "▶ Secrets Management",
                "  No secrets in code — use Foundry secrets manager or environment variables",
                "  Secret rotation: quarterly rotation for service account credentials",
                "  Secret scanning: automated scan on every commit",
                "▶ Static Application Security Testing (SAST)",
                "  Bandit (Python), Semgrep — run in CI pipeline on every PR",
                "  SAST findings block merge until resolved or risk-accepted with documentation",
                "▶ Authorized Penetration Testing",
                "  Annual pentest of production MSS applications required",
                "  SL 5L supports pentest scoping and reviews findings report",
            ]),
            ("PLATFORM LEADERSHIP AND ATO SUPPORT (CHAPTER 8)", [
                "▶ Architecture Review",
                "  SL 5L reviews all new application architectures before build begins",
                "  Review criteria: performance, security, tenant isolation, governance compliance",
                "  Architecture decision records (ADR): document why, not just what",
                "▶ Code Review Standards",
                "  Every production code change reviewed by SL 5L or designated reviewer",
                "  Review checklist: security, performance, maintainability, operational readiness",
                "▶ Developer Onboarding",
                "  Standardized onboarding: environment setup, coding standards, first PR review",
                "  Pair programming: SL 5L pairs with new SL 4L practitioners on first deployment",
                "▶ ATO Support Documentation",
                "  SL 5L owns: system architecture diagram, data flow diagrams, security controls",
                "  Continuous monitoring contributions: automated evidence collection for ATO",
                "▶ Managing Platform Upgrades",
                "  Upgrade testing in staging before production; rollback plan required",
                "  Upgrade communication: notify all dependent teams before production rollout",
            ]),
        ]
    )


# ══════════════════════════════════════════════════════════════════════════════
# MAIN
# ══════════════════════════════════════════════════════════════════════════════

def main():
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    print(f"Building 10 new decks → {OUTPUT_DIR}/\n")

    deck_wff_overview()
    deck_tm40m()
    deck_tm40j()
    deck_tm50_overview()
    deck_tm50g()
    deck_tm50h()
    deck_tm50m()
    deck_tm50j()
    deck_tm50k()
    deck_tm50l()

    print("\nDone.")


if __name__ == "__main__":
    main()
